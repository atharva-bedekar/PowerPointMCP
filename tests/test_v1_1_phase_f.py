"""Unit tests for PowerPoint MCP v1.1 Phase F: Verification & End-to-End Workflow."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.pptx.inspector import inspect_slide
from powerpoint_mcp.rendering.visual_compare import compare_slides
from powerpoint_mcp.tools.editing import (
    ppt_align_shapes,
    ppt_apply_style,
    ppt_create_flow_diagram,
    ppt_scale_slide_typography,
)
from powerpoint_mcp.tools.inspection import ppt_analyze_slide_structure, ppt_validate_slide
from powerpoint_mcp.tools.rendering import ppt_render_slide
from powerpoint_mcp.tools.versioning import open_presentation


@pytest.fixture
def diff_deck(tmp_path: Path) -> Path:
    """Create a test deck with two slides showing clear differences."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Original
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(5.0), Inches(1.0))
    tb1.name = "Title"
    tb1.text_frame.text = "Original Slide Header"
    tb1.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

    # Slide 2: Modified (text changed, repositioned, styled)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb2 = s2.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(5.0), Inches(1.0))
    tb2.name = "Title"
    tb2.text_frame.text = "Updated Strategic Slide Header"
    tb2.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    deck_path = tmp_path / "diff_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_structural_slide_diff(diff_deck: Path):
    """Test comprehensive structural difference detection between two slides."""
    model1 = inspect_slide(diff_deck, 1)
    model2 = inspect_slide(diff_deck, 2)

    diff_res = compare_slides(model1, model2)
    assert diff_res.slide_a_number == 1
    assert diff_res.slide_b_number == 2
    assert len(diff_res.layout_differences) >= 1  # dx=0.5, dy=0.2 shift
    assert len(diff_res.typography_differences) >= 1  # text changed + font size 24->28

    # Verify typography difference recorded
    text_change = next(d for d in diff_res.typography_differences if d["type"] == "text_content_change")
    assert text_change["text_a"] == "Original Slide Header"
    assert text_change["text_b"] == "Updated Strategic Slide Header"


def test_render_slide_metadata_and_caching(diff_deck: Path, tmp_path: Path):
    """Test render metadata fields and second-run caching behavior."""
    session = open_presentation(str(diff_deck))
    out_png = tmp_path / "test_slide_1.png"

    # First render: fresh execution
    res1 = ppt_render_slide(slide_number=1, output_path=str(out_png), renderer="mock")
    assert res1["success"] is True
    assert res1["cached"] is False
    assert "render_time_ms" in res1
    assert "dimensions" in res1
    assert out_png.exists()

    # Second render: should hit cache immediately
    res2 = ppt_render_slide(slide_number=1, output_path=str(out_png), renderer="mock")
    assert res2["success"] is True
    assert res2["cached"] is True
    assert res2["renderer"] == "cache"


def test_end_to_end_v1_1_deck_assembly(tmp_path: Path):
    """Complete end-to-end test assembling a professional slide using all v1.1 features."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slides.add_slide(prs.slide_layouts[6])
    deck_path = tmp_path / "e2e_deck.pptx"
    prs.save(str(deck_path))

    session = open_presentation(str(deck_path))

    # 1. Create a 3-step pipeline flow diagram
    diagram_res = ppt_create_flow_diagram(
        slide_number=1,
        steps=[
            {"title": "1. Ingestion", "description": "Kafka events & API ingest"},
            {"title": "2. Processing", "description": "Real-time stream enrichment"},
            {"title": "3. Analytics", "description": "Interactive BI dashboards"},
        ],
        direction="horizontal",
        shape_type="rounded_rectangle",
        style_preset="card_accent",
        start_x=1.0,
        start_y=2.0,
        total_width=11.333,
    )
    assert diagram_res["success"] is True

    # 2. Scale typography across slide proportionally
    scale_res = ppt_scale_slide_typography(slide_number=1, scale_factor=1.1, min_pt=9.0, max_pt=20.0)
    assert scale_res["success"] is True

    # 3. Analyze semantic structure
    struct_res = ppt_analyze_slide_structure(slide_number=1)
    assert struct_res["success"] is True
    assert struct_res["total_elements"] >= 5

    # 4. Validate slide quality
    val_res = ppt_validate_slide(slide_number=1)
    assert val_res["success"] is True
    assert val_res["is_valid"] is True
    assert val_res["summary"]["overlaps"] == 0
    assert val_res["summary"]["boundary_violations"] == 0

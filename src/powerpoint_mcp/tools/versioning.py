"""Non-destructive session management, working-copy tracking, and timestamped backups."""

from dataclasses import asdict, dataclass, field
from datetime import datetime, timezone
import json
import os
from pathlib import Path
import shutil
from typing import Any, Dict, List, Optional, Tuple, Union
import uuid

from pptx import Presentation

from powerpoint_mcp.rendering.com_lifecycle import defensive_file_operation
from powerpoint_mcp.utils.logging import get_logger
from powerpoint_mcp.utils.paths import (
    DEFAULT_WORKING_FILENAME,
    ensure_session_dirs,
    generate_backup_filename,
    get_session_backups_dir,
    get_session_dir,
    get_session_metadata_path,
    get_session_original_path,
    get_session_working_path,
    get_workspace_dir,
)

logger = get_logger("powerpoint_mcp.versioning")


import hashlib

def _iso_now() -> str:
    """Return current UTC time in ISO 8601 format."""
    return datetime.now(timezone.utc).isoformat()


def compute_file_hash(path: Union[str, Path]) -> str:
    """Compute SHA256 hex digest of a file."""
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ""
    h = hashlib.sha256()
    with open(p, "rb") as f:
        while chunk := f.read(65536):
            h.update(chunk)
    return h.hexdigest()


@dataclass
class Session:
    """Model representing an active editing session workspace."""
    session_id: str
    source_path: str
    working_path: str
    created_at: str
    last_modified_at: str
    slide_count: int = 0
    mutation_count: int = 0
    initial_working_hash: Optional[str] = None
    initial_source_hash: Optional[str] = None
    last_working_hash: Optional[str] = None
    backups: List[Dict[str, Any]] = field(default_factory=list)
    renders: List[Dict[str, Any]] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        """Serialize session to dictionary matching metadata.json schema."""
        return {
            "session_id": self.session_id,
            "source_path": str(self.source_path),
            "working_path": str(self.working_path),
            "created_at": self.created_at,
            "last_modified_at": self.last_modified_at,
            "slide_count": self.slide_count,
            "mutation_count": self.mutation_count,
            "initial_working_hash": self.initial_working_hash,
            "initial_source_hash": self.initial_source_hash,
            "last_working_hash": self.last_working_hash,
            "backups": self.backups,
            "renders": self.renders,
            "metadata": self.metadata,
        }

    def save_metadata(self) -> None:
        """Write current session state to metadata.json."""
        meta_path = get_session_metadata_path(self.session_id)
        meta_path.parent.mkdir(parents=True, exist_ok=True)
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(self.to_dict(), f, indent=2)

    def record_mutation(self) -> None:
        """Record that a mutation was performed on the working copy."""
        self.mutation_count += 1
        now_str = _iso_now()
        self.last_modified_at = now_str
        self.last_working_hash = compute_file_hash(self.working_path)
        self.save_metadata()


class SessionManager:
    """Manages isolated session workspaces, working copies, and timestamped backups."""

    def __init__(self, workspace_dir: Optional[Union[str, Path]] = None) -> None:
        self.workspace_dir = get_workspace_dir(workspace_dir)
        self._active_session_id: Optional[str] = None
        self._sessions: Dict[str, Session] = {}

    @property
    def active_session_id(self) -> Optional[str]:
        """Return ID of the currently active session."""
        return self._active_session_id

    @active_session_id.setter
    def active_session_id(self, session_id: Optional[str]) -> None:
        self._active_session_id = session_id

    def open_presentation(
        self,
        presentation_path: Union[str, Path],
        session_id: Optional[str] = None,
    ) -> Session:
        """Open a presentation, initialize an isolated session workspace, and create working copy.

        Args:
            presentation_path: Path to the target .pptx file.
            session_id: Optional custom session ID (defaults to new UUID4).

        Returns:
            Session instance.
        """
        src_path = Path(presentation_path).resolve()
        if not src_path.exists():
            raise FileNotFoundError(f"Presentation not found at path: {src_path}")
        if not src_path.is_file():
            raise ValueError(f"Presentation path is not a file: {src_path}")

        sid = session_id or str(uuid.uuid4())
        dirs = ensure_session_dirs(sid, self.workspace_dir.parent)

        working_path = dirs["working_path"]
        original_path = dirs["original_path"]

        # Copy source presentation to original.pptx and working.pptx
        defensive_file_operation(lambda: shutil.copy2(src_path, original_path), original_path, action_name="open_copy_original")
        defensive_file_operation(lambda: shutil.copy2(src_path, working_path), working_path, action_name="open_copy_working")

        src_hash = compute_file_hash(src_path)
        working_hash = compute_file_hash(working_path)

        # Inspect slide count from working copy
        try:
            prs = Presentation(str(working_path))
            slide_count = len(prs.slides)
        except Exception as e:
            logger.warning(f"Could not inspect slide count from {working_path}: {e}")
            slide_count = 0

        now_str = _iso_now()
        session = Session(
            session_id=sid,
            source_path=str(src_path),
            working_path=str(working_path),
            created_at=now_str,
            last_modified_at=now_str,
            slide_count=slide_count,
            mutation_count=0,
            initial_working_hash=working_hash,
            initial_source_hash=src_hash,
            last_working_hash=working_hash,
            backups=[],
            renders=[],
            metadata={
                "source_filename": src_path.name,
                "file_size_bytes": src_path.stat().st_size,
            },
        )

        session.save_metadata()
        self._sessions[sid] = session
        self._active_session_id = sid

        logger.info(f"Opened presentation session {sid} for {src_path.name}")
        return session

    def get_session(self, session_id: Optional[str] = None) -> Optional[Session]:
        """Retrieve a session by ID or return the currently active session.

        Args:
            session_id: Session identifier. If None, returns active session.

        Returns:
            Session instance or None if not found.
        """
        sid = session_id or self._active_session_id
        if not sid:
            return None

        if sid in self._sessions:
            return self._sessions[sid]

        # Try loading from metadata.json on disk
        meta_path = get_session_metadata_path(sid, self.workspace_dir.parent)
        if meta_path.exists():
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                session = Session(
                    session_id=data["session_id"],
                    source_path=data["source_path"],
                    working_path=data["working_path"],
                    created_at=data["created_at"],
                    last_modified_at=data["last_modified_at"],
                    slide_count=data.get("slide_count", 0),
                    mutation_count=data.get("mutation_count", 0),
                    initial_working_hash=data.get("initial_working_hash"),
                    initial_source_hash=data.get("initial_source_hash"),
                    last_working_hash=data.get("last_working_hash"),
                    backups=data.get("backups", []),
                    renders=data.get("renders", []),
                    metadata=data.get("metadata", {}),
                )
                self._sessions[sid] = session
                return session
            except Exception as e:
                logger.error(f"Error loading session metadata from {meta_path}: {e}")
                return None

        return None

        return None

    def get_current_session(self) -> Optional[Session]:
        """Alias for get_session(None)."""
        return self.get_session(None)

    def create_backup(
        self,
        session_id_or_path: Optional[str] = None,
        operation: str = "manual_backup",
        label: Optional[str] = None,
        details: Optional[Dict[str, Any]] = None,
    ) -> str:
        """Create a timestamped backup of the session working copy or standalone presentation file.

        Args:
            session_id_or_path: Session ID, file path, or None (uses active session).
            operation: Description of the operation triggering the backup.
            label: Optional human-readable backup label.
            details: Optional extra metadata dictionary.

        Returns:
            Absolute path string of the created backup file.
        """
        # Check if backups are disabled by environment
        backup_enabled = os.environ.get("PPT_BACKUP_ENABLED", "true").lower() in ("true", "1", "yes")
        if not backup_enabled:
            logger.info("Backups disabled via PPT_BACKUP_ENABLED")

        # Case 1: Active session or session ID
        session = None
        if session_id_or_path is None or (
            isinstance(session_id_or_path, str) and session_id_or_path in self._sessions
        ):
            session = self.get_session(session_id_or_path)
        elif isinstance(session_id_or_path, str):
            # Check if it looks like a session ID
            meta_path = get_session_metadata_path(session_id_or_path, self.workspace_dir.parent)
            if meta_path.exists():
                session = self.get_session(session_id_or_path)

        if session is not None:
            working_file = Path(session.working_path)
            if not working_file.exists():
                raise FileNotFoundError(f"Working copy not found: {working_file}")

            backups_dir = get_session_backups_dir(session.session_id, self.workspace_dir.parent)
            backups_dir.mkdir(parents=True, exist_ok=True)

            source_stem = Path(session.source_path).stem or "presentation"
            backup_filename = generate_backup_filename(source_stem)
            backup_path = backups_dir / backup_filename

            # If filename collision within same second, append microsecond
            if backup_path.exists():
                ts = datetime.now().strftime("%Y%m%d-%H%M%S-%f")
                backup_path = backups_dir / f"{source_stem}.backup-{ts}.pptx"

            defensive_file_operation(lambda: shutil.copy2(working_file, backup_path), backup_path, action_name="create_backup")

            now_str = _iso_now()
            backup_record = {
                "timestamp": now_str,
                "backup_path": str(backup_path),
                "operation": operation,
                "label": label or operation,
                "details": details or {},
            }
            session.backups.append(backup_record)
            session.last_modified_at = now_str
            session.save_metadata()

            logger.info(f"Created session backup at {backup_path}")
            return str(backup_path)

        # Case 2: Standalone file path provided
        target_path = Path(session_id_or_path or "").resolve()
        if not target_path.exists():
            raise FileNotFoundError(f"Cannot create backup: file does not exist at {target_path}")

        parent_dir = target_path.parent
        backup_filename = generate_backup_filename(target_path.stem)
        backup_path = parent_dir / backup_filename
        if backup_path.exists():
            ts = datetime.now().strftime("%Y%m%d-%H%M%S-%f")
            backup_path = parent_dir / f"{target_path.stem}.backup-{ts}.pptx"

        defensive_file_operation(lambda: shutil.copy2(target_path, backup_path), backup_path, action_name="create_standalone_backup")
        logger.info(f"Created standalone file backup at {backup_path}")
        return str(backup_path)

    def revert_session(
        self,
        session_id: Optional[str] = None,
        backup_path_or_target: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Revert the session working copy to the original file or a specific backup snapshot.

        Args:
            session_id: Session identifier (defaults to active session).
            backup_path_or_target: "original", backup timestamp string, or backup file path.

        Returns:
            Dictionary describing the revert operation outcome.
        """
        session = self.get_session(session_id)
        if session is None:
            raise ValueError(f"Session not found: {session_id or self._active_session_id}")

        working_file = Path(session.working_path)
        original_file = Path(get_session_original_path(session.session_id, self.workspace_dir.parent))

        target_str = (backup_path_or_target or "original").strip()
        reverted_from: str = ""

        if target_str == "original" or not target_str:
            if not original_file.exists():
                raise FileNotFoundError(f"Original presentation snapshot not found at {original_file}")
            defensive_file_operation(lambda: shutil.copy2(original_file, working_file), working_file, action_name="revert_original")
            reverted_from = "original"
        else:
            # Check if target_str is a direct file path
            candidate_path = Path(target_str)
            if candidate_path.is_file() and candidate_path.exists():
                defensive_file_operation(lambda: shutil.copy2(candidate_path, working_file), working_file, action_name="revert_custom_path")
                reverted_from = str(candidate_path)
            else:
                # Search within session backups directory or backup records
                backups_dir = get_session_backups_dir(session.session_id, self.workspace_dir.parent)
                found = None
                for b_file in backups_dir.glob("*.pptx"):
                    if target_str in b_file.name:
                        found = b_file
                        break

                if found is None:
                    # Check backup records by timestamp
                    for record in session.backups:
                        if target_str in record.get("timestamp", "") or target_str in record.get("backup_path", ""):
                            found = Path(record["backup_path"])
                            break

                if found is None or not found.exists():
                    raise FileNotFoundError(f"Backup snapshot matching '{target_str}' not found in session backups")

                defensive_file_operation(lambda: shutil.copy2(found, working_file), working_file, action_name="revert_backup")
                reverted_from = str(found)

        now_str = _iso_now()
        session.last_modified_at = now_str
        session.metadata["last_revert"] = {
            "timestamp": now_str,
            "reverted_from": reverted_from,
        }
        session.save_metadata()

        logger.info(f"Reverted session {session.session_id} to {reverted_from}")
        return {
            "success": True,
            "session_id": session.session_id,
            "reverted_to": reverted_from,
            "working_path": str(working_file),
            "timestamp": now_str,
        }

    def save_session(
        self,
        session_id: Optional[str] = None,
        destination_path: Optional[Union[str, Path]] = None,
    ) -> Dict[str, Any]:
        """Commit session changes to the original source file (creating pre-save backup).

        Args:
            session_id: Session identifier (defaults to active session).
            destination_path: Optional override destination path. Defaults to session source_path.

        Returns:
            Dictionary detailing saved path and pre-save backup path.
        """
        session = self.get_session(session_id)
        if session is None:
            raise ValueError(f"Session not found: {session_id or self._active_session_id}")

        working_file = Path(session.working_path)
        if not working_file.exists() or working_file.stat().st_size == 0:
            raise FileNotFoundError(f"Session working copy missing or empty at {working_file}")

        # Integrity Check 1: Verify working copy can be parsed as a valid Presentation
        try:
            test_prs = Presentation(str(working_file))
            if len(test_prs.slides) == 0 and session.slide_count > 0:
                raise ValueError("Working copy has 0 slides but original had slides.")
        except Exception as exc:
            raise ValueError(f"Session working copy is corrupted or invalid PPTX: {exc}")

        current_working_hash = compute_file_hash(working_file)

        target_path = Path(destination_path).resolve() if destination_path else Path(session.source_path).resolve()
        target_path.parent.mkdir(parents=True, exist_ok=True)

        backup_path: Optional[str] = None
        # If target file exists, create a pre-save backup before overwriting
        if target_path.exists():
            backup_path = self.create_backup(
                session.session_id,
                operation="pre_save_backup",
                label="Pre-Save Backup",
                details={"saved_to": str(target_path)},
            )

        defensive_file_operation(lambda: shutil.copy2(working_file, target_path), target_path, action_name="save_session")

        # Integrity Check 2: Verify target file was successfully saved
        try:
            Presentation(str(target_path))
        except Exception as exc:
            raise ValueError(f"Saved presentation failed verification after save: {exc}")

        saved_hash = compute_file_hash(target_path)
        now_str = _iso_now()
        session.last_modified_at = now_str
        session.metadata["last_save"] = {
            "timestamp": now_str,
            "saved_to": str(target_path),
            "backup_path": backup_path,
            "working_hash": current_working_hash,
            "saved_hash": saved_hash,
            "mutation_count": session.mutation_count,
        }
        session.save_metadata()

        logger.info(f"Saved session {session.session_id} to {target_path}")
        return {
            "success": True,
            "session_id": session.session_id,
            "saved_path": str(target_path),
            "backup_path": backup_path,
            "mutation_count": session.mutation_count,
            "timestamp": now_str,
        }

    def save_as(
        self,
        session_id: Optional[str] = None,
        output_path: Union[str, Path] = None,
        overwrite: bool = False,
    ) -> Dict[str, Any]:
        """Save the session working copy to a new output path without modifying the original.

        Args:
            session_id: Session identifier (defaults to active session).
            output_path: Target destination path for the copy.
            overwrite: If False, raises FileExistsError if output_path exists.

        Returns:
            Dictionary detailing save-as confirmation.
        """
        if output_path is None:
            raise ValueError("output_path must be specified for save_as")

        session = self.get_session(session_id)
        if session is None:
            raise ValueError(f"Session not found: {session_id or self._active_session_id}")

        working_file = Path(session.working_path)
        if not working_file.exists() or working_file.stat().st_size == 0:
            raise FileNotFoundError(f"Session working copy missing or empty at {working_file}")

        # Integrity Check: verify working copy validity
        try:
            Presentation(str(working_file))
        except Exception as exc:
            raise ValueError(f"Session working copy is corrupted or invalid PPTX: {exc}")

        out_path = Path(output_path).resolve()
        if out_path.exists() and not overwrite:
            raise FileExistsError(f"Destination file already exists and overwrite is False: {out_path}")

        out_path.parent.mkdir(parents=True, exist_ok=True)

        backup_path: Optional[str] = None
        if out_path.exists() and overwrite:
            backup_path = self.create_backup(
                session.session_id,
                operation="pre_save_as_overwrite_backup",
                label="Pre-Save-As Overwrite Backup",
                details={"target": str(out_path)},
            )

        defensive_file_operation(lambda: shutil.copy2(working_file, out_path), out_path, action_name="save_as")

        # Integrity Check: verify written output file
        try:
            Presentation(str(out_path))
        except Exception as exc:
            raise ValueError(f"Saved-as presentation failed verification: {exc}")

        saved_hash = compute_file_hash(out_path)
        now_str = _iso_now()
        session.last_modified_at = now_str
        session.metadata["last_save_as"] = {
            "timestamp": now_str,
            "saved_as": str(out_path),
            "backup_path": backup_path,
            "saved_hash": saved_hash,
            "mutation_count": session.mutation_count,
        }
        session.save_metadata()

        logger.info(f"Saved-as session {session.session_id} to {out_path}")
        return {
            "success": True,
            "session_id": session.session_id,
            "saved_path": str(out_path),
            "backup_path": backup_path,
            "mutation_count": session.mutation_count,
            "timestamp": now_str,
        }

    def close_session(self, session_id: Optional[str] = None, cleanup: bool = False) -> bool:
        """Close session and optionally remove session directory from disk."""
        sid = session_id or self._active_session_id
        if not sid:
            return False

        if sid in self._sessions:
            del self._sessions[sid]

        if self._active_session_id == sid:
            self._active_session_id = None

        if cleanup:
            from powerpoint_mcp.utils.paths import cleanup_session
            cleanup_session(sid, self.workspace_dir.parent)

        return True


# Global default SessionManager instance
_DEFAULT_MANAGER: Optional[SessionManager] = None


def get_session_manager(workspace_dir: Optional[Union[str, Path]] = None) -> SessionManager:
    """Return the global default SessionManager instance."""
    global _DEFAULT_MANAGER
    if _DEFAULT_MANAGER is None:
        _DEFAULT_MANAGER = SessionManager(workspace_dir)
    elif workspace_dir is not None:
        _DEFAULT_MANAGER.workspace_dir = get_workspace_dir(workspace_dir)
    return _DEFAULT_MANAGER


def open_presentation(
    presentation_path: Union[str, Path],
    session_id: Optional[str] = None,
) -> Session:
    """Open a presentation session via the default SessionManager."""
    return get_session_manager().open_presentation(presentation_path, session_id=session_id)


def create_backup(
    session_id_or_path: Optional[str] = None,
    operation: str = "manual_backup",
    label: Optional[str] = None,
) -> str:
    """Create a backup via the default SessionManager."""
    return get_session_manager().create_backup(session_id_or_path, operation=operation, label=label)


def revert_session(
    session_id: Optional[str] = None,
    backup_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Revert session via the default SessionManager."""
    return get_session_manager().revert_session(session_id, backup_path_or_target=backup_path)


def save_session(
    session_id: Optional[str] = None,
    destination_path: Optional[Union[str, Path]] = None,
) -> Dict[str, Any]:
    """Save session changes back to the source file via the default SessionManager."""
    return get_session_manager().save_session(session_id, destination_path=destination_path)


def save_as(
    session_id: Optional[str] = None,
    output_path: Union[str, Path] = None,
    overwrite: bool = False,
) -> Dict[str, Any]:
    """Save session changes to a new destination file via the default SessionManager."""
    return get_session_manager().save_as(session_id, output_path=output_path, overwrite=overwrite)


def get_session(session_id: Optional[str] = None) -> Optional[Session]:
    """Get session by ID or active session from the default SessionManager."""
    return get_session_manager().get_session(session_id)


def get_current_session() -> Optional[Session]:
    """Get active session from the default SessionManager."""
    return get_session_manager().get_current_session()


def resolve_active_target(
    presentation_path: Optional[str] = None,
    require_session: bool = False,
    mutation: bool = False,
    operation: Optional[str] = None,
) -> Tuple[str, Optional[Session]]:
    """Canonical presentation target resolution mechanism across all tools.

    Rules:
    1. If an active session exists:
       - If mutation=True:
         - All mutations MUST target session.working_path.
         - If explicit presentation_path is provided:
           - If it matches session.source_path or session.working_path (normalized), redirect seamlessly to working_path.
           - If it points to a completely different presentation, raise ValueError (conflicting presentation_path).
       - If mutation=False (read/inspect/render):
         - If presentation_path is None or matches source/working copy, resolve to session.working_path.
         - If explicit presentation_path points to another file, resolve to that other file.
       - If mutation=True, create safety backup before mutation and increment session mutation tracking.
       - Return (resolved_path_str, session).
    2. If no active session exists:
       - If require_session=True, raise ValueError.
       - If presentation_path is provided:
         - Verify file exists.
         - If mutation=True, create standalone file backup.
         - Return (resolved_path_str, None).
       - Else:
         - Raise ValueError.
    """
    mgr = get_session_manager()
    session = mgr.get_current_session()

    if session is not None and session.working_path and Path(session.working_path).exists():
        working_p = Path(session.working_path).resolve()
        source_p = Path(session.source_path).resolve()

        if presentation_path:
            req_p = Path(presentation_path).resolve()
            # If presentation_path matches source or working copy, target working copy
            if req_p == source_p or req_p == working_p:
                target_p = working_p
            else:
                if mutation:
                    raise ValueError(
                        f"Active session '{session.session_id}' is targeting '{session.source_path}'. "
                        f"Conflicting presentation_path '{presentation_path}' provided for mutation. "
                        f"Cannot mutate another presentation while an active session exists. "
                        f"Please close the session or omit presentation_path to edit the active session."
                    )
                else:
                    if not req_p.exists():
                        raise FileNotFoundError(f"Presentation not found: {req_p}")
                    return str(req_p), session
        else:
            target_p = working_p

        # Target is session working copy
        target_path_str = str(target_p)
        if mutation:
            try:
                mgr.create_backup(
                    session.session_id,
                    operation=operation or "mutation",
                    label=f"Pre-mutation backup ({operation or 'edit'})",
                )
            except Exception:
                pass
            session.record_mutation()

        return target_path_str, session

    # No active session
    if require_session:
        raise ValueError("No active editing session found. Please call ppt_open first.")

    if presentation_path:
        p = Path(presentation_path).resolve()
        if not p.exists():
            raise FileNotFoundError(f"Presentation not found: {p}")
        target_path_str = str(p)
        if mutation:
            try:
                mgr.create_backup(target_path_str, operation=operation or "standalone_mutation")
            except Exception:
                pass
        return target_path_str, None

    raise ValueError("No presentation path provided and no active editing session found. Please call ppt_open first.")

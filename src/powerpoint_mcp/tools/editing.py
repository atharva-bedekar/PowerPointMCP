"""Editing tools for PowerPoint MCP server."""

from functools import wraps
import os
from pathlib import Path
import traceback
from typing import Any, Dict, List, Optional, Union

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import (
    AlignmentType,
    DistributionMode,
    SpacingMode,
    emu_to_inches,
    inches_to_emu,
)
from powerpoint_mcp.pptx.diagrams import create_flow_diagram
from powerpoint_mcp.pptx.editor import (
    copy_shape,
    delete_shape,
    modify_shape,
    modify_text,
    move_shape,
    resize_shape,
    scale_slide_typography,
)
from powerpoint_mcp.pptx.geometry import (
    align_shapes,
    distribute_shapes,
    equalize_dimensions,
    space_shapes,
)
from powerpoint_mcp.pptx.structure import (
    move_container,
    reflow_container,
    resize_container,
)
from powerpoint_mcp.pptx.ooxml import (
    NAMESPACES,
    get_raw_shape_xml,
    safe_modify_xml,
    set_drop_shadow,
    set_gradient_fill,
    set_shape_transparency,
)
from powerpoint_mcp.pptx.styles import (
    STYLE_PRESETS,
    apply_style_to_shape,
    extract_complete_shape_style,
)
from powerpoint_mcp.tools.inspection import handle_tool_errors
from powerpoint_mcp.tools.versioning import get_session_manager, resolve_active_target


def _get_target_presentation(
    presentation_path: Optional[str] = None, operation: str = "edit"
) -> Any:
    """Resolve presentation file using canonical resolve_active_target and ensure pre-mutation backup is made.

    Returns:
        tuple of (resolved_path_str, Presentation_instance, session_or_None)
    """
    target_path_str, session = resolve_active_target(
        presentation_path=presentation_path,
        require_session=False,
        mutation=True,
        operation=operation,
    )
    prs = Presentation(target_path_str)
    return target_path_str, prs, session



# Simple helper type for annotation
Tuple_Target = Any


@handle_tool_errors
def ppt_modify_shape(
    slide_number: int,
    shape_id: int,
    x: Optional[float] = None,
    y: Optional[float] = None,
    width: Optional[float] = None,
    height: Optional[float] = None,
    rotation: Optional[float] = None,
    z_order: Optional[Union[int, str]] = None,
    dx: Optional[float] = None,
    dy: Optional[float] = None,
    dwidth: Optional[float] = None,
    dheight: Optional[float] = None,
    drotation: Optional[float] = None,
    align: Optional[str] = None,
    distribute: Optional[str] = None,
    target_shape_ids: Optional[List[int]] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Modify shape coordinates, dimensions, rotation, z-order, alignment, or distribution.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: ID of primary target shape.
        x: Absolute X coordinate in inches.
        y: Absolute Y coordinate in inches.
        width: Absolute width in inches. Must be > 0.
        height: Absolute height in inches. Must be > 0.
        rotation: Absolute rotation in degrees (0..360).
        z_order: Z-order index or action ('bring_to_front', 'send_to_back', 'bring_forward', 'send_backward').
        dx: Delta X shift in inches.
        dy: Delta Y shift in inches.
        dwidth: Delta width in inches.
        dheight: Delta height in inches.
        drotation: Delta rotation in degrees.
        align: Alignment mode ('left', 'center', 'right', 'top', 'middle', 'bottom').
        distribute: Distribution mode ('horizontal', 'vertical').
        target_shape_ids: Additional shape IDs to align/distribute with primary shape.
        presentation_path: Presentation path. If omitted, uses active session working copy.

    Returns:
        Structured dictionary detailing updated shape geometry.
    """
    if width is not None and width <= 0:
        raise ValueError(f"Shape width must be positive, got {width}")
    if height is not None and height <= 0:
        raise ValueError(f"Shape height must be positive, got {height}")

    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"modify_shape_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]

    # Check multi-shape alignment / distribution
    if align or distribute:
        all_ids = [shape_id] + (target_shape_ids or [])
        selected_shapes = []
        for sid in all_ids:
            found = False
            for s in slide.shapes:
                if s.shape_id == sid:
                    selected_shapes.append(s)
                    found = True
                    break
            if not found:
                raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")

        if align:
            align_shapes(selected_shapes, align)
        if distribute:
            distribute_shapes(selected_shapes, distribute)

        prs.save(target_path)
        if session:
            session.save_metadata()

        return {
            "success": True,
            "shape_id": shape_id,
            "operation": "multi_shape_geometry",
            "aligned": align,
            "distributed": distribute,
            "affected_shape_ids": all_ids,
        }

    # Single shape modification
    res = modify_shape(
        slide,
        shape_id,
        x=x,
        y=y,
        width=width,
        height=height,
        rotation=rotation,
        z_order=z_order,
        dx=dx,
        dy=dy,
        dwidth=dwidth,
        dheight=dheight,
        drotation=drotation,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "shape_id": shape_id,
        "updated_properties": {
            "x_inches": res["x"],
            "y_inches": res["y"],
            "width_inches": res["width"],
            "height_inches": res["height"],
            "rotation": res["rotation"],
        },
        "shape": res,
    }


@handle_tool_errors
def ppt_modify_text(
    slide_number: int,
    shape_id: int,
    text: Optional[str] = None,
    font_family: Optional[str] = None,
    font_name: Optional[str] = None,
    font_size: Optional[float] = None,
    font_size_pt: Optional[float] = None,
    font_size_delta: Optional[float] = None,
    font_size_scale: Optional[float] = None,
    min_font_size: Optional[float] = None,
    max_font_size: Optional[float] = None,
    min_pt: Optional[float] = None,
    max_pt: Optional[float] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    color: Optional[str] = None,
    color_rgb: Optional[str] = None,
    alignment: Optional[str] = None,
    paragraph_spacing: Optional[float] = None,
    space_before: Optional[float] = None,
    space_after: Optional[float] = None,
    line_spacing: Optional[float] = None,
    margins: Optional[Dict[str, float]] = None,
    paragraph_index: Optional[int] = None,
    run_index: Optional[int] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Modify text content, typography, colors, and margins while preserving surrounding styles.

    Supports absolute font sizing, relative point deltas, and scale multipliers with min/max bounds.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: Target shape ID.
        text: New text string.
        font_family / font_name: Font name (e.g. 'Calibri', 'Arial', 'Aptos').
        font_size / font_size_pt: Font point size (absolute).
        font_size_delta: Relative point delta (+2, -2) to adjust current font size.
        font_size_scale: Scale multiplier (e.g. 1.15) to proportionally scale font size.
        min_font_size / min_pt: Lower bound clamp for font size in points.
        max_font_size / max_pt: Upper bound clamp for font size in points.
        bold: Bold weight flag.
        italic: Italic flag.
        underline: Underline flag.
        color / color_rgb: Hex RGB color string (e.g. '#1F497D' or '1F497D').
        alignment: Text alignment ('left', 'center', 'right', 'justify').
        paragraph_spacing / space_before: Space before paragraph in points.
        space_after: Space after paragraph in points.
        line_spacing: Line spacing in points.
        margins: Margin dict in inches: {'left': 0.1, 'top': 0.05, 'right': 0.1, 'bottom': 0.05}.
        paragraph_index: Optional 0-indexed paragraph to modify.
        run_index: Optional 0-indexed run to modify.
        presentation_path: Presentation path.

    Returns:
        Structured dictionary detailing updated text, font sizes, and paragraph counts.
    """
    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"modify_text_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]

    res = modify_text(
        slide,
        shape_id,
        text=text,
        font_family=font_family or font_name,
        font_size=font_size if font_size is not None else font_size_pt,
        font_size_delta=font_size_delta,
        font_size_scale=font_size_scale,
        min_font_size=min_font_size if min_font_size is not None else min_pt,
        max_font_size=max_font_size if max_font_size is not None else max_pt,
        bold=bold,
        italic=italic,
        underline=underline,
        color=color or color_rgb,
        alignment=alignment,
        paragraph_spacing=paragraph_spacing if paragraph_spacing is not None else space_before,
        space_after=space_after,
        line_spacing=line_spacing,
        margins=margins,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "shape_id": shape_id,
        "text_summary": (res.get("text") or "")[:100],
        "paragraph_count": res.get("paragraph_count", 1),
        "font_size": res.get("font_size"),
        "original_font_size": res.get("original_font_size"),
        "resulting_font_size": res.get("resulting_font_size"),
        "font_family": res.get("font_name"),
    }


@handle_tool_errors
def ppt_copy_shape(
    slide_number: int,
    shape_id: int,
    target_slide_number: Optional[int] = None,
    x_offset: float = 0.2,
    y_offset: float = 0.2,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Clone an existing shape with all styling and relationships onto the same or a different slide.

    Args:
        slide_number: 1-indexed source slide number.
        shape_id: ID of shape to clone.
        target_slide_number: Destination slide number (defaults to same slide).
        x_offset: Horizontal offset in inches for copied shape.
        y_offset: Vertical offset in inches for copied shape.
        presentation_path: Presentation path.

    Returns:
        Structured dictionary containing new_shape_id and target_slide number.
    """
    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"copy_shape_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    dest_slide_num = target_slide_number if target_slide_number is not None else slide_number
    if dest_slide_num < 1 or dest_slide_num > len(prs.slides):
        raise IndexError(f"Target slide number {dest_slide_num} is out of range (1..{len(prs.slides)})")

    src_slide = prs.slides[slide_number - 1]
    dest_slide = prs.slides[dest_slide_num - 1]

    new_shape_id = copy_shape(
        src_slide,
        shape_id,
        target_slide=dest_slide,
        offset_x_inches=x_offset,
        offset_y_inches=y_offset,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "new_shape_id": new_shape_id,
        "source_shape_id": shape_id,
        "target_slide": dest_slide_num,
    }


@handle_tool_errors
def ppt_move_shape(
    slide_number: int,
    shape_id: int,
    dx: Optional[float] = None,
    dy: Optional[float] = None,
    x: Optional[float] = None,
    y: Optional[float] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Move a shape by absolute coordinates (x, y) or relative deltas (dx, dy) in inches.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: Target shape ID.
        dx: Relative delta X in inches.
        dy: Relative delta Y in inches.
        x: Absolute X coordinate in inches.
        y: Absolute Y coordinate in inches.
        presentation_path: Presentation path.

    Returns:
        Updated shape position dictionary.
    """
    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"move_shape_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = move_shape(slide, shape_id, dx=dx, dy=dy, x=x, y=y)

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "shape_id": shape_id,
        "x_inches": res["x"],
        "y_inches": res["y"],
    }


@handle_tool_errors
def ppt_resize_shape(
    slide_number: int,
    shape_id: int,
    width: Optional[float] = None,
    height: Optional[float] = None,
    scale_width: Optional[float] = None,
    scale_height: Optional[float] = None,
    scale_x: Optional[float] = None,
    scale_y: Optional[float] = None,
    lock_aspect_ratio: Optional[bool] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Resize a shape using absolute dimensions (width, height) or scale multipliers.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: Target shape ID.
        width: Absolute width in inches. Must be > 0.
        height: Absolute height in inches. Must be > 0.
        scale_width / scale_x: Width scale multiplier (e.g. 1.2 = +20%).
        scale_height / scale_y: Height scale multiplier.
        lock_aspect_ratio: Maintain aspect ratio when scaling.
        presentation_path: Presentation path.

    Returns:
        Updated shape dimension dictionary.
    """
    if width is not None and width <= 0:
        raise ValueError(f"Width must be positive, got {width}")
    if height is not None and height <= 0:
        raise ValueError(f"Height must be positive, got {height}")

    eff_scale_x = scale_width if scale_width is not None else scale_x
    eff_scale_y = scale_height if scale_height is not None else scale_y

    if lock_aspect_ratio and eff_scale_x is not None and eff_scale_y is None:
        eff_scale_y = eff_scale_x
    elif lock_aspect_ratio and eff_scale_y is not None and eff_scale_x is None:
        eff_scale_x = eff_scale_y

    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"resize_shape_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = resize_shape(slide, shape_id, width=width, height=height, scale_x=eff_scale_x, scale_y=eff_scale_y)

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "shape_id": shape_id,
        "width_inches": res["width"],
        "height_inches": res["height"],
    }


@handle_tool_errors
def ppt_delete_shape(
    slide_number: int,
    shape_id: int,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Delete a shape cleanly from a slide.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: ID of shape to delete.
        presentation_path: Presentation path.

    Returns:
        Summary of deleted shape and remaining shape count/IDs.
    """
    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"delete_shape_{shape_id}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    deleted = delete_shape(slide, shape_id)

    if not deleted:
        raise ValueError(f"Failed to delete shape {shape_id}")

    prs.save(target_path)
    if session:
        session.save_metadata()

    remaining_ids = [s.shape_id for s in slide.shapes]

    return {
        "success": True,
        "deleted_shape_id": shape_id,
        "remaining_shape_count": len(remaining_ids),
        "remaining_shape_ids": remaining_ids,
    }


@handle_tool_errors
def ppt_modify_ooxml(
    slide_number: int,
    shape_id: Optional[int] = None,
    operation: str = "set_attribute",
    xpath: Optional[str] = None,
    attributes: Optional[Dict[str, str]] = None,
    xml_fragment: Optional[str] = None,
    transparency_percent: Optional[float] = None,
    gradient_start: Optional[str] = None,
    gradient_end: Optional[str] = None,
    gradient_angle: float = 90.0,
    shadow_blur_pt: Optional[float] = None,
    shadow_dist_pt: Optional[float] = None,
    shadow_color: Optional[str] = None,
    shadow_alpha: Optional[float] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Controlled low-level OOXML manipulation helper for gradients, transparency, drop shadows, and XML tags.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: Optional target shape ID (if modifying a specific shape).
        operation: Operation type ('set_attribute', 'insert_element', 'replace_element', 'transparency', 'gradient', 'shadow').
        xpath: XPath targeting the element to mutate.
        attributes: Key-value attributes to apply.
        xml_fragment: Valid XML snippet to insert or replace.
        transparency_percent: Transparency percentage (0.0 = solid, 100.0 = transparent).
        gradient_start: Hex start color for gradient fill (e.g. 'FFFFFF').
        gradient_end: Hex end color for gradient fill (e.g. '000000').
        gradient_angle: Linear gradient angle in degrees (default 90.0).
        shadow_blur_pt: Drop shadow blur radius in points.
        shadow_dist_pt: Drop shadow distance offset in points.
        shadow_color: Drop shadow color hex.
        shadow_alpha: Drop shadow opacity percentage (0..100).
        presentation_path: Presentation path.

    Returns:
        Result summary with XML snippet.
    """
    target_path, prs, session = _get_target_presentation(presentation_path, operation=f"ooxml_{operation}")

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]

    # Target shape or slide root
    target_obj = None
    if shape_id is not None:
        for s in slide.shapes:
            if s.shape_id == shape_id:
                target_obj = s
                break
        if target_obj is None:
            raise ValueError(f"Shape with ID {shape_id} not found on slide {slide_number}")
    else:
        target_obj = slide._element

    # 1. Transparency helper
    if transparency_percent is not None or operation == "transparency":
        tp = transparency_percent if transparency_percent is not None else 50.0
        set_shape_transparency(target_obj, tp)

    # 2. Gradient helper
    elif gradient_start is not None or gradient_end is not None or operation == "gradient":
        set_gradient_fill(
            target_obj,
            start_hex=gradient_start or "FFFFFF",
            end_hex=gradient_end or "000000",
            angle_deg=gradient_angle,
        )

    # 3. Drop shadow helper
    elif (
        shadow_blur_pt is not None
        or shadow_dist_pt is not None
        or shadow_color is not None
        or operation == "shadow"
    ):
        set_drop_shadow(
            target_obj,
            blur_rad_pt=shadow_blur_pt if shadow_blur_pt is not None else 4.0,
            dist_pt=shadow_dist_pt if shadow_dist_pt is not None else 3.0,
            color_hex=shadow_color or "000000",
            alpha_percent=shadow_alpha if shadow_alpha is not None else 40.0,
        )

    # 4. Custom XML mutation via XPath
    elif xpath or attributes or xml_fragment:
        def mutate_fn(elem: etree._Element):
            target_nodes = elem.xpath(xpath, namespaces=NAMESPACES) if xpath else [elem]
            if not target_nodes:
                raise ValueError(f"XPath '{xpath}' did not match any elements")

            for node in target_nodes:
                if operation == "set_attribute" and attributes:
                    for k, v in attributes.items():
                        node.set(k, str(v))
                elif operation == "insert_element" and xml_fragment:
                    frag = etree.fromstring(xml_fragment.encode("utf-8"))
                    node.append(frag)
                elif operation == "replace_element" and xml_fragment:
                    frag = etree.fromstring(xml_fragment.encode("utf-8"))
                    p = node.getparent()
                    if p is not None:
                        idx = p.index(node)
                        p.remove(node)
                        p.insert(idx, frag)

        safe_modify_xml(target_obj, mutate_fn)

    prs.save(target_path)
    if session:
        session.save_metadata()

    snippet = get_raw_shape_xml(target_obj) if hasattr(target_obj, "_element") else ""
    return {
        "success": True,
        "operation": operation,
        "shape_id": shape_id,
        "xml_snippet": snippet[:500] if snippet else "<modified/>",
    }


@handle_tool_errors
def ppt_batch_modify_text(
    slide_number: int,
    operations: List[Dict[str, Any]],
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Modify multiple text shapes on a slide in a single transaction with pre-validation.

    Args:
        slide_number: 1-indexed slide number.
        operations: List of shape edit dicts. Each dict supports:
            - shape_id: Target shape ID (int, required)
            - text: New text content (preserves paragraphs/bullets unless explicitly changed)
            - font_family / font_name: Font family name
            - font_size / font_size_pt: Absolute font size in points
            - font_size_delta: Relative point delta (+2, -2) to adjust font size
            - font_size_scale: Scale multiplier (e.g. 1.15) to scale font size
            - min_font_size / min_pt: Lower bound clamp for font size in points
            - max_font_size / max_pt: Upper bound clamp for font size in points
            - bold, italic, underline: Typography flags
            - color / color_rgb: Hex RGB color string (e.g. '#1F497D')
            - alignment: Text alignment ('left', 'center', 'right', 'justify')
            - paragraph_spacing / space_before: Points
            - space_after: Points
            - line_spacing: Points
            - margins: Margin dict in inches
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Structured batch summary detailing applied changes per shape ID.
    """
    if not operations:
        raise ValueError("Operations list cannot be empty")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"batch_modify_text_s{slide_number}_n{len(operations)}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    # Step 1: Pre-validation of all operations before mutating
    for idx, op in enumerate(operations):
        if not isinstance(op, dict):
            raise ValueError(f"Operation at index {idx} must be a dictionary")
        sid = op.get("shape_id")
        if sid is None:
            raise ValueError(f"Operation at index {idx} is missing required 'shape_id'")
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")
        shape = shape_map[sid]
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(f"Shape with ID {sid} does not support a text frame")

    # Step 2: Apply mutations
    results = []
    for op in operations:
        sid = op["shape_id"]
        res = modify_text(
            slide,
            sid,
            text=op.get("text"),
            font_family=op.get("font_family") or op.get("font_name"),
            font_size=op.get("font_size") if op.get("font_size") is not None else op.get("font_size_pt"),
            font_size_delta=op.get("font_size_delta"),
            font_size_scale=op.get("font_size_scale"),
            min_font_size=op.get("min_font_size") if op.get("min_font_size") is not None else op.get("min_pt"),
            max_font_size=op.get("max_font_size") if op.get("max_font_size") is not None else op.get("max_pt"),
            bold=op.get("bold"),
            italic=op.get("italic"),
            underline=op.get("underline"),
            color=op.get("color") or op.get("color_rgb"),
            alignment=op.get("alignment"),
            paragraph_spacing=op.get("paragraph_spacing") if op.get("paragraph_spacing") is not None else op.get("space_before"),
            space_after=op.get("space_after"),
            line_spacing=op.get("line_spacing"),
            margins=op.get("margins"),
        )
        results.append({
            "shape_id": sid,
            "success": True,
            "text_summary": res.get("text_summary") or (res.get("text") or "")[:80],
            "font_size": res.get("font_size"),
            "original_font_size": res.get("original_font_size"),
            "resulting_font_size": res.get("resulting_font_size"),
            "font_family": res.get("font_name"),
        })

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "operations_applied": len(results),
        "total_operations": len(operations),
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        "results": results,
    }


@handle_tool_errors
def ppt_scale_slide_typography(
    slide_number: int,
    scale_factor: float = 1.0,
    font_size_delta: Optional[float] = None,
    min_pt: Optional[float] = None,
    max_pt: Optional[float] = None,
    min_font_size: Optional[float] = None,
    max_font_size: Optional[float] = None,
    include_shape_ids: Optional[List[int]] = None,
    exclude_shape_ids: Optional[List[int]] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Proportionally scale or shift typography across all text-bearing shapes on a slide while preserving hierarchy.

    Args:
        slide_number: 1-indexed slide number.
        scale_factor: Scale factor multiplier for font sizes (e.g. 1.15 for +15%, 0.85 for -15%).
        font_size_delta: Point shift added to font sizes (e.g. +2.0, -1.0).
        min_pt / min_font_size: Minimum resulting font size clamp in points.
        max_pt / max_font_size: Maximum resulting font size clamp in points.
        include_shape_ids: Optional list of shape IDs to exclusively scale.
        exclude_shape_ids: Optional list of shape IDs to skip.
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Structured summary detailing modified shapes with old/new sizes, skipped shapes, and reasons.
    """
    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"scale_typography_s{slide_number}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]

    eff_min = min_pt if min_pt is not None else min_font_size
    eff_max = max_pt if max_pt is not None else max_font_size

    res = scale_slide_typography(
        slide,
        scale_factor=scale_factor,
        font_size_delta=font_size_delta,
        min_pt=eff_min,
        max_pt=eff_max,
        include_shape_ids=include_shape_ids,
        exclude_shape_ids=exclude_shape_ids,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        **res,
    }


@handle_tool_errors
def ppt_batch_modify_shapes(
    slide_number: int,
    operations: List[Dict[str, Any]],
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Modify multiple shape geometries (positions, sizes, rotations, z-orders) on a slide in a single transaction.

    Args:
        slide_number: 1-indexed slide number.
        operations: List of shape modification dicts. Each dict supports:
            - shape_id: Target shape ID (int, required)
            - changes: Optional dict with geometry properties, OR geometry properties directly at top level:
                - x, y, width, height: Absolute coordinates/dimensions in inches
                - dx, dy, dwidth, dheight: Relative offsets in inches
                - rotation, drotation: Absolute or relative rotation in degrees
                - z_order: 'bring_to_front', 'send_to_back', 'bring_forward', 'send_backward', or int
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Structured batch summary with per-shape updated coordinates.
    """
    if not operations:
        raise ValueError("Operations list cannot be empty")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"batch_modify_shapes_s{slide_number}_n{len(operations)}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    # Step 1: Pre-validation of all operations before mutating
    for idx, op in enumerate(operations):
        if not isinstance(op, dict):
            raise ValueError(f"Operation at index {idx} must be a dictionary")
        sid = op.get("shape_id")
        if sid is None:
            raise ValueError(f"Operation at index {idx} is missing required 'shape_id'")
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")

    # Step 2: Apply geometry mutations
    results = []
    for op in operations:
        sid = op["shape_id"]
        changes = op.get("changes", {})
        if not isinstance(changes, dict):
            changes = {}

        # Merge changes dict with top-level keys
        merged = {**op, **changes}

        res = modify_shape(
            slide,
            sid,
            x=merged.get("x", merged.get("x_inches")),
            y=merged.get("y", merged.get("y_inches")),
            width=merged.get("width", merged.get("width_inches")),
            height=merged.get("height", merged.get("height_inches")),
            rotation=merged.get("rotation"),
            z_order=merged.get("z_order"),
            dx=merged.get("dx", merged.get("delta_x")),
            dy=merged.get("dy", merged.get("delta_y")),
            dwidth=merged.get("dwidth", merged.get("delta_width")),
            dheight=merged.get("dheight", merged.get("delta_height")),
            drotation=merged.get("drotation", merged.get("delta_rotation")),
        )
        results.append({
            "shape_id": sid,
            "success": True,
            "x": res["x"],
            "y": res["y"],
            "width": res["width"],
            "height": res["height"],
            "rotation": res["rotation"],
        })

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "operations_applied": len(results),
        "total_operations": len(operations),
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        "results": results,
    }


@handle_tool_errors
def ppt_align_shapes(
    slide_number: int,
    shape_ids: List[int],
    alignment: str,
    reference_shape_id: Optional[int] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Align multiple shapes along a common edge or center line without manual arithmetic.

    Args:
        slide_number: 1-indexed slide number.
        shape_ids: List of shape IDs to align (minimum 2).
        alignment: 'left', 'center', 'right', 'top', 'middle', 'bottom' (or 'align_left', etc.).
        reference_shape_id: Optional shape ID to align against (defaults to leftmost/topmost shape).
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing aligned shape IDs and resulting coordinates.
    """
    if len(shape_ids) < 2:
        raise ValueError("At least 2 shape IDs are required for alignment")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"align_shapes_s{slide_number}_{alignment}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    selected = []
    for sid in shape_ids:
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")
        selected.append(shape_map[sid])

    ref_shape = shape_map.get(reference_shape_id) if reference_shape_id else None

    align_shapes(selected, alignment=alignment, reference_shape=ref_shape)

    prs.save(target_path)
    if session:
        session.save_metadata()

    updated = [
        {
            "shape_id": s.shape_id,
            "x": emu_to_inches(int(s.left)),
            "y": emu_to_inches(int(s.top)),
            "width": emu_to_inches(int(s.width)),
            "height": emu_to_inches(int(s.height)),
        }
        for s in selected
    ]

    return {
        "success": True,
        "slide_number": slide_number,
        "operation": f"align_{alignment.replace('align_', '')}",
        "aligned_count": len(selected),
        "shapes": updated,
    }


@handle_tool_errors
def ppt_distribute_shapes(
    slide_number: int,
    shape_ids: List[int],
    direction: str = "horizontal",
    spacing_mode: str = "equal_gaps",
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Distribute shapes evenly across a horizontal or vertical axis.

    Args:
        slide_number: 1-indexed slide number.
        shape_ids: List of shape IDs to distribute (minimum 3).
        direction: 'horizontal' or 'vertical'.
        spacing_mode: 'equal_gaps' (default) or 'equal_centers'.
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing distributed shapes and updated coordinates.
    """
    if len(shape_ids) < 3:
        raise ValueError("At least 3 shape IDs are required for distribution")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"distribute_shapes_s{slide_number}_{direction}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    selected = []
    for sid in shape_ids:
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")
        selected.append(shape_map[sid])

    distributed = distribute_shapes(selected, mode=direction, spacing=spacing_mode)

    prs.save(target_path)
    if session:
        session.save_metadata()

    updated = [
        {
            "shape_id": s.shape_id,
            "x": emu_to_inches(int(s.left)),
            "y": emu_to_inches(int(s.top)),
            "width": emu_to_inches(int(s.width)),
            "height": emu_to_inches(int(s.height)),
        }
        for s in distributed
    ]

    return {
        "success": True,
        "slide_number": slide_number,
        "direction": direction,
        "spacing_mode": spacing_mode,
        "distributed_count": len(distributed),
        "shapes": updated,
    }


@handle_tool_errors
def ppt_space_shapes(
    slide_number: int,
    shape_ids: List[int],
    gap_inches: float,
    direction: str = "horizontal",
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Set an exact fixed spacing gap between consecutive shapes.

    Args:
        slide_number: 1-indexed slide number.
        shape_ids: List of shape IDs to space (minimum 2).
        gap_inches: Fixed gap distance in inches between adjacent shape boundaries.
        direction: 'horizontal' or 'vertical'.
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing spaced shapes and updated positions.
    """
    if len(shape_ids) < 2:
        raise ValueError("At least 2 shape IDs are required for spacing")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"space_shapes_s{slide_number}_{direction}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    selected = []
    for sid in shape_ids:
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")
        selected.append(shape_map[sid])

    spaced = space_shapes(selected, gap_inches=gap_inches, direction=direction)

    prs.save(target_path)
    if session:
        session.save_metadata()

    updated = [
        {
            "shape_id": s.shape_id,
            "x": emu_to_inches(int(s.left)),
            "y": emu_to_inches(int(s.top)),
            "width": emu_to_inches(int(s.width)),
            "height": emu_to_inches(int(s.height)),
        }
        for s in spaced
    ]

    return {
        "success": True,
        "slide_number": slide_number,
        "gap_inches": gap_inches,
        "direction": direction,
        "spaced_count": len(spaced),
        "shapes": updated,
    }


@handle_tool_errors
def ppt_equalize_sizes(
    slide_number: int,
    shape_ids: List[int],
    equalize_width: bool = True,
    equalize_height: bool = True,
    target_width: Optional[float] = None,
    target_height: Optional[float] = None,
    mode: str = "first",
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Equalize widths and/or heights across multiple shapes deterministically.

    Args:
        slide_number: 1-indexed slide number.
        shape_ids: List of shape IDs to equalize.
        equalize_width: Whether to equalize widths (default True).
        equalize_height: Whether to equalize heights (default True).
        target_width: Explicit target width in inches (overrides mode).
        target_height: Explicit target height in inches (overrides mode).
        mode: Strategy for target dimension if not explicitly given ('first', 'max', 'min', 'avg').
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing equalized dimensions and updated shape geometries.
    """
    if len(shape_ids) < 2:
        raise ValueError("At least 2 shape IDs are required for equalization")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"equalize_sizes_s{slide_number}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    selected = []
    for sid in shape_ids:
        if sid not in shape_map:
            raise ValueError(f"Shape with ID {sid} not found on slide {slide_number}")
        selected.append(shape_map[sid])

    equalized = equalize_dimensions(
        selected,
        equalize_width=equalize_width,
        equalize_height=equalize_height,
        target_width_inches=target_width,
        target_height_inches=target_height,
        mode=mode,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    updated = [
        {
            "shape_id": s.shape_id,
            "x": emu_to_inches(int(s.left)),
            "y": emu_to_inches(int(s.top)),
            "width": emu_to_inches(int(s.width)),
            "height": emu_to_inches(int(s.height)),
        }
        for s in equalized
    ]

    return {
        "success": True,
        "slide_number": slide_number,
        "equalized_count": len(equalized),
        "equalized_width": equalize_width,
        "equalized_height": equalize_height,
        "resulting_width": updated[0]["width"] if equalize_width else None,
        "resulting_height": updated[0]["height"] if equalize_height else None,
        "shapes": updated,
    }


@handle_tool_errors
def ppt_move_container(
    slide_number: int,
    container_id: int,
    x: Optional[float] = None,
    y: Optional[float] = None,
    dx: Optional[float] = None,
    dy: Optional[float] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Move a logical container (card) and all its nested child shapes atomically.

    Args:
        slide_number: 1-indexed slide number.
        container_id: ID of container shape.
        x: Absolute destination X coordinate in inches.
        y: Absolute destination Y coordinate in inches.
        dx: Relative horizontal delta shift in inches.
        dy: Relative vertical delta shift in inches.
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing moved container position and all moved children.
    """
    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"move_container_c{container_id}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = move_container(slide, container_id=container_id, x=x, y=y, dx=dx, dy=dy)

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        **res,
    }


@handle_tool_errors
def ppt_resize_container(
    slide_number: int,
    container_id: int,
    width: Optional[float] = None,
    height: Optional[float] = None,
    dwidth: Optional[float] = None,
    dheight: Optional[float] = None,
    scale_width: Optional[float] = None,
    scale_height: Optional[float] = None,
    reflow_children: bool = True,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Resize a logical container/card and adjust child layout proportionally to prevent overflow.

    Args:
        slide_number: 1-indexed slide number.
        container_id: ID of container shape.
        width: Absolute width in inches.
        height: Absolute height in inches.
        dwidth: Relative delta width in inches.
        dheight: Relative delta height in inches.
        scale_width: Width scale multiplier (e.g. 1.2 for +20%).
        scale_height: Height scale multiplier (e.g. 1.1 for +10%).
        reflow_children: Proportionally adjust child bounds (default True).
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing new container dimensions and updated child shapes.
    """
    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"resize_container_c{container_id}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = resize_container(
        slide,
        container_id=container_id,
        width=width,
        height=height,
        dwidth=dwidth,
        dheight=dheight,
        scale_width=scale_width,
        scale_height=scale_height,
        reflow_children=reflow_children,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        **res,
    }


@handle_tool_errors
def ppt_reflow_container(
    slide_number: int,
    container_id: int,
    padding_inches: float = 0.2,
    item_spacing_inches: float = 0.15,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Deterministically stack and organize child elements vertically inside a container with clean padding.

    Args:
        slide_number: 1-indexed slide number.
        container_id: ID of container shape.
        padding_inches: Margin padding inside container edges in inches (default 0.2).
        item_spacing_inches: Vertical spacing gap between items in inches (default 0.15).
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing reflowed child elements and coordinates.
    """
    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"reflow_container_c{container_id}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = reflow_container(
        slide,
        container_id=container_id,
        padding_inches=padding_inches,
        item_spacing_inches=item_spacing_inches,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        **res,
    }


@handle_tool_errors
def ppt_apply_style(
    slide_number: int,
    shape_id: Optional[int] = None,
    shape_ids: Optional[List[int]] = None,
    source_shape_id: Optional[int] = None,
    source_slide_number: Optional[int] = None,
    preset: Optional[str] = None,
    fill_color: Optional[str] = None,
    line_color: Optional[str] = None,
    line_width_pt: Optional[float] = None,
    font_family: Optional[str] = None,
    font_size_pt: Optional[float] = None,
    font_color: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Apply style presets or transfer styles from another shape without re-typing text.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: Target shape ID (or pass list in shape_ids).
        shape_ids: List of target shape IDs to style in batch.
        source_shape_id: Shape ID to copy style from (fill, line, typography).
        source_slide_number: Slide number containing source shape (defaults to same slide).
        preset: Standard design preset name ('card_default', 'card_accent', 'badge_neutral',
                'badge_success', 'badge_warning', 'badge_danger', 'title_hero', 'title_section', 'metric_kpi').
        fill_color: Override fill hex color (e.g. '#F8FAFC').
        line_color: Override border line hex color.
        line_width_pt: Override border line thickness in points.
        font_family: Override font family name.
        font_size_pt: Override font size in points.
        font_color: Override font text hex color.
        bold: Override bold flag.
        italic: Override italic flag.
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing styled shapes and applied attributes.
    """
    target_ids = list(shape_ids) if shape_ids else ([shape_id] if shape_id is not None else [])
    if not target_ids:
        raise ValueError("Either shape_id or shape_ids must be provided")

    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"apply_style_s{slide_number}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    shape_map = {s.shape_id: s for s in slide.shapes}

    # Extract base style from source shape or preset
    style_props: Dict[str, Any] = {}

    if source_shape_id is not None:
        src_slide_num = source_slide_number or slide_number
        if src_slide_num < 1 or src_slide_num > len(prs.slides):
            raise IndexError(f"Source slide number {src_slide_num} is out of range")
        src_slide = prs.slides[src_slide_num - 1]
        src_shape = next((s for s in src_slide.shapes if s.shape_id == source_shape_id), None)
        if not src_shape:
            raise ValueError(f"Source shape ID {source_shape_id} not found on slide {src_slide_num}")
        style_props = extract_complete_shape_style(src_shape)

    elif preset:
        preset_key = preset.strip().lower()
        if preset_key not in STYLE_PRESETS:
            raise ValueError(f"Unknown preset '{preset}'. Available: {list(STYLE_PRESETS.keys())}")
        style_props = dict(STYLE_PRESETS[preset_key])

    # Override with explicitly passed parameters
    effective_fill = fill_color if fill_color is not None else style_props.get("fill_color")
    effective_line = line_color if line_color is not None else style_props.get("line_color")
    effective_line_w = line_width_pt if line_width_pt is not None else style_props.get("line_width_pt")
    effective_font_fam = font_family if font_family is not None else style_props.get("font_family")
    effective_font_sz = font_size_pt if font_size_pt is not None else style_props.get("font_size_pt")
    effective_font_col = font_color if font_color is not None else style_props.get("font_color")
    effective_bold = bold if bold is not None else style_props.get("bold")
    effective_italic = italic if italic is not None else style_props.get("italic")

    results = []
    for sid in target_ids:
        if sid not in shape_map:
            raise ValueError(f"Target shape ID {sid} not found on slide {slide_number}")
        shape = shape_map[sid]
        res = apply_style_to_shape(
            shape,
            fill_color=effective_fill,
            line_color=effective_line,
            line_width_pt=effective_line_w,
            font_family=effective_font_fam,
            font_size_pt=effective_font_sz,
            font_color=effective_font_col,
            bold=effective_bold,
            italic=effective_italic,
        )
        results.append(res)

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "preset_applied": preset,
        "source_shape_id": source_shape_id,
        "shapes_styled_count": len(results),
        "results": results,
    }


@handle_tool_errors
def ppt_create_flow_diagram(
    slide_number: int,
    steps: List[Union[str, Dict[str, Any]]],
    direction: str = "horizontal",
    shape_type: str = "rounded_rectangle",
    start_x: float = 1.0,
    start_y: float = 2.2,
    total_width: Optional[float] = None,
    total_height: Optional[float] = None,
    node_width: Optional[float] = None,
    node_height: Optional[float] = None,
    node_gap: Optional[float] = None,
    style_preset: str = "card_default",
    connector_style: str = "arrow",
    connector_color: str = "#94A3B8",
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Create a structured multi-step flow diagram with connecting arrows and typography.

    Args:
        slide_number: 1-indexed slide number.
        steps: List of step strings (e.g. ['Plan', 'Build', 'Test', 'Deploy']) or dicts with title, description, badge.
        direction: 'horizontal' (default) or 'vertical'.
        shape_type: 'rounded_rectangle', 'rectangle', 'chevron', 'oval'.
        start_x: Diagram origin X position in inches (default 1.0).
        start_y: Diagram origin Y position in inches (default 2.2).
        total_width: Total span width in inches.
        total_height: Total span height in inches.
        node_width: Explicit node width in inches.
        node_height: Explicit node height in inches.
        node_gap: Gap distance between nodes in inches.
        style_preset: Preset style ('card_default', 'card_accent', 'badge_primary').
        connector_style: 'arrow' (default), 'line', or 'none'.
        connector_color: Hex color for connectors (default '#94A3B8').
        presentation_path: Presentation path (defaults to active session).

    Returns:
        Summary detailing created node shapes and connecting arrows.
    """
    target_path, prs, session = _get_target_presentation(
        presentation_path, operation=f"create_flow_diagram_s{slide_number}"
    )

    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number {slide_number} is out of range (1..{len(prs.slides)})")

    slide = prs.slides[slide_number - 1]
    res = create_flow_diagram(
        slide,
        steps=steps,
        direction=direction,
        shape_type=shape_type,
        start_x=start_x,
        start_y=start_y,
        total_width=total_width,
        total_height=total_height,
        node_width=node_width,
        node_height=node_height,
        node_gap=node_gap,
        style_preset=style_preset,
        connector_style=connector_style,
        connector_color=connector_color,
    )

    prs.save(target_path)
    if session:
        session.save_metadata()

    return {
        "success": True,
        "slide_number": slide_number,
        "session_id": session.session_id if session else None,
        "target": "working" if session else "standalone",
        **res,
    }


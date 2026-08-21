"""Comprehensive test suite for rule-based slide and presentation validation (VAL-01 to VAL-07)."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import (
    BoundingBox,
    ParagraphModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    TextFrameModel,
    TextRunModel,
    TextStyle,
)
from powerpoint_mcp.models.slide import SlideModel
from powerpoint_mcp.pptx.inspector import inspect_presentation, inspect_slide
from powerpoint_mcp.utils.validation import (
    IssueSeverity,
    PresentationValidationResult,
    SlideIssue,
    SlideValidationResult,
    validate_presentation,
    validate_slide,
)


class TestSyntheticDeckValidation:
    """Test validation engine on the synthetic 3-slide presentation deck."""

    def test_slide_1_clean_performance_overview(self, synthetic_deck_path: Path):
        """Slide 1 is a clean, well-formatted dashboard layout with 0 validation issues."""
        slide1_model = inspect_slide(str(synthetic_deck_path), 1)
        result = validate_slide(slide1_model, slide_number=1)

        assert isinstance(result, SlideValidationResult)
        assert result.slide_number == 1
        assert result.is_valid is True
        assert result.warning_count == 0
        assert result.error_count == 0
        assert len(result.issues) == 0
        assert result.metrics["shape_count"] >= 4
        assert result.metrics["image_shape_count"] >= 1

    def test_slide_2_operational_architecture_clean(self, synthetic_deck_path: Path):
        """Slide 2 is a clean 2-column + chevron diagram layout with 0 critical errors."""
        slide2_model = inspect_slide(str(synthetic_deck_path), 2)
        result = validate_slide(slide2_model, slide_number=2)

        assert result.slide_number == 2
        assert result.error_count == 0
        # Should have 0 overlap / clipping / overflow defects
        overlap_issues = [i for i in result.issues if i.rule_id == "VAL-01"]
        clipping_issues = [i for i in result.issues if i.rule_id == "VAL-02"]
        overflow_issues = [i for i in result.issues if i.rule_id == "VAL-03"]
        assert len(overlap_issues) == 0
        assert len(clipping_issues) == 0
        assert len(overflow_issues) == 0

    def test_slide_3_detects_all_four_intentional_defects(self, synthetic_deck_path: Path):
        """Slide 3 contains 4 intentional defects: VAL-01 (overlap), VAL-02 (clipping), VAL-03 (overflow), VAL-04 (tiny font)."""
        slide3_model = inspect_slide(str(synthetic_deck_path), 3)
        result = validate_slide(slide3_model, slide_number=3)

        assert result.slide_number == 3
        assert result.is_valid is False
        assert result.warning_count >= 4

        rule_ids = {i.rule_id for i in result.issues}

        # 1. Defect 1: VAL-01 Heavy Overlap between Box A and Box B (Area = 3.0 sq in > 0.5 sq in)
        assert "VAL-01" in rule_ids
        val01_issues = [i for i in result.issues if i.rule_id == "VAL-01"]
        assert len(val01_issues) >= 1
        overlap_details = val01_issues[0].details
        assert overlap_details["overlap_area_sq_in"] >= 0.5
        assert len(val01_issues[0].shape_ids) == 2

        # 2. Defect 2: VAL-02 Boundary clipping on Box C (right edge extends past 13.333 in limit)
        assert "VAL-02" in rule_ids
        val02_issues = [i for i in result.issues if i.rule_id == "VAL-02"]
        assert len(val02_issues) >= 1
        assert "right" in val02_issues[0].details["protrusions"]
        assert val02_issues[0].details["protrusions"]["right"] > 1.0

        # 3. Defect 3: VAL-04 Suspiciously tiny font on Box D (5.5 pt < 8.0 pt threshold)
        assert "VAL-04" in rule_ids
        val04_issues = [i for i in result.issues if i.rule_id == "VAL-04"]
        assert len(val04_issues) >= 1
        assert val04_issues[0].details["font_size_pt"] == 5.5
        assert "LEGAL NOTICE" in val04_issues[0].details["text_preview"]

        # 4. Defect 4: VAL-03 Text overflow on Box E (large high-point sentence overflowing container)
        assert "VAL-03" in rule_ids
        val03_issues = [i for i in result.issues if i.rule_id == "VAL-03"]
        assert len(val03_issues) >= 1
        assert val03_issues[0].details["overflow_ratio"] > 1.15


class TestIndividualValidationRules:
    """Test each validation rule with synthetic shape models in isolation."""

    def test_val_01_overlap_detection_and_background_exclusion(self):
        """VAL-01 detects overlapping shapes and ignores full-slide background rectangles."""
        # Full canvas background shape
        bg_shape = ShapeModel(
            shape_id=1,
            name="Slide Background",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.UNKNOWN,
            bbox=BoundingBox.from_inches(0.0, 0.0, 13.333, 7.5),
            z_order=0,
        )
        # Shape 1
        box1 = ShapeModel(
            shape_id=2,
            name="Box 1",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(2.0, 2.0, 3.0, 2.0),
        )
        # Shape 2 (Overlapping Box 1 by 1.0 x 1.0 = 1.0 sq in)
        box2 = ShapeModel(
            shape_id=3,
            name="Box 2",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(4.0, 3.0, 3.0, 2.0),
        )
        # Shape 3 (Non-overlapping)
        box3 = ShapeModel(
            shape_id=4,
            name="Box 3",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(8.0, 2.0, 2.0, 2.0),
        )

        slide = SlideModel(
            slide_number=1,
            slide_id=10,
            width_inches=13.333,
            height_inches=7.5,
            shapes=[bg_shape, box1, box2, box3],
        )

        res = validate_slide(slide, rules=["VAL-01"])
        assert len(res.issues) == 1
        issue = res.issues[0]
        assert issue.rule_id == "VAL-01"
        assert set(issue.shape_ids) == {2, 3}
        assert issue.details["overlap_area_sq_in"] == 1.0

    def test_val_02_off_slide_boundary_clipping(self):
        """VAL-02 detects shapes protruding across left, top, right, and bottom boundaries."""
        shape_left_clip = ShapeModel(
            shape_id=1,
            name="Left Clip",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(-0.5, 2.0, 2.0, 2.0),
        )
        shape_bottom_clip = ShapeModel(
            shape_id=2,
            name="Bottom Clip",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(2.0, 7.0, 2.0, 2.0),
        )
        shape_in_bounds = ShapeModel(
            shape_id=3,
            name="In Bounds",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(2.0, 2.0, 3.0, 2.0),
        )

        slide = SlideModel(
            slide_number=1,
            slide_id=10,
            width_inches=13.333,
            height_inches=7.5,
            shapes=[shape_left_clip, shape_bottom_clip, shape_in_bounds],
        )

        res = validate_slide(slide, rules=["VAL-02"])
        assert len(res.issues) == 2
        shape_ids_clipped = {i.shape_ids[0] for i in res.issues}
        assert shape_ids_clipped == {1, 2}

    def test_val_03_text_overflow_detection(self):
        """VAL-03 calculates text height vs frame height and flags overflow condition."""
        # Overflowing text shape: 1000 characters of 18pt text in a tiny 2x1 inch box
        overflow_runs = [TextRunModel(text="A" * 600, style=TextStyle(font_size_pt=18.0))]
        overflow_para = ParagraphModel(text="A" * 600, runs=overflow_runs)
        overflow_tf = TextFrameModel(text="A" * 600, paragraphs=[overflow_para])

        overflow_shape = ShapeModel(
            shape_id=1,
            name="Overflow Shape",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(1.0, 1.0, 2.0, 1.0),
            text_frame=overflow_tf,
        )

        # Non-overflowing normal shape
        normal_runs = [TextRunModel(text="Short title", style=TextStyle(font_size_pt=14.0))]
        normal_para = ParagraphModel(text="Short title", runs=normal_runs)
        normal_tf = TextFrameModel(text="Short title", paragraphs=[normal_para])
        normal_shape = ShapeModel(
            shape_id=2,
            name="Normal Shape",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(5.0, 1.0, 4.0, 2.0),
            text_frame=normal_tf,
        )

        slide = SlideModel(
            slide_number=1,
            slide_id=10,
            shapes=[overflow_shape, normal_shape],
        )

        res = validate_slide(slide, rules=["VAL-03"])
        assert len(res.issues) == 1
        assert res.issues[0].rule_id == "VAL-03"
        assert res.issues[0].shape_ids == [1]

    def test_val_04_tiny_font_detection(self):
        """VAL-04 detects font sizes below 8.0 pt."""
        tiny_run = TextRunModel(text="Tiny disclaimer", style=TextStyle(font_size_pt=6.0))
        tiny_para = ParagraphModel(text="Tiny disclaimer", runs=[tiny_run])
        tiny_tf = TextFrameModel(text="Tiny disclaimer", paragraphs=[tiny_para])
        tiny_shape = ShapeModel(
            shape_id=5,
            name="Disclaimer Box",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.FOOTER,
            bbox=BoundingBox.from_inches(1.0, 6.0, 4.0, 1.0),
            text_frame=tiny_tf,
        )

        slide = SlideModel(slide_number=1, slide_id=10, shapes=[tiny_shape])
        res = validate_slide(slide, rules=["VAL-04"])

        assert len(res.issues) == 1
        assert res.issues[0].rule_id == "VAL-04"
        assert res.issues[0].details["font_size_pt"] == 6.0

    def test_val_05_inconsistent_title_position(self):
        """VAL-05 detects title positions deviating from presentation baseline."""
        title_shape = ShapeModel(
            shape_id=1,
            name="Title 1",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.TITLE,
            bbox=BoundingBox.from_inches(2.5, 1.5, 8.0, 1.0),
        )
        slide = SlideModel(slide_number=2, slide_id=11, shapes=[title_shape])

        # Baseline title position on slide 1: x=1.0, y=0.8
        baseline = {"x": 1.0, "y": 0.8, "width": 11.333, "height": 0.9}
        res = validate_slide(slide, rules=["VAL-05"], title_baseline=baseline)

        assert len(res.issues) == 1
        assert res.issues[0].rule_id == "VAL-05"
        assert res.issues[0].details["delta_x_inches"] == 1.5
        assert res.issues[0].details["delta_y_inches"] == 0.7

    def test_val_06_duplicate_objects(self):
        """VAL-06 detects duplicate superimposed objects with matching coordinates and content."""
        box1 = ShapeModel(
            shape_id=10,
            name="Button A",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(3.0, 3.0, 2.0, 1.0),
            text_frame=TextFrameModel(text="Click Here"),
        )
        box2 = ShapeModel(
            shape_id=11,
            name="Button A Copy",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(3.0, 3.0, 2.0, 1.0),
            text_frame=TextFrameModel(text="Click Here"),
        )

        slide = SlideModel(slide_number=1, slide_id=10, shapes=[box1, box2])
        res = validate_slide(slide, rules=["VAL-06"])

        assert len(res.issues) == 1
        assert res.issues[0].rule_id == "VAL-06"
        assert set(res.issues[0].shape_ids) == {10, 11}

    def test_val_07_extreme_rotations(self):
        """VAL-07 detects shapes with non-standard / irregular rotation angles."""
        rotated_shape = ShapeModel(
            shape_id=20,
            name="Tilted Banner",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(4.0, 4.0, 3.0, 1.0),
            rotation=23.5,
        )
        normal_rot_shape = ShapeModel(
            shape_id=21,
            name="Standard 90 Banner",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(1.0, 1.0, 3.0, 1.0),
            rotation=90.0,
        )

        slide = SlideModel(slide_number=1, slide_id=10, shapes=[rotated_shape, normal_rot_shape])
        res = validate_slide(slide, rules=["VAL-07"])

        assert len(res.issues) == 1
        assert res.issues[0].rule_id == "VAL-07"
        assert res.issues[0].shape_ids == [20]
        assert res.issues[0].details["rotation_degrees"] == 23.5


class TestPresentationValidationAndSerialization:
    """Test full presentation validation and JSON schema compliance."""

    def test_validate_full_presentation(self, synthetic_deck_path: Path):
        """validate_presentation aggregates results across all slides and detects issues on slide 3."""
        prs_model = inspect_presentation(str(synthetic_deck_path))
        pres_result = validate_presentation(prs_model)

        assert isinstance(pres_result, PresentationValidationResult)
        assert pres_result.slide_count == 3
        assert pres_result.is_valid is False
        assert pres_result.total_warnings >= 4
        assert len(pres_result.slide_results) == 3

        # Slide 1 clean
        assert pres_result.slide_results[0].is_valid is True
        # Slide 3 has defects
        assert pres_result.slide_results[2].is_valid is False

    def test_slide_validation_result_serialization(self, synthetic_deck_path: Path):
        """to_dict() matches the SlideValidationReport schema."""
        slide3_model = inspect_slide(str(synthetic_deck_path), 3)
        res = validate_slide(slide3_model, slide_number=3)
        data = res.to_dict()

        assert "slide_number" in data
        assert "is_valid" in data
        assert "warning_count" in data
        assert "error_count" in data
        assert "warnings" in data
        assert "metrics" in data
        assert isinstance(data["warnings"], list)
        assert len(data["warnings"]) >= 4

        first_warning = data["warnings"][0]
        assert "rule_id" in first_warning
        assert "severity" in first_warning
        assert "shape_ids" in first_warning
        assert "message" in first_warning
        assert "details" in first_warning

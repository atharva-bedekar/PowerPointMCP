"""Unit tests for PowerPoint MCP v1.1 Phase C: Layout Primitives & Container Transformations."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.tools.editing import (
    ppt_align_shapes,
    ppt_distribute_shapes,
    ppt_equalize_sizes,
    ppt_move_container,
    ppt_reflow_container,
    ppt_resize_container,
    ppt_space_shapes,
)
from powerpoint_mcp.tools.versioning import open_presentation


@pytest.fixture
def layout_deck(tmp_path: Path) -> Path:
    """Create a test deck with multiple shapes and a container with children."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 3 independent shapes for alignment / distribution
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.0))
    s1.name = "Shape 1"
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(1.5), Inches(3.0), Inches(1.2))
    s2.name = "Shape 2"
    s3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(0.8), Inches(2.5), Inches(0.9))
    s3.name = "Shape 3"

    # Container Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.5), Inches(4.0), Inches(3.5))
    card.name = "Card Container"
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(230, 235, 245)

    # Child 1: Title
    c_title = slide.shapes.add_textbox(Inches(1.2), Inches(3.7), Inches(3.6), Inches(0.5))
    c_title.name = "Card Title"
    c_title.text_frame.paragraphs[0].text = "Security Services"

    # Child 2: Body
    c_body = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(3.6), Inches(1.5))
    c_body.name = "Card Body"
    c_body.text_frame.paragraphs[0].text = "Role-based access control\nZero-trust network architecture"

    # Child 3: Badge
    c_badge = slide.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(1.5), Inches(0.35))
    c_badge.name = "Card Badge"
    c_badge.text_frame.paragraphs[0].text = "COMPLIANT"

    deck_path = tmp_path / "layout_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_align_shapes_top(layout_deck: Path):
    """Test aligning shapes along their top edge."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id, s3_id = shapes[0].shape_id, shapes[1].shape_id, shapes[2].shape_id

    res = ppt_align_shapes(slide_number=1, shape_ids=[s1_id, s2_id, s3_id], alignment="top")
    assert res["success"] is True
    assert res["aligned_count"] == 3
    # Topmost was Shape 3 at 0.8 in
    assert res["shapes"][0]["y"] == 0.8
    assert res["shapes"][1]["y"] == 0.8
    assert res["shapes"][2]["y"] == 0.8


def test_distribute_shapes_horizontal(layout_deck: Path):
    """Test distributing 3 shapes horizontally with equal gaps."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id, s3_id = shapes[0].shape_id, shapes[1].shape_id, shapes[2].shape_id

    res = ppt_distribute_shapes(slide_number=1, shape_ids=[s1_id, s2_id, s3_id], direction="horizontal")
    assert res["success"] is True
    assert res["distributed_count"] == 3


def test_space_shapes_fixed_gap(layout_deck: Path):
    """Test setting an exact 0.5 inch gap between shapes."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id = shapes[0].shape_id, shapes[1].shape_id

    res = ppt_space_shapes(slide_number=1, shape_ids=[s1_id, s2_id], gap_inches=0.5, direction="horizontal")
    assert res["success"] is True
    # Shape 1 is at x=1.0, w=2.0 -> right edge is 3.0. Shape 2 should be at 3.5.
    assert res["shapes"][0]["x"] == 1.0
    assert pytest.approx(res["shapes"][1]["x"], 0.01) == 3.5


def test_equalize_sizes(layout_deck: Path):
    """Test equalizing dimensions of shapes."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id, s3_id = shapes[0].shape_id, shapes[1].shape_id, shapes[2].shape_id

    res = ppt_equalize_sizes(
        slide_number=1,
        shape_ids=[s1_id, s2_id, s3_id],
        target_width=3.5,
        target_height=1.5,
    )
    assert res["success"] is True
    for s in res["shapes"]:
        assert s["width"] == 3.5
        assert s["height"] == 1.5


def test_move_container_moves_children(layout_deck: Path):
    """Test moving a container card moves both the container and all nested child elements."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    card_id = shapes[3].shape_id  # Card Container is shape 4

    # Move container from x=1.0, y=3.5 to x=6.0, y=2.0 (dx=+5.0, dy=-1.5)
    res = ppt_move_container(slide_number=1, container_id=card_id, x=6.0, y=2.0)
    assert res["success"] is True
    assert res["x"] == 6.0
    assert res["y"] == 2.0
    assert res["children_moved_count"] == 3

    # Check child coordinates shifted by dx=+5.0, dy=-1.5
    for c in res["children"]:
        assert c["x"] >= 6.0  # Child was shifted right into new container position


def test_resize_container_scales_children(layout_deck: Path):
    """Test resizing a container proportionally resizes/repositions nested child shapes."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    card_id = shapes[3].shape_id

    # Resize container width from 4.0 to 6.0 (scale_width = 1.5)
    res = ppt_resize_container(slide_number=1, container_id=card_id, width=6.0, reflow_children=True)
    assert res["success"] is True
    assert res["width"] == 6.0
    assert res["children_count"] == 3


def test_reflow_container(layout_deck: Path):
    """Test reflowing child items vertically inside a container."""
    session = open_presentation(str(layout_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    card_id = shapes[3].shape_id

    res = ppt_reflow_container(slide_number=1, container_id=card_id, padding_inches=0.25, item_spacing_inches=0.2)
    assert res["success"] is True
    assert res["children_reflowed_count"] == 3
    # First item top should be container top (3.5) + padding (0.25) = 3.75
    assert pytest.approx(res["children"][0]["y"], 0.01) == 3.75

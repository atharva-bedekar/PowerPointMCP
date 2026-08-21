"""Adversarial stress-testing script for M1 Inspection Engine."""

import os
from pathlib import Path
import sys

SRC_DIR = Path(__file__).resolve().parent.parent.parent / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    inches_to_emu,
    emu_to_inches,
)
from powerpoint_mcp.pptx import (
    PPTXInspector,
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    match_shapes,
)

def test_stress_extreme_coordinates():
    print("Stress 1: Extreme Coordinates")
    # Coordinates in billions of EMUs (e.g. 50,000 inches)
    huge_emu = 50000 * 914400
    bb = BoundingBox(left_emu=huge_emu, top_emu=huge_emu, width_emu=huge_emu, height_emu=huge_emu)
    assert bb.left_inches == 50000.0
    assert bb.width_inches == 50000.0
    assert bb.right_emu == huge_emu * 2
    d = bb.to_dict()
    assert d["left_inches"] == 50000.0

def test_stress_unicode_and_emojis():
    print("Stress 2: Unicode and Emojis")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb.text_frame.text = "🚀 Q3 Results — 日本語, العربية, and Greek: Σ(α+β) = γ 🎉"
    
    slide_model = inspect_slide(prs, 1)
    assert "🚀 Q3 Results" in slide_model.shapes[0].text_frame.text
    assert "日本語" in slide_model.shapes[0].text_frame.text
    assert "العربية" in slide_model.shapes[0].text_frame.text

def test_stress_zero_and_negative_dimensions():
    print("Stress 3: Zero and Negative Dimensions")
    bb_zero = BoundingBox(left_emu=0, top_emu=0, width_emu=0, height_emu=0)
    assert bb_zero.width_inches == 0.0
    assert bb_zero.height_inches == 0.0
    assert bb_zero.center_x_inches == 0.0

    bb_neg = BoundingBox(left_emu=-914400, top_emu=-1828800, width_emu=914400, height_emu=1828800)
    assert bb_neg.left_inches == -1.0
    assert bb_neg.top_inches == -2.0
    assert bb_neg.right_inches == 0.0
    assert bb_neg.bottom_inches == 0.0

def test_stress_empty_presentation():
    print("Stress 4: Empty Presentation")
    prs = Presentation()
    prs_model = inspect_presentation(prs)
    assert prs_model.slide_count == 0
    assert len(prs_model.slides) == 0
    assert len(prs_model.slide_titles) == 0

def test_stress_bipartite_matching_scale():
    print("Stress 5: Bipartite Matching with 50 shapes")
    shapes_a = []
    shapes_b = []
    for i in range(50):
        sa = ShapeModel(
            shape_id=i + 1,
            name=f"Element_{i}",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(float(i % 10), float(i // 10), 0.8, 0.8),
        )
        sb = ShapeModel(
            shape_id=1000 + i + 1,
            name=f"Element_{i}",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(float(i % 10), float(i // 10), 0.8, 0.8),
        )
        shapes_a.append(sa)
        shapes_b.append(sb)

    sl_a = SlideModel(slide_number=1, slide_id=1, shapes=shapes_a)
    sl_b = SlideModel(slide_number=2, slide_id=2, shapes=shapes_b)

    matches = match_shapes(sl_a, sl_b, min_confidence=0.40)
    assert len(matches) == 50
    # Check that each element matched its exact counterpart
    for m in matches:
        assert m["shape_b_id"] == 1000 + m["shape_a_id"]
        assert m["confidence_score"] >= 0.90

def main():
    test_stress_extreme_coordinates()
    test_stress_unicode_and_emojis()
    test_stress_zero_and_negative_dimensions()
    test_stress_empty_presentation()
    test_stress_bipartite_matching_scale()
    print("\nALL ADVERSARIAL STRESS TESTS PASSED CLEANLY")

if __name__ == "__main__":
    main()

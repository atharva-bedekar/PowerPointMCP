"""Inspection tools for PowerPoint MCP server."""

from functools import wraps
import os
from pathlib import Path
import traceback
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation

from powerpoint_mcp.models.shape import ShapeModel, emu_to_inches
from powerpoint_mcp.models.slide import SlideModel
from powerpoint_mcp.models.presentation import PresentationModel
from powerpoint_mcp.pptx.components import inspect_components
from powerpoint_mcp.pptx.cross_slide import compare_cross_slides
from powerpoint_mcp.pptx.inspector import (
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    match_shapes,
)
from powerpoint_mcp.pptx.structure import (
    analyze_containers,
    analyze_slide_structure,
)
from powerpoint_mcp.rendering.visual_compare import compare_slides
from powerpoint_mcp.tools.versioning import get_session_manager, resolve_active_target
from powerpoint_mcp.utils.validation import validate_slide


def handle_tool_errors(func):
    """Decorator to catch exceptions and return structured JSON error payloads."""
    @wraps(func)
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except FileNotFoundError as exc:
            return {
                "success": False,
                "error_type": "FileNotFound",
                "message": str(exc),
                "details": {"path": str(kwargs.get("presentation_path", ""))},
            }
        except IndexError as exc:
            return {
                "success": False,
                "error_type": "SlideNotFound",
                "message": str(exc),
                "details": {"error": str(exc)},
            }
        except ValueError as exc:
            msg = str(exc)
            err_type = "ValidationError"
            if "Shape with ID" in msg or "Shape ID" in msg:
                err_type = "ShapeNotFound"
            elif "Slide number" in msg or "Slide" in msg and "out of range" in msg:
                err_type = "SlideNotFound"
            elif "Session not found" in msg or "active" in msg.lower():
                err_type = "SessionNotFound"
            elif "coordinate" in msg.lower() or "dimension" in msg.lower():
                err_type = "InvalidCoordinate"
            elif "text frame" in msg.lower():
                err_type = "InvalidTextFrame"
            elif "xml" in msg.lower():
                err_type = "OOXMLValidationError"
            return {
                "success": False,
                "error_type": err_type,
                "message": msg,
                "details": {"error": msg},
            }
        except Exception as exc:
            return {
                "success": False,
                "error_type": type(exc).__name__,
                "message": str(exc),
                "details": {"traceback": traceback.format_exc()},
            }
    return wrapper


def _resolve_presentation_path(presentation_path: Optional[str] = None) -> str:
    """Resolve presentation path from argument or active session using resolve_active_target."""
    target_path, _ = resolve_active_target(
        presentation_path=presentation_path,
        require_session=False,
        mutation=False,
    )
    return target_path


@handle_tool_errors
def ppt_inspect_presentation(
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Inspect presentation metadata, slide count, dimensions, layouts, and slide summaries.

    Args:
        presentation_path: Path to presentation. If omitted, uses the active session's working copy.

    Returns:
        Structured dictionary detailing presentation dimensions, slide count, layouts, and titles.
    """
    target_path = _resolve_presentation_path(presentation_path)
    model: PresentationModel = inspect_presentation(target_path)

    return {
        "success": True,
        "path": model.path,
        "slide_count": model.slide_count,
        "width_inches": model.width_inches,
        "height_inches": model.height_inches,
        "slide_width_emu": model.width_emu,
        "slide_height_emu": model.height_emu,
        "layouts": model.layouts,
        "theme": model.theme_name or "Default",
        "metadata": model.metadata.to_dict() if model.metadata else {},
        "titles": [s.title for s in model.slides if s.title],
        "slides": model.slide_titles if model.slide_titles else [
            {
                "slide_number": s.slide_number,
                "slide_id": s.slide_id,
                "title": s.title,
                "layout_name": s.layout_name,
                "shape_count": s.shape_count,
            }
            for s in model.slides
        ],
    }


@handle_tool_errors
def ppt_inspect_slide(
    slide_number: int,
    presentation_path: Optional[str] = None,
    detail: str = "summary",
    text_only: bool = False,
    include_geometry: bool = True,
    include_style: bool = True,
    include_xml: bool = False,
    include_images: bool = True,
    shape_types: Optional[List[str]] = None,
    semantic_roles: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Inspect shapes on a specific slide with filtering and configurable detail levels.

    Args:
        slide_number: 1-indexed slide number.
        presentation_path: Path to presentation. If omitted, uses active session.
        detail: 'summary' (default, agent-friendly concise output) or 'full' (exhaustive shape tree).
        text_only: If True, return only shapes containing text.
        include_geometry: Whether to include coordinates/dimensions.
        include_style: Whether to include colors/font properties.
        include_xml: Whether to include raw XML snippets (only in full mode).
        include_images: Whether to include image/picture shapes.
        shape_types: Optional list of shape types to filter (e.g. ['auto_shape', 'text_box']).
        semantic_roles: Optional list of semantic roles to filter (e.g. ['title', 'body']).

    Returns:
        Structured dictionary containing slide metadata and filtered shape collections.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    slide_model: SlideModel = inspect_slide(target_path, slide_number)

    filtered_shapes = []
    normalized_types = [t.strip().lower() for t in shape_types] if shape_types else None
    normalized_roles = [r.strip().lower() for r in semantic_roles] if semantic_roles else None

    for s in slide_model.shapes:
        # Filter: text_only
        has_text = bool(s.text_frame and s.text_frame.text and s.text_frame.text.strip())
        if text_only and not has_text:
            continue

        # Filter: include_images
        is_image = s.shape_type.value in ("picture", "image") or s.image_metadata is not None
        if not include_images and is_image:
            continue

        # Filter: shape_types
        if normalized_types and s.shape_type.value.lower() not in normalized_types:
            continue

        # Filter: semantic_roles
        if normalized_roles and s.semantic_role.value.lower() not in normalized_roles:
            continue

        # Format shape according to detail level and inclusion flags
        if str(detail).lower() == "full":
            d = s.to_dict()
            if not include_xml:
                d.pop("raw_xml", None)
            if not include_geometry:
                for k in ("bbox", "x", "y", "width", "height", "right", "bottom", "rotation"):
                    d.pop(k, None)
            if not include_style:
                for k in (
                    "fill", "fill_color", "line", "line_color", "font_family", "font_name",
                    "font_size", "font_size_pt", "bold", "italic", "underline", "color",
                    "color_rgb", "alignment"
                ):
                    d.pop(k, None)
        else:
            d = s.to_summary_dict()
            if not include_geometry:
                for k in ("bbox", "x", "y", "width", "height", "right", "bottom", "rotation"):
                    d.pop(k, None)
            if not include_style:
                for k in (
                    "fill_color", "line_color", "line_width_pt", "font_family", "font_name",
                    "font_size", "font_size_pt", "bold", "color", "color_rgb", "alignment"
                ):
                    d.pop(k, None)

        filtered_shapes.append(d)

    return {
        "success": True,
        "slide_number": slide_model.slide_number,
        "slide_id": slide_model.slide_id,
        "layout_name": slide_model.layout_name,
        "title": slide_model.title,
        "total_shape_count": slide_model.shape_count,
        "shape_count": len(filtered_shapes),
        "width_inches": slide_model.width_inches,
        "height_inches": slide_model.height_inches,
        "width_emu": slide_model.width_emu,
        "height_emu": slide_model.height_emu,
        "has_notes": slide_model.has_notes,
        "notes": slide_model.notes,
        "detail": detail,
        "filters": {
            "text_only": text_only,
            "include_geometry": include_geometry,
            "include_style": include_style,
            "include_images": include_images,
            "shape_types": shape_types,
            "semantic_roles": semantic_roles,
        },
        "shapes": filtered_shapes,
    }


@handle_tool_errors
def ppt_inspect_text(
    slide_number: int,
    include_geometry: bool = True,
    include_style: bool = True,
    include_paragraph_metadata: bool = False,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Efficiently inspect all text-bearing shapes on a slide without full shape-tree overhead.

    Args:
        slide_number: 1-indexed slide number.
        include_geometry: Whether to include coordinates and bounding box dimensions.
        include_style: Whether to include typography details (font, size, weight, color).
        include_paragraph_metadata: Whether to include per-paragraph level and bullet metadata.
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured list of all text objects with IDs, semantic roles, text, and styles.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    slide_model: SlideModel = inspect_slide(target_path, slide_number)

    text_shapes_out = []
    for s in slide_model.shapes:
        if not s.text_frame or not s.text_frame.text or not s.text_frame.text.strip():
            continue

        raw_text = s.text_frame.text.strip()
        item: Dict[str, Any] = {
            "shape_id": s.shape_id,
            "name": s.name,
            "semantic_role": s.semantic_role.value,
            "text": raw_text,
        }

        if include_geometry:
            item["bbox"] = {
                "x": s.bbox.left_inches,
                "y": s.bbox.top_inches,
                "width": s.bbox.width_inches,
                "height": s.bbox.height_inches,
            }
            item["x"] = s.bbox.left_inches
            item["y"] = s.bbox.top_inches
            item["width"] = s.bbox.width_inches
            item["height"] = s.bbox.height_inches

        if include_style:
            first_style = None
            alignment = None
            for p in s.text_frame.paragraphs:
                if p.runs:
                    first_style = p.runs[0].style
                    alignment = p.alignment
                    break

            item["font_family"] = first_style.font_name if (first_style and first_style.font_name) else "Default"
            item["font_size"] = first_style.font_size_pt if (first_style and first_style.font_size_pt is not None) else None
            item["bold"] = first_style.bold if first_style else False
            item["italic"] = first_style.italic if first_style else False
            item["color"] = first_style.color_rgb if (first_style and first_style.color_rgb) else None
            item["alignment"] = alignment or "left"

        if include_paragraph_metadata:
            paras = []
            for p in s.text_frame.paragraphs:
                paras.append({
                    "level": p.level,
                    "alignment": p.alignment or "left",
                    "text": p.text,
                    "runs_count": len(p.runs),
                })
            item["paragraph_count"] = len(paras)
            item["paragraphs"] = paras

        # Rough overflow estimation: typical character area in pt^2 vs box area in sq pt
        if include_geometry and include_style and item.get("font_size"):
            f_size = item["font_size"]
            box_area_sq_pt = (s.bbox.width_inches * 72) * (s.bbox.height_inches * 72)
            est_text_area_sq_pt = len(raw_text) * (f_size * 0.6) * (f_size * 1.2)
            if est_text_area_sq_pt > box_area_sq_pt * 1.25 and len(raw_text) > 30:
                item["overflow_warning"] = True

        text_shapes_out.append(item)

    return {
        "success": True,
        "slide_number": slide_model.slide_number,
        "slide_title": slide_model.title,
        "total_slide_shapes": slide_model.shape_count,
        "text_shape_count": len(text_shapes_out),
        "shapes": text_shapes_out,
    }




@handle_tool_errors
def ppt_inspect_shape(
    slide_number: int,
    shape_id: int,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Get deep inspection details for a single shape on a slide.

    Args:
        slide_number: 1-indexed slide number.
        shape_id: ID of the shape to inspect.
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Complete shape details with bounding box, text frames, paragraphs, runs, fills, and lines.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    shape_model: ShapeModel = inspect_shape(target_path, slide_number, shape_id)

    return {
        "success": True,
        "slide_number": slide_number,
        "shape_id": shape_id,
        "shape": shape_model.to_dict(),
    }


@handle_tool_errors
def ppt_inspect_components(
    slide_number: int,
    detail: str = "summary",
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Inspect semantic visual components (headers, footers, steppers, cards, content containers) on a slide.

    Args:
        slide_number: 1-indexed slide number.
        detail: 'summary' (default concise component overview) or 'full' (with child shape details).
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured dictionary detailing detected components, bounding boxes, properties, and constituent shape IDs.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    return inspect_components(target_path, slide_number, detail=detail)


@handle_tool_errors
def ppt_compare_slides(
    slide_a: Optional[int] = None,
    slide_b: Optional[int] = None,
    reference_slide: Optional[int] = None,
    target_slides: Optional[List[int]] = None,
    aspects: Optional[List[str]] = None,
    match_shapes_flag: bool = True,
    render_diff: bool = False,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Compare geometric, typographic, and semantic layout properties between slides.

    Supports both cross-slide multi-target comparison against a reference slide (v1.2)
    and legacy two-slide comparison (v1.1).

    Args:
        slide_a: Legacy 1-indexed reference slide number.
        slide_b: Legacy 1-indexed comparison slide number.
        reference_slide: 1-indexed reference slide number for multi-slide comparison.
        target_slides: List of 1-indexed target slide numbers to compare against reference slide.
        aspects: List of comparison aspects (['components', 'geometry', 'typography', 'colors', 'spacing']).
        match_shapes_flag: Whether to perform multi-factor semantic shape matching (for 2-slide mode).
        render_diff: Whether to perform pixel visual diffing if renderers are available.
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Slide comparison report detailing matched shapes/components, layout differences, and summary text.
    """
    target_path = _resolve_presentation_path(presentation_path)

    # Cross-slide mode: if reference_slide and target_slides provided, or target_slides provided
    if reference_slide is not None and target_slides is not None:
        if reference_slide < 1:
            raise IndexError(f"Reference slide number must be >= 1, got {reference_slide}")
        for t in target_slides:
            if t < 1:
                raise IndexError(f"Target slide number must be >= 1, got {t}")
        return compare_cross_slides(target_path, reference_slide, target_slides, aspects=aspects)

    # If slide_a is provided and target_slides is provided
    if slide_a is not None and target_slides is not None:
        if slide_a < 1:
            raise IndexError(f"Reference slide number must be >= 1, got {slide_a}")
        return compare_cross_slides(target_path, slide_a, target_slides, aspects=aspects)

    # Legacy two-slide mode
    ref_s = reference_slide if reference_slide is not None else slide_a
    tgt_s = slide_b

    if ref_s is None or tgt_s is None:
        raise ValueError("Either specify (reference_slide, target_slides) or (slide_a, slide_b)")

    if ref_s < 1 or tgt_s < 1:
        raise IndexError(f"Slide numbers must be >= 1, got ref_s={ref_s}, tgt_s={tgt_s}")

    slide_a_model = inspect_slide(target_path, ref_s)
    slide_b_model = inspect_slide(target_path, tgt_s)

    comp_result = compare_slides(slide_a_model, slide_b_model)
    res_dict = comp_result.to_dict()

    return {
        "success": True,
        **res_dict,
    }


@handle_tool_errors
def ppt_validate_slide(
    slide_number: int,
    rules: Optional[List[str]] = None,
    presentation_path: Optional[str] = None,
    detail: str = "summary",
) -> Dict[str, Any]:
    """Run rule-based geometric and typographic validation on a slide.

    Detects overlaps (VAL-01), boundary clipping (VAL-02), off-slide objects (VAL-03),
    text overflow (VAL-04), tiny fonts (VAL-05), irregular rotation (VAL-08), and duplicate shapes (VAL-07).

    Args:
        slide_number: 1-indexed slide number.
        rules: Optional list of specific rule IDs to check (e.g. ['VAL-01', 'VAL-02']).
        presentation_path: Path to presentation. If omitted, uses active session.
        detail: 'summary' (default, concise report) or 'full' (deep issue detail dictionaries).

    Returns:
        Slide validation report with is_valid, valid, summary counts, warning_count, warnings list, and slide metrics.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    slide_model = inspect_slide(target_path, slide_number)
    val_result = validate_slide(slide_model, rules=rules)

    return {
        "success": True,
        "detail": detail,
        **val_result.to_dict(detail=detail),
    }


@handle_tool_errors
def ppt_analyze_slide_structure(
    slide_number: int,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Analyze the complete semantic layout hierarchy, roles, and card containers on a slide.

    Args:
        slide_number: 1-indexed slide number.
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured layout tree with roles (title, card, card_title, metric, badge, body), confidence, and containers.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    return analyze_slide_structure(target_path, slide_number)


@handle_tool_errors
def ppt_analyze_containers(
    slide_number: int,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Identify and return logical containers/cards and their nested children on a slide.

    Args:
        slide_number: 1-indexed slide number.
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured list of containers with bounding boxes and child shape IDs.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    return analyze_containers(target_path, slide_number)


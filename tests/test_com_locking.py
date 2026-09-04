"""Dedicated test suite for PowerPoint COM locking prevention, process isolation, and defensive retry behavior."""

import gc
from pathlib import Path
import shutil
import sys
import time
import pytest
from pptx import Presentation

from powerpoint_mcp.rendering.com_lifecycle import (
    cleanup_mcp_com_processes,
    com_powerpoint_session,
    defensive_file_operation,
    ensure_mcp_pid_closed,
    get_active_mcp_pids,
    get_powerpoint_pids,
    is_file_locked_error,
    register_mcp_pid,
    terminate_mcp_pid,
    unregister_mcp_pid,
)
from powerpoint_mcp.rendering.renderer import PowerPointRenderer
from powerpoint_mcp.tools.versioning import open_presentation, save_session, save_as, revert_session, create_backup


@pytest.fixture
def sample_presentation(tmp_path):
    """Create a temporary test presentation with 2 slides."""
    file_path = tmp_path / "com_lock_test.pptx"
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "COM Lock Slide 1"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "COM Lock Slide 2"
    prs.save(str(file_path))
    return file_path


class TestCOMLifecycleAndProcessIsolation:
    """Test COM process management, process ID tracking, and process isolation."""

    def test_pid_enumeration_returns_set(self):
        pids = get_powerpoint_pids()
        assert isinstance(pids, set)

    def test_mcp_pid_registry_tracking(self):
        dummy_pid = 9999999
        register_mcp_pid(dummy_pid)
        assert dummy_pid in get_active_mcp_pids()

        # Refuse to terminate if not in registry
        assert terminate_mcp_pid(8888888) is False

        unregister_mcp_pid(dummy_pid)
        assert dummy_pid not in get_active_mcp_pids()

    def test_is_file_locked_error(self):
        assert is_file_locked_error(PermissionError("[WinError 32] Locked")) is True
        os_err = OSError("The process cannot access the file because it is being used by another process")
        assert is_file_locked_error(os_err) is True
        os_err32 = OSError()
        os_err32.winerror = 32
        assert is_file_locked_error(os_err32) is True
        assert is_file_locked_error(ValueError("Not a lock error")) is False

    def test_com_session_invokes_couninitialize_in_outermost_finally(self):
        """Verify pythoncom.CoUninitialize is always called in the outermost finally block."""
        from unittest.mock import MagicMock, patch

        if sys.platform != "win32":
            pytest.skip("Windows only test")

        import pythoncom
        couninit_mock = MagicMock()
        coinit_mock = MagicMock()

        with patch.object(pythoncom, "CoInitialize", coinit_mock), \
             patch.object(pythoncom, "CoUninitialize", couninit_mock), \
             patch("win32com.client.DispatchEx") as dispatch_mock:
            mock_app = MagicMock()
            dispatch_mock.return_value = mock_app

            # Scenario 1: Normal exit
            with com_powerpoint_session() as (app, pid):
                assert app == mock_app

            assert couninit_mock.call_count == 1

            # Scenario 2: Exception inside with block
            couninit_mock.reset_mock()
            with pytest.raises(RuntimeError, match="Inner error"):
                with com_powerpoint_session() as (app, pid):
                    raise RuntimeError("Inner error")

            assert couninit_mock.call_count == 1

    def test_com_session_handles_rpc_disconnect_during_quit(self):
        """Verify com_powerpoint_session handles RPC_E_DISCONNECTED on ppt_app.Quit() gracefully."""
        from unittest.mock import MagicMock, patch
        import pywintypes

        if sys.platform != "win32":
            pytest.skip("Windows only test")

        with patch("win32com.client.DispatchEx") as dispatch_mock, \
             patch("powerpoint_mcp.rendering.com_lifecycle.get_powerpoint_pids", return_value={100}), \
             patch("powerpoint_mcp.rendering.com_lifecycle.ensure_mcp_pid_closed"):
            mock_app = MagicMock()
            # Simulate RPC_E_DISCONNECTED (-2147417848 / 0x80010108)
            mock_app.Quit.side_effect = pywintypes.com_error(-2147417848, "The object invoked has disconnected from its clients.", None, None)
            dispatch_mock.return_value = mock_app

            # Should not raise exception; should exit cleanly
            with com_powerpoint_session() as (app, pid):
                pass


    @pytest.mark.skipif(
        sys.platform != "win32" or not PowerPointRenderer().is_available,
        reason="PowerPoint COM is not available on this platform",
    )
    def test_com_render_releases_file_lock_immediately(self, sample_presentation, tmp_path):
        """Verify rendering a slide leaves 0 locks on the source presentation."""
        renderer = PowerPointRenderer()
        out_png = tmp_path / "rendered_slide_1.png"

        # Render slide 1
        renderer.render_slide(
            presentation_path=sample_presentation,
            slide_number=1,
            output_path=out_png,
        )
        assert out_png.exists()

        # Immediate overwrite / write test: should succeed with no WinError 32
        shutil.copy2(sample_presentation, tmp_path / "copied.pptx")
        with open(sample_presentation, "r+b") as f:
            data = f.read()
            f.seek(0)
            f.write(data)

    @pytest.mark.skipif(
        sys.platform != "win32" or not PowerPointRenderer().is_available,
        reason="PowerPoint COM is not available on this platform",
    )
    def test_com_render_presentation_all_slides_clean_cleanup(self, sample_presentation, tmp_path):
        """Verify rendering all slides in deck terminates COM and releases locks."""
        renderer = PowerPointRenderer()
        out_dir = tmp_path / "renders"

        rendered = renderer.render_presentation(
            presentation_path=sample_presentation,
            output_dir=out_dir,
        )
        assert len(rendered) == 2

        # Verify no orphaned MCP processes remain in registry
        assert len(get_active_mcp_pids()) == 0

        # Overwrite source file
        with open(sample_presentation, "r+b") as f:
            data = f.read()
            f.seek(0)
            f.write(data)

    @pytest.mark.skipif(
        sys.platform != "win32" or not PowerPointRenderer().is_available,
        reason="PowerPoint COM is not available on this platform",
    )
    def test_com_render_out_of_bounds_no_rpc_disconnect_crash(self, sample_presentation, tmp_path):
        """Verify out of bounds slide raises IndexError without 0x80010108 RPC crash."""
        renderer = PowerPointRenderer()
        with pytest.raises(IndexError, match="out of range"):
            renderer.render_slide(
                presentation_path=sample_presentation,
                slide_number=999,
                output_path=tmp_path / "err.png",
            )

        # Confirm source file is unlocked and writable
        with open(sample_presentation, "r+b") as f:
            data = f.read()
            f.seek(0)
            f.write(data)


class TestDefensiveRetryBehavior:
    """Test defensive retry mechanisms for file operations."""

    def test_defensive_retry_succeeds_on_transient_lock(self, tmp_path):
        target = tmp_path / "transient.txt"
        attempts = 0

        def flaky_op():
            nonlocal attempts
            attempts += 1
            if attempts < 3:
                err = PermissionError("[WinError 32] The process cannot access the file because it is being used by another process")
                err.winerror = 32
                raise err
            target.write_text("success", encoding="utf-8")
            return "ok"

        res = defensive_file_operation(flaky_op, target, action_name="test_write", max_retries=4, initial_delay=0.01)
        assert res == "ok"
        assert attempts == 3
        assert target.read_text(encoding="utf-8") == "success"

    def test_defensive_retry_raises_actionable_error_when_persistently_locked(self, tmp_path):
        target = tmp_path / "locked.txt"

        def permanently_locked_op():
            err = PermissionError("[WinError 32] The process cannot access the file because it is being used by another process")
            err.winerror = 32
            raise err

        with pytest.raises(PermissionError) as exc_info:
            defensive_file_operation(
                permanently_locked_op,
                target,
                action_name="save_presentation",
                max_retries=2,
                initial_delay=0.01,
            )

        err_msg = str(exc_info.value)
        assert "Cannot save_presentation presentation" in err_msg
        assert "[WinError 32]" in err_msg
        assert "Please ensure no other application is accessing the file" in err_msg or "PowerPoint process" in err_msg

    def test_defensive_retry_propagates_non_lock_errors_immediately(self, tmp_path):
        target = tmp_path / "value_error.txt"
        attempts = 0

        def error_op():
            nonlocal attempts
            attempts += 1
            raise ValueError("Corrupt data")

        with pytest.raises(ValueError, match="Corrupt data"):
            defensive_file_operation(error_op, target, max_retries=3)

        assert attempts == 1


class TestEndToEndSessionSaveAfterRender:
    """Test full editing session lifecycle: open -> render -> mutate -> save."""

    @pytest.mark.skipif(
        sys.platform != "win32" or not PowerPointRenderer().is_available,
        reason="PowerPoint COM is not available on this platform",
    )
    def test_session_open_render_save_cycle(self, sample_presentation, tmp_path):
        # 1. Open session
        session = open_presentation(sample_presentation)
        assert session is not None
        assert Path(session.working_path).exists()

        # 2. Render slide 1 via COM from session working copy
        renderer = PowerPointRenderer()
        render_png = tmp_path / "session_render_1.png"
        renderer.render_slide(
            presentation_path=session.working_path,
            slide_number=1,
            output_path=render_png,
        )
        assert render_png.exists()

        # 3. Save session back to original presentation
        save_res = save_session(session.session_id)
        assert save_res["success"] is True
        assert Path(save_res["saved_path"]).exists()

        # 4. Save as a new destination file
        save_as_path = tmp_path / "exported_deck.pptx"
        save_as_res = save_as(session.session_id, output_path=save_as_path)
        assert save_as_res["success"] is True
        assert save_as_path.exists()

"""Adversarial stress tests for shapes, geometry, unicode, nested shapes, and edge cases."""

import io
import math
import os
from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    ParagraphModel,
    PresentationModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    TextFrameModel,
    TextRunModel,
    TextStyle,
    apply_delta_inches,
    emu_to_inches,
    emu_to_pt,
    inches_to_emu,
    pt_to_emu,
)
from powerpoint_mcp.pptx import (
    PPTXInspector,
    extract_alignment_name,
    extract_embedded_images,
    extract_fill_style,
    extract_font_style,
    extract_line_style,
    extract_paragraph,
    extract_rgb_hex,
    extract_text_frame,
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    map_shape_type,
    match_shapes,
)


class TestAdversarialBoundingBox:
    """Stress testing BoundingBox math with extreme and adversarial values."""

    def test_zero_dimensions(self):
        bbox = BoundingBox(left_emu=0, top_emu=0, width_emu=0, height_emu=0)
        assert bbox.left_inches == 0.0
        assert bbox.top_inches == 0.0
        assert bbox.width_inches == 0.0
        assert bbox.height_inches == 0.0
        assert bbox.right_inches == 0.0
        assert bbox.bottom_inches == 0.0
        assert bbox.center_x_inches == 0.0
        assert bbox.center_y_inches == 0.0
        d = bbox.to_dict()
        assert d["x"] == 0.0
        assert d["width"] == 0.0

    def test_extreme_negative_coordinates(self):
        bbox = BoundingBox(
            left_emu=-9144000,   # -10 inches
            top_emu=-18288000,   # -20 inches
            width_emu=914400,    # 1 inch
            height_emu=1828800,  # 2 inches
        )
        assert bbox.left_inches == -10.0
        assert bbox.top_inches == -20.0
        assert bbox.right_inches == -9.0
        assert bbox.bottom_inches == -18.0
        assert bbox.center_x_inches == -9.5
        assert bbox.center_y_inches == -19.0

    def test_massive_coordinates(self):
        # 1 billion EMUs (~1,093 inches)
        bbox = BoundingBox(
            left_emu=1000000000,
            top_emu=2000000000,
            width_emu=500000000,
            height_emu=500000000,
        )
        assert bbox.left_inches == pytest.approx(1093.6133, abs=0.001)
        assert bbox.width_inches == pytest.approx(546.8066, abs=0.001)
        assert bbox.right_inches == pytest.approx(1640.42, abs=0.01)

    def test_odd_integer_emu_center_rounding(self):
        bbox = BoundingBox(left_emu=1, top_emu=3, width_emu=5, height_emu=7)
        # center_x = 1 + 5 // 2 = 3
        # center_y = 3 + 7 // 2 = 6
        assert bbox.center_x_emu == 3
        assert bbox.center_y_emu == 6

    def test_cumulative_delta_precision(self):
        # Apply small delta (0.0001 inches) 10,000 times
        emu = 0
        delta = 0.0001
        for _ in range(10000):
            emu = apply_delta_inches(emu, delta)
        # 10,000 * 0.0001 = 1.0 inch = 914400 EMU (or within discrete step rounding)
        inches = emu_to_inches(emu)
        assert inches == pytest.approx(1.0, abs=0.01)


class TestAdversarialShapesPresentation:
    """Generate and inspect presentations with adversarial PPTX elements."""

    @pytest.fixture
    def adversarial_deck(self, tmp_path):
        """Build a presentation containing extreme and unusual shapes."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        blank_layout = prs.slide_layouts[6]

        # -------------------------------------------------------------
        # Slide 1: Zero-area, massive coordinates, negative coordinates
        # -------------------------------------------------------------
        slide1 = prs.slides.add_slide(blank_layout)

        # 1. Zero area shape (1x0 EMU point/line box)
        slide1.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            0,
            0,
        )

        # 2. Line connector with 0 width (vertical line)
        from pptx.enum.shapes import MSO_CONNECTOR
        slide1.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(2),
            Inches(1),
            Inches(2),
            Inches(4),
        )

        # 3. Off-canvas negative position shape
        neg_box = slide1.shapes.add_textbox(
            Inches(-3),
            Inches(-2),
            Inches(2),
            Inches(1),
        )
        neg_box.text_frame.text = "Offscreen Top-Left"

        # 4. Far off-canvas positive position shape
        far_box = slide1.shapes.add_textbox(
            Inches(25),
            Inches(30),
            Inches(5),
            Inches(2),
        )
        far_box.text_frame.text = "Far Away Shape"

        # -------------------------------------------------------------
        # Slide 2: Extreme Unicode, RTL, Emojis, XML Injection, Long Text
        # -------------------------------------------------------------
        slide2 = prs.slides.add_slide(blank_layout)

        # 1. Unicode multilingual text box
        u_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
        tf_u = u_box.text_frame
        p_cjk = tf_u.paragraphs[0]
        p_cjk.text = "日本語テスト 中文测试 한국어 繁體字"
        p_rtl = tf_u.add_paragraph()
        p_rtl.text = "مرحبا بالعالم - שלום עולם (RTL text)"
        p_emoji = tf_u.add_paragraph()
        p_emoji.text = "🚀 🤖 💻 📊 📈 🔥 🧑🏽‍💻 🏳️‍🌈"

        # 2. XML entity & injection test
        xml_box = slide2.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
        tf_xml = xml_box.text_frame
        p_xml = tf_xml.paragraphs[0]
        p_xml.text = "<script>alert('XSS')</script> &amp; <test attr=\"val\"> 'quote' & < >"

        # 3. Very long text content (10,000 characters across 50 paragraphs)
        long_box = slide2.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.8))
        tf_long = long_box.text_frame
        for i in range(50):
            p = tf_long.paragraphs[0] if i == 0 else tf_long.add_paragraph()
            p.text = f"Paragraph {i:03d}: " + ("Adversarial stress text " * 10)

        # -------------------------------------------------------------
        # Slide 3: Complex Styles, Obscure Fonts, Lines, Fills
        # -------------------------------------------------------------
        slide3 = prs.slides.add_slide(blank_layout)

        # 1. Obscure font with fractional size and custom color
        font_box = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        p_f = font_box.text_frame.paragraphs[0]
        r_f = p_f.add_run()
        r_f.text = "Special Font Run"
        r_f.font.name = "NonExistentCrazyFont_12345!@#"
        r_f.font.size = Pt(13.75)
        r_f.font.bold = True
        r_f.font.italic = False
        r_f.font.color.rgb = RGBColor(12, 34, 56)

        # 2. Empty text frame
        empty_box = slide3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(1))
        empty_box.text_frame.text = ""

        # 3. Whitespace-only text frame
        ws_box = slide3.shapes.add_textbox(Inches(5), Inches(2.5), Inches(3), Inches(1))
        ws_box.text_frame.text = "\t  \n \u00A0 \u200B \n   "

        # 4. Colored shape with thick border
        shape_styled = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1),
            Inches(4),
            Inches(4),
            Inches(1),
        )
        shape_styled.fill.solid()
        shape_styled.fill.fore_color.rgb = RGBColor(200, 100, 50)
        shape_styled.line.color.rgb = RGBColor(50, 100, 200)
        shape_styled.line.width = Pt(4.5)

        deck_path = tmp_path / "adversarial_stress_deck.pptx"
        prs.save(str(deck_path))
        return str(deck_path)

    def test_inspect_adversarial_deck_presentation_level(self, adversarial_deck):
        prs_model = inspect_presentation(adversarial_deck)
        assert prs_model.slide_count == 3
        assert len(prs_model.slides) == 3
        assert prs_model.width_inches == 10.0
        assert prs_model.height_inches == 5.625

    def test_inspect_zero_area_and_negative_coordinates(self, adversarial_deck):
        slide1 = inspect_slide(adversarial_deck, 1)
        assert slide1.shape_count >= 4

        # Zero width/height shape
        s0 = slide1.shapes[0]
        assert s0.bbox.width_emu >= 0
        assert s0.bbox.height_emu >= 0

        # Negative coordinate shape
        neg_s = next(s for s in slide1.shapes if s.bbox.left_inches < 0)
        assert neg_s.bbox.left_inches == -3.0
        assert neg_s.bbox.top_inches == -2.0
        assert "Offscreen" in neg_s.text_frame.text

        # Far positive coordinate shape
        far_s = next(s for s in slide1.shapes if s.bbox.left_inches > 20)
        assert far_s.bbox.left_inches == 25.0
        assert "Far Away" in far_s.text_frame.text

        # Serialization to dict works cleanly without errors
        d1 = slide1.to_dict()
        assert len(d1["shapes"]) == slide1.shape_count

    def test_inspect_unicode_rtl_emoji_and_long_text(self, adversarial_deck):
        slide2 = inspect_slide(adversarial_deck, 2)
        assert slide2.shape_count >= 3

        # Unicode text shape
        u_shape = slide2.shapes[0]
        assert "日本語テスト" in u_shape.text_frame.text
        assert "مرحبا بالعالم" in u_shape.text_frame.text
        assert "🚀" in u_shape.text_frame.text

        # XML injection characters preserved as raw text
        xml_shape = slide2.shapes[1]
        assert "<script>alert('XSS')</script>" in xml_shape.text_frame.text
        assert "<test attr=\"val\">" in xml_shape.text_frame.text

        # Long text shape (50 paragraphs)
        long_shape = slide2.shapes[2]
        assert len(long_shape.text_frame.paragraphs) == 50
        assert len(long_shape.text_frame.text) > 5000

        # Dict serialization handles all unicode cleanly
        d2 = slide2.to_dict()
        assert "日本語テスト" in d2["shapes"][0]["text"]
        assert "🚀" in d2["shapes"][0]["text"]

    def test_inspect_obscure_fonts_and_styles(self, adversarial_deck):
        slide3 = inspect_slide(adversarial_deck, 3)

        # Obscure font
        font_shape = slide3.shapes[0]
        run_style = font_shape.text_frame.paragraphs[0].runs[0].style
        assert run_style.font_name == "NonExistentCrazyFont_12345!@#"
        assert run_style.font_size_pt == 13.75
        assert run_style.color_rgb == "0C2238"

        # Styled shape with solid fill & thick line
        styled_shape = slide3.shapes[3]
        assert styled_shape.fill["type"] == "solid"
        assert styled_shape.fill_color == "C86432"
        assert styled_shape.line_color == "3264C8"
        assert styled_shape.line_width_pt == 4.5

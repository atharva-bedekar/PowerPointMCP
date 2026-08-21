"""Adversarial stress tests for shape matching stability, perturbations, and swapped order."""

import copy
import random
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    ParagraphModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    TextFrameModel,
    TextRunModel,
    TextStyle,
)
from powerpoint_mcp.pptx import inspect_slide, match_shapes


class TestShapeMatchingStability:
    """Stress test shape matching under ordering perturbations, spatial shifts, and structural variations."""

    @pytest.fixture
    def sample_slide(self):
        """Create a SlideModel with 6 diverse shapes."""
        s1 = ShapeModel(
            shape_id=1,
            name="Title Box",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.TITLE,
            bbox=BoundingBox.from_inches(1.0, 0.8, 8.0, 1.0),
            text_frame=TextFrameModel(
                text="Quarterly Financial Report",
                paragraphs=[ParagraphModel(text="Quarterly Financial Report", runs=[TextRunModel("Quarterly Financial Report")])],
            ),
        )
        s2 = ShapeModel(
            shape_id=2,
            name="Subtitle Box",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.SUBTITLE,
            bbox=BoundingBox.from_inches(1.0, 1.8, 8.0, 0.6),
            text_frame=TextFrameModel(
                text="Q3 Consolidated Performance Summary",
                paragraphs=[ParagraphModel(text="Q3 Consolidated Performance Summary", runs=[TextRunModel("Q3 Consolidated Performance Summary")])],
            ),
        )
        s3 = ShapeModel(
            shape_id=3,
            name="Card Left",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(1.0, 2.8, 3.8, 3.0),
            text_frame=TextFrameModel(
                text="Revenue growth reached 18.5% YoY with strong enterprise expansion.",
                paragraphs=[ParagraphModel(text="Revenue growth reached 18.5% YoY with strong enterprise expansion.", runs=[TextRunModel("Revenue growth reached 18.5% YoY with strong enterprise expansion.")])],
            ),
        )
        s4 = ShapeModel(
            shape_id=4,
            name="Card Right",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=BoundingBox.from_inches(5.2, 2.8, 3.8, 3.0),
            text_frame=TextFrameModel(
                text="Operating margin expanded 240bps due to cloud efficiency gains.",
                paragraphs=[ParagraphModel(text="Operating margin expanded 240bps due to cloud efficiency gains.", runs=[TextRunModel("Operating margin expanded 240bps due to cloud efficiency gains.")])],
            ),
        )
        s5 = ShapeModel(
            shape_id=5,
            name="Chart Placeholder",
            shape_type=ShapeType.CHART,
            semantic_role=SemanticRole.CHART,
            bbox=BoundingBox.from_inches(9.2, 2.8, 3.0, 3.0),
            chart_metadata={"chart_type": "COLUMN_CLUSTERED"},
        )
        s6 = ShapeModel(
            shape_id=6,
            name="Footer Text",
            shape_type=ShapeType.TEXT_BOX,
            semantic_role=SemanticRole.FOOTER,
            bbox=BoundingBox.from_inches(1.0, 6.8, 11.0, 0.4),
            text_frame=TextFrameModel(
                text="Confidential - Internal Use Only",
                paragraphs=[ParagraphModel(text="Confidential - Internal Use Only", runs=[TextRunModel("Confidential - Internal Use Only")])],
            ),
        )
        return SlideModel(
            slide_number=1,
            slide_id=101,
            layout_name="Custom",
            title="Quarterly Financial Report",
            width_inches=13.3333,
            height_inches=7.5,
            width_emu=12192000,
            height_emu=6858000,
            shapes=[s1, s2, s3, s4, s5, s6],
        )

    def test_match_shapes_order_invariance(self, sample_slide):
        """Shape matching should produce identical 1-to-1 pairs regardless of shape list order."""
        slide_a = sample_slide
        # Create Slide B with shapes in reverse order and different shape_ids
        shapes_reversed = []
        for orig in reversed(sample_slide.shapes):
            s_copy = copy.deepcopy(orig)
            s_copy.shape_id = orig.shape_id + 100  # Shift IDs
            shapes_reversed.append(s_copy)

        slide_b = SlideModel(
            slide_number=2,
            slide_id=102,
            layout_name=sample_slide.layout_name,
            title=sample_slide.title,
            width_inches=sample_slide.width_inches,
            height_inches=sample_slide.height_inches,
            width_emu=sample_slide.width_emu,
            height_emu=sample_slide.height_emu,
            shapes=shapes_reversed,
        )

        matches = match_shapes(slide_a, slide_b, min_confidence=0.50)
        assert len(matches) == len(sample_slide.shapes)

        # Check every shape matched with its corresponding reversed shape
        for sa in slide_a.shapes:
            m = next((m for m in matches if m["shape_a_id"] == sa.shape_id), None)
            assert m is not None
            assert m["shape_b_id"] == sa.shape_id + 100
            assert m["confidence_score"] >= 0.95

    def test_match_shapes_random_permutations(self, sample_slide):
        """Test stability across 10 random permutations of the shape list."""
        rng = random.Random(42)
        slide_a = sample_slide

        for trial in range(10):
            shuffled_shapes = copy.deepcopy(sample_slide.shapes)
            rng.shuffle(shuffled_shapes)
            for s in shuffled_shapes:
                s.shape_id = s.shape_id + 200

            slide_b = SlideModel(
                slide_number=trial + 2,
                slide_id=100 + trial,
                shapes=shuffled_shapes,
                width_emu=sample_slide.width_emu,
                height_emu=sample_slide.height_emu,
            )

            matches = match_shapes(slide_a, slide_b, min_confidence=0.50)
            assert len(matches) == len(sample_slide.shapes)
            for sa in slide_a.shapes:
                m = next(m for m in matches if m["shape_a_id"] == sa.shape_id)
                assert m["shape_b_id"] == sa.shape_id + 200

    def test_match_shapes_spatial_and_text_perturbation(self, sample_slide):
        """Test matching when coordinates are shifted and text has typos."""
        slide_a = sample_slide
        perturbed_shapes = copy.deepcopy(sample_slide.shapes)

        # Shift title slightly (0.3 inches right) and change one word
        title_b = perturbed_shapes[0]
        title_b.shape_id = 1001
        title_b.bbox.left_emu += 274320  # 0.3 inches
        title_b.text_frame.text = "Quarterly Financial Overview"  # Report -> Overview

        # Shift Card Left
        card_l_b = perturbed_shapes[2]
        card_l_b.shape_id = 1003
        card_l_b.bbox.top_emu += 182880  # 0.2 inches down

        slide_b = SlideModel(
            slide_number=2,
            slide_id=102,
            shapes=perturbed_shapes,
            width_emu=sample_slide.width_emu,
            height_emu=sample_slide.height_emu,
        )

        matches = match_shapes(slide_a, slide_b, min_confidence=0.40)
        assert len(matches) == len(sample_slide.shapes)

        # Title match should still be top confidence despite text edit and shift
        title_match = next(m for m in matches if m["shape_a_id"] == 1)
        assert title_match["shape_b_id"] == 1001
        assert title_match["confidence_score"] >= 0.80

    def test_match_shapes_swapped_positions(self, sample_slide):
        """When two cards swap positions, the text and names should preserve the true logical identity."""
        slide_a = sample_slide
        shapes_b = copy.deepcopy(sample_slide.shapes)

        # Swap bounding boxes of Card Left (id 3) and Card Right (id 4)
        card_3_b = next(s for s in shapes_b if s.shape_id == 3)
        card_4_b = next(s for s in shapes_b if s.shape_id == 4)
        card_3_b.shape_id = 303
        card_4_b.shape_id = 404

        bbox_temp = card_3_b.bbox
        card_3_b.bbox = card_4_b.bbox
        card_4_b.bbox = bbox_temp

        slide_b = SlideModel(
            slide_number=2,
            slide_id=102,
            shapes=shapes_b,
            width_emu=sample_slide.width_emu,
            height_emu=sample_slide.height_emu,
        )

        matches = match_shapes(slide_a, slide_b, min_confidence=0.40)
        assert len(matches) == len(sample_slide.shapes)

        # Card Left (id 1) should match with card_3_b (id 303) based on text similarity and name
        match_3 = next(m for m in matches if m["shape_a_id"] == 3)
        assert match_3["shape_b_id"] == 303

        # Card Right (id 4) should match with card_4_b (id 404)
        match_4 = next(m for m in matches if m["shape_a_id"] == 4)
        assert match_4["shape_b_id"] == 404

    def test_match_shapes_asymmetric_shape_counts(self, sample_slide):
        """Matching when slide A has 6 shapes and slide B has only 2 shapes."""
        slide_a = sample_slide
        # Slide B only has Title and Footer
        shapes_b = [
            copy.deepcopy(sample_slide.shapes[0]),
            copy.deepcopy(sample_slide.shapes[5]),
        ]
        shapes_b[0].shape_id = 501
        shapes_b[1].shape_id = 502

        slide_b = SlideModel(
            slide_number=2,
            slide_id=102,
            shapes=shapes_b,
            width_emu=sample_slide.width_emu,
            height_emu=sample_slide.height_emu,
        )

        matches = match_shapes(slide_a, slide_b, min_confidence=0.40)
        assert len(matches) == 2
        matched_a_ids = {m["shape_a_id"] for m in matches}
        assert 1 in matched_a_ids  # Title matched
        assert 6 in matched_a_ids  # Footer matched

    def test_match_shapes_duplicate_identical_blank_cards(self):
        """Slide with 4 identical blank cards at different positions should match 1-to-1 based on coordinates."""
        shapes_a = []
        shapes_b = []
        for i in range(4):
            x = 1.0 + i * 2.5
            sa = ShapeModel(
                shape_id=i + 1,
                name=f"Rectangle {i + 1}",
                shape_type=ShapeType.AUTO_SHAPE,
                semantic_role=SemanticRole.BODY,
                bbox=BoundingBox.from_inches(x, 2.0, 2.0, 3.0),
            )
            # Slide B has the exact same positions but different IDs
            sb = ShapeModel(
                shape_id=i + 101,
                name=f"Rectangle {i + 1}",
                shape_type=ShapeType.AUTO_SHAPE,
                semantic_role=SemanticRole.BODY,
                bbox=BoundingBox.from_inches(x, 2.0, 2.0, 3.0),
            )
            shapes_a.append(sa)
            shapes_b.append(sb)

        slide_a = SlideModel(slide_number=1, slide_id=1, shapes=shapes_a, width_emu=12192000, height_emu=6858000)
        slide_b = SlideModel(slide_number=2, slide_id=2, shapes=shapes_b, width_emu=12192000, height_emu=6858000)

        matches = match_shapes(slide_a, slide_b, min_confidence=0.40)
        assert len(matches) == 4

        # Every blank card should match to its exact positional counterpart
        for i in range(4):
            m = next(m for m in matches if m["shape_a_id"] == i + 1)
            assert m["shape_b_id"] == i + 101
            assert m["confidence_score"] >= 0.90

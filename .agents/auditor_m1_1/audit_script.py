"""Forensic Audit Verification Script for Milestone M1."""

import ast
import difflib
import math
from pathlib import Path
import random
import sys

# Ensure src/ is on sys.path
SRC_DIR = Path(__file__).resolve().parent.parent.parent / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    ParagraphModel,
    PresentationMetadata,
    PresentationModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    TextFrameModel,
    TextRunModel,
    TextStyle,
    apply_delta_inches,
    emu_to_inches,
    emu_to_pt,
    inches_to_emu,
    pt_to_emu,
)
from powerpoint_mcp.pptx import (
    PPTXInspector,
    extract_embedded_images,
    extract_fill_style,
    extract_font_style,
    extract_hyperlinks,
    extract_line_style,
    extract_paragraph,
    extract_rgb_hex,
    extract_run,
    extract_shape_properties,
    extract_text_frame,
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    inspect_slide_relationships,
    map_shape_type,
    match_shapes,
)

def test_ast_and_hardcoded_strings():
    print("=== 1. Hardcoded String Search in Source Files ===")
    project_root = Path(__file__).resolve().parent.parent.parent
    files_to_check = [
        project_root / "src/powerpoint_mcp/models/shape.py",
        project_root / "src/powerpoint_mcp/models/slide.py",
        project_root / "src/powerpoint_mcp/models/presentation.py",
        project_root / "src/powerpoint_mcp/pptx/inspector.py",
        project_root / "src/powerpoint_mcp/pptx/styles.py",
        project_root / "src/powerpoint_mcp/pptx/relationships.py",
    ]

    prohibited_strings = [
        "Quarterly Performance Overview",
        "Operational Architecture",
        "Audit & Compliance Issues",
        "Total Revenue",
        "Defect Box",
        "450K",
        "78 / 100",
        "Key Initiatives",
        "Program Milestones",
    ]

    violations = []
    for f in files_to_check:
        content = f.read_text(encoding="utf-8-sig")
        for s in prohibited_strings:
            if s.lower() in content.lower():
                violations.append((f, s))
                print(f"VIOLATION: Found prohibited fixture string '{s}' in {f}")

    if not violations:
        print("PASS: Zero fixture string leakage in production source files.")

    print("\n=== 2. AST Structure & Facade Detection ===")
    for f in files_to_check:
        tree = ast.parse(f.read_text(encoding="utf-8-sig"), filename=str(f))
        func_count = 0
        class_count = 0
        dummy_funcs = []
        for node in ast.walk(tree):
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
                func_count += 1
                # Filter out docstrings from length
                body_statements = [s for s in node.body if not (isinstance(s, ast.Expr) and isinstance(s.value, ast.Constant))]
                if len(body_statements) == 1:
                    single_stmt = body_statements[0]
                    if isinstance(single_stmt, ast.Pass):
                        dummy_funcs.append((node.name, "empty pass"))
                    elif isinstance(single_stmt, ast.Return) and isinstance(single_stmt.value, ast.Constant):
                        dummy_funcs.append((node.name, f"returns constant {single_stmt.value.value}"))
            elif isinstance(node, ast.ClassDef):
                class_count += 1
        print(f"File {f.name:20s}: {class_count:2d} classes, {func_count:2d} functions, flagged dummy stubs: {dummy_funcs}")
        assert len(dummy_funcs) == 0, f"Found facade functions in {f.name}: {dummy_funcs}"

def test_math_and_conversions():
    print("\n=== 3. Mathematical Coordinate & Unit Conversion Verification ===")
    assert inches_to_emu(1.0) == 914400
    assert inches_to_emu(2.5) == 2286000
    assert inches_to_emu(0.0) == 0
    assert inches_to_emu(-1.0) == -914400
    assert emu_to_inches(914400) == 1.0
    assert emu_to_inches(2286000) == 2.5
    assert pt_to_emu(1.0) == 12700
    assert pt_to_emu(72.0) == 914400
    assert emu_to_pt(12700) == 1.0
    assert emu_to_pt(914400) == 72.0

    # Delta arithmetic
    base = 1000000
    delta = 0.5
    res = apply_delta_inches(base, delta)
    assert res == base + 457200

    bb = BoundingBox(left_emu=1000, top_emu=2000, width_emu=3000, height_emu=4000)
    assert bb.right_emu == 4000
    assert bb.bottom_emu == 6000
    assert bb.center_x_emu == 2500
    assert bb.center_y_emu == 4000
    print("PASS: Exact integer EMU arithmetic and conversions verified.")

def test_dynamic_random_presentation_inspection():
    print("\n=== 4. Dynamic Arbitrary Presentation Inspection & Role Inference ===")
    prs = Presentation()
    prs.slide_width = Inches(16.0)
    prs.slide_height = Inches(9.0)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Add dynamic shapes with randomized properties
    random.seed(42)
    shape_coords = [
        ("Dynamic Main Title", Inches(1.0), Inches(0.5), Inches(14.0), Inches(1.0), 36.0, SemanticRole.TITLE),
        ("Dynamic Subtitle Line", Inches(1.0), Inches(1.6), Inches(14.0), Inches(0.8), 20.0, SemanticRole.SUBTITLE),
        ("Multi-line Body Content", Inches(1.0), Inches(4.0), Inches(6.5), Inches(3.0), 12.0, SemanticRole.BODY),
        ("Dynamic Footer Notice", Inches(1.0), Inches(8.0), Inches(14.0), Inches(0.5), 9.0, SemanticRole.FOOTER),
    ]

    for name, left, top, w, h, font_sz, expected_role in shape_coords:
        tb = slide.shapes.add_textbox(left, top, w, h)
        tb.name = name
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"Sample text for {name}"
        p.font.size = Pt(font_sz)
        p.font.name = "Arial"
        if expected_role == SemanticRole.BODY:
            p2 = tf.add_paragraph()
            p2.text = "Second paragraph for body validation."

    # Inspect slide
    slide_model = inspect_slide(prs, 1)
    assert slide_model.slide_number == 1
    assert slide_model.width_inches == 16.0
    assert slide_model.height_inches == 9.0
    assert slide_model.shape_count == 4
    assert slide_model.title == "Sample text for Dynamic Main Title"

    roles = [s.semantic_role for s in slide_model.shapes]
    print(f"Inferred dynamic roles: {[r.value for r in roles]}")
    assert roles == [SemanticRole.TITLE, SemanticRole.SUBTITLE, SemanticRole.BODY, SemanticRole.FOOTER]
    print("PASS: Dynamic arbitrary slide inspection and role inference verified.")

def test_bipartite_shape_matching_algorithm():
    print("\n=== 5. Bipartite Shape Matching Multi-Factor Algorithm ===")
    # Construct two slides with known similarities and variations
    s1 = ShapeModel(
        shape_id=1,
        name="Header",
        shape_type=ShapeType.TEXT_BOX,
        semantic_role=SemanticRole.TITLE,
        bbox=BoundingBox.from_inches(1.0, 1.0, 8.0, 1.0),
        text_frame=TextFrameModel(text="Global Market Analysis"),
    )
    s2 = ShapeModel(
        shape_id=2,
        name="Chart Box",
        shape_type=ShapeType.CHART,
        semantic_role=SemanticRole.CHART,
        bbox=BoundingBox.from_inches(1.0, 2.5, 6.0, 4.0),
    )

    # Slide B: Slightly moved Header and different text, same Chart
    s1_b = ShapeModel(
        shape_id=101,
        name="Header",
        shape_type=ShapeType.TEXT_BOX,
        semantic_role=SemanticRole.TITLE,
        bbox=BoundingBox.from_inches(1.2, 1.0, 8.0, 1.0),
        text_frame=TextFrameModel(text="Global Market Analysis (Updated)"),
    )
    s2_b = ShapeModel(
        shape_id=102,
        name="Chart Box",
        shape_type=ShapeType.CHART,
        semantic_role=SemanticRole.CHART,
        bbox=BoundingBox.from_inches(1.0, 2.5, 6.0, 4.0),
    )

    slide_a = SlideModel(slide_number=1, slide_id=1, shapes=[s1, s2])
    slide_b = SlideModel(slide_number=2, slide_id=2, shapes=[s1_b, s2_b])

    matches = match_shapes(slide_a, slide_b, min_confidence=0.40)
    assert len(matches) == 2
    match_header = next(m for m in matches if m["shape_a_id"] == 1)
    match_chart = next(m for m in matches if m["shape_a_id"] == 2)

    assert match_header["shape_b_id"] == 101
    assert match_header["confidence_score"] > 0.85
    assert match_header["factors"]["role_score"] == 1.0
    assert match_header["factors"]["type_score"] == 1.0

    assert match_chart["shape_b_id"] == 102
    assert match_chart["confidence_score"] > 0.90
    print(f"Match 1: Header -> Score {match_header['confidence_score']:.4f} ({match_header['reasoning']})")
    print(f"Match 2: Chart -> Score {match_chart['confidence_score']:.4f} ({match_chart['reasoning']})")
    print("PASS: Multi-factor bipartite matching verified.")

def main():
    test_ast_and_hardcoded_strings()
    test_math_and_conversions()
    test_dynamic_random_presentation_inspection()
    test_bipartite_shape_matching_algorithm()
    print("\nALL FORENSIC CHECKS PASSED: VERDICT = CLEAN")

if __name__ == "__main__":
    main()

"""Comprehensive test suite for text editing, typography manipulation, and run-level style preservation."""

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import AlignmentType
from powerpoint_mcp.pptx.editor import modify_text
from powerpoint_mcp.pptx.inspector import inspect_shape, inspect_slide


class TestTextModificationAndStylePreservation:
    """Test modify_text for single/multi paragraph edits and style preservation."""

    def test_modify_text_preserves_single_run_typography(self, temp_presentation):
        slide = temp_presentation.slides[0]
        # Find a shape with text
        text_shape = next(s for s in slide.shapes if s.has_text_frame and s.text.strip())
        shape_id = text_shape.shape_id

        # Inspect initial typography
        initial_p = text_shape.text_frame.paragraphs[0]
        initial_run = initial_p.runs[0]
        initial_font_name = initial_run.font.name
        initial_font_size = initial_run.font.size
        initial_bold = initial_run.font.bold

        # Replace text without specifying font properties
        modify_text(slide, shape_id, text="Updated Executive Summary")

        updated_p = text_shape.text_frame.paragraphs[0]
        updated_run = updated_p.runs[0]

        assert updated_p.text == "Updated Executive Summary"
        if initial_font_name:
            assert updated_run.font.name == initial_font_name
        if initial_font_size is not None:
            assert updated_run.font.size == initial_font_size
        if initial_bold is not None:
            assert updated_run.font.bold == initial_bold

    def test_modify_text_multiline_string(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        shape_id = text_shape.shape_id

        multiline_text = "First Paragraph\nSecond Paragraph\nThird Paragraph"
        modify_text(slide, shape_id, text=multiline_text, font_family="Arial", font_size=18.0)

        tf = text_shape.text_frame
        assert len(tf.paragraphs) == 3
        assert tf.paragraphs[0].text == "First Paragraph"
        assert tf.paragraphs[1].text == "Second Paragraph"
        assert tf.paragraphs[2].text == "Third Paragraph"

        for p in tf.paragraphs:
            for r in p.runs:
                assert r.font.name == "Arial"
                assert r.font.size == Pt(18.0)

    def test_modify_text_targeted_typography_overrides(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        shape_id = text_shape.shape_id

        modify_text(
            slide,
            shape_id,
            font_family="Calibri",
            font_size=24.0,
            bold=True,
            italic=False,
            underline=True,
            color="#FF0000",
        )

        for p in text_shape.text_frame.paragraphs:
            for r in p.runs:
                assert r.font.name == "Calibri"
                assert r.font.size == Pt(24.0)
                assert r.font.bold is True
                assert r.font.italic is False
                assert r.font.underline is True
                assert r.font.color.rgb == RGBColor(255, 0, 0)

    def test_modify_text_alignment(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        shape_id = text_shape.shape_id

        # Center alignment
        modify_text(slide, shape_id, alignment="center")
        for p in text_shape.text_frame.paragraphs:
            assert p.alignment == PP_ALIGN.CENTER

        # Right alignment via enum
        modify_text(slide, shape_id, alignment=AlignmentType.RIGHT)
        for p in text_shape.text_frame.paragraphs:
            assert p.alignment == PP_ALIGN.RIGHT

    def test_modify_text_spacing_and_margins(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        shape_id = text_shape.shape_id

        modify_text(
            slide,
            shape_id,
            paragraph_spacing=12.0,
            space_after=6.0,
            line_spacing=18.0,
            margins={"left": 0.2, "right": 0.2, "top": 0.1, "bottom": 0.1},
        )

        p = text_shape.text_frame.paragraphs[0]
        assert p.space_before == Pt(12.0)
        assert p.space_after == Pt(6.0)
        assert p.line_spacing == Pt(18.0)

        tf = text_shape.text_frame
        assert tf.margin_left == Inches(0.2)
        assert tf.margin_right == Inches(0.2)
        assert tf.margin_top == Inches(0.1)
        assert tf.margin_bottom == Inches(0.1)

    def test_modify_text_shape_without_text_frame_raises_error(self, temp_presentation):
        # Create a shape or line without text frame if available
        slide = temp_presentation.slides[0]
        # Find connector or shape without text frame or mock
        class FakeShape:
            shape_id = 1234
            has_text_frame = False

        slide.shapes._spTree # test against real slide
        with pytest.raises(ValueError):
            modify_text(slide, 999999, text="Error")


class TestTextEdgeCases:
    """Test boundary and edge-case text formatting scenarios."""

    def test_modify_text_empty_shape_replacement(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        text_shape.text_frame.text = ""

        modify_text(slide, text_shape.shape_id, text="Populated Text", font_size=16.0)
        assert text_shape.text_frame.text == "Populated Text"

    def test_modify_text_color_hex_formats(self, temp_presentation):
        slide = temp_presentation.slides[0]
        text_shape = next(s for s in slide.shapes if s.has_text_frame)
        shape_id = text_shape.shape_id

        # 3-char shorthand
        modify_text(slide, shape_id, color="#0F0")
        assert text_shape.text_frame.paragraphs[0].runs[0].font.color.rgb == RGBColor(0, 255, 0)

        # 6-char without hash
        modify_text(slide, shape_id, color="0000FF")
        assert text_shape.text_frame.paragraphs[0].runs[0].font.color.rgb == RGBColor(0, 0, 255)


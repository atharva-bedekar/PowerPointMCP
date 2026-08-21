"""Adversarial stress tests for Semantic Role Inference and Shape Inspection."""

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    ParagraphModel,
    PresentationModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    TextFrameModel,
    TextRunModel,
    TextStyle,
)
from powerpoint_mcp.pptx import (
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    map_shape_type,
)


class TestAdversarialSemanticRoleInference:
    """Stress test the 5-stage semantic role inference cascade."""

    @pytest.fixture
    def blank_slide_setup(self):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        return prs, slide, int(prs.slide_width), int(prs.slide_height)

    def test_infer_role_conflicting_name_vs_position(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        # Shape positioned at bottom (norm_top = 0.90) but named "Slide Title" with large font
        box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(0.5))
        box.name = "Slide Title"
        p = box.text_frame.paragraphs[0]
        p.text = "Bottom Title Note"
        p.font.size = Pt(28)

        # Rule 3C checks norm_top >= 0.85 -> FOOTER
        role = infer_semantic_role(box, slide_w, slide_h)
        assert role == SemanticRole.FOOTER

    def test_infer_role_offscreen_top_shape(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        # Offscreen top box (norm_top < 0) with 24pt font
        box = slide.shapes.add_textbox(Inches(1), Inches(-2.0), Inches(8), Inches(1.0))
        p = box.text_frame.paragraphs[0]
        p.text = "Offscreen Scratch Title"
        p.font.size = Pt(24)

        # norm_top < 0.22 and font >= 24 -> TITLE
        role = infer_semantic_role(box, slide_w, slide_h)
        assert role == SemanticRole.TITLE

    def test_infer_role_single_run_vs_paragraph_font_size(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        # Font size set on Run level rather than Paragraph level
        box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "Run-Level Title"
        r.font.size = Pt(30)  # On run, not paragraph

        role = infer_semantic_role(box, slide_w, slide_h)
        assert role == SemanticRole.TITLE

    def test_infer_role_multi_paragraph_mid_slide(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        # Multi-paragraph text in upper-mid slide (top = 1.0 in / 5.625 = 0.177)
        box = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(8), Inches(2.0))
        tf = box.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "Bullet item 1"
        p1.font.size = Pt(14)
        p2 = tf.add_paragraph()
        p2.text = "Bullet item 2"
        p2.font.size = Pt(14)

        role = infer_semantic_role(box, slide_w, slide_h)
        # norm_top < 0.38 and 14 <= max_font_size < 24 -> SUBTITLE
        # Wait, Rule 3B checks 0.15 <= norm_top < 0.38 and 14 <= max_font_size < 24 -> SUBTITLE
        assert role in (SemanticRole.SUBTITLE, SemanticRole.BODY)

    def test_infer_role_table_and_chart_types(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(4), Inches(2))
        assert infer_semantic_role(table_shape, slide_w, slide_h) == SemanticRole.TABLE

        # Connector
        from pptx.enum.shapes import MSO_CONNECTOR
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(3), Inches(3))
        assert infer_semantic_role(connector, slide_w, slide_h) == SemanticRole.DIAGRAM

    def test_infer_role_group_shapes(self, blank_slide_setup):
        prs, slide, slide_w, slide_h = blank_slide_setup

        # Group shape detection
        # Create shapes to group
        s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1))
        s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(1), Inches(1), Inches(1))
        group = slide.shapes.add_group_shape((s1, s2))

        assert map_shape_type(group) == ShapeType.GROUP
        assert infer_semantic_role(group, slide_w, slide_h) == SemanticRole.DIAGRAM

    def test_infer_role_zero_slide_dimensions(self):
        # Degenerate slide dimensions (0, 0) passed directly to infer_semantic_role
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = "Zero Dimension Test"
        box.text_frame.paragraphs[0].font.size = Pt(24)

        # Should not raise ZeroDivisionError and should infer a role safely
        role = infer_semantic_role(box, 0, 0)
        assert isinstance(role, SemanticRole)

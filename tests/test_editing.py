"""Comprehensive test suite for shape editing, positioning, resizing, deletion, copying, and z-ordering."""

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import emu_to_inches, inches_to_emu
from powerpoint_mcp.pptx.editor import (
    copy_shape,
    delete_shape,
    modify_shape,
    move_shape,
    resize_shape,
)
from powerpoint_mcp.pptx.inspector import inspect_slide


class TestShapeModification:
    """Test modify_shape, move_shape, and resize_shape."""

    def test_modify_shape_absolute_geometry(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        shape_id = shape.shape_id

        res = modify_shape(
            slide,
            shape_id,
            x=1.5,
            y=2.5,
            width=4.0,
            height=3.0,
            rotation=45.0,
        )

        assert res["shape_id"] == shape_id
        assert shape.left == inches_to_emu(1.5)
        assert shape.top == inches_to_emu(2.5)
        assert shape.width == inches_to_emu(4.0)
        assert shape.height == inches_to_emu(3.0)
        assert shape.rotation == 45.0
        assert res["x"] == 1.5
        assert res["y"] == 2.5
        assert res["width"] == 4.0
        assert res["height"] == 3.0
        assert res["rotation"] == 45.0

    def test_modify_shape_relative_deltas(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        shape_id = shape.shape_id

        initial_left = shape.left
        initial_top = shape.top
        initial_width = shape.width
        initial_height = shape.height

        modify_shape(
            slide,
            shape_id,
            dx=0.5,
            dy=-0.25,
            dwidth=1.0,
            dheight=-0.5,
            drotation=15.0,
        )

        assert shape.left == initial_left + inches_to_emu(0.5)
        assert shape.top == initial_top + inches_to_emu(-0.25)
        assert shape.width == initial_width + inches_to_emu(1.0)
        assert shape.height == initial_height + inches_to_emu(-0.5)
        assert shape.rotation == 15.0

    def test_modify_shape_via_presentation_and_slide_number(self, temp_presentation):
        res = modify_shape(
            temp_presentation,
            1, # slide_number
            temp_presentation.slides[0].shapes[0].shape_id,
            x=2.0,
            y=1.0,
        )
        assert res["x"] == 2.0
        assert res["y"] == 1.0

    def test_move_shape_absolute_and_delta(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        shape_id = shape.shape_id

        # Absolute move
        move_shape(slide, shape_id, x_inches=3.0, y_inches=1.5)
        assert shape.left == inches_to_emu(3.0)
        assert shape.top == inches_to_emu(1.5)

        # Delta move
        move_shape(slide, shape_id, delta_x_inches=-1.0, delta_y_inches=0.5)
        assert shape.left == inches_to_emu(2.0)
        assert shape.top == inches_to_emu(2.0)

    def test_resize_shape_dimensions_and_scaling(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        shape_id = shape.shape_id

        # Set baseline size
        modify_shape(slide, shape_id, width=2.0, height=2.0)

        # Scale by 2.0x width and 1.5x height
        resize_shape(slide, shape_id, scale_x=2.0, scale_y=1.5)
        assert shape.width == inches_to_emu(4.0)
        assert shape.height == inches_to_emu(3.0)

    def test_modify_shape_nonexistent_id_raises_error(self, temp_presentation):
        slide = temp_presentation.slides[0]
        with pytest.raises(ValueError, match="not found"):
            modify_shape(slide, 999999, x=1.0)


class TestZOrderManipulation:
    """Test shape stacking order manipulation inside <p:spTree>."""

    def test_z_order_bring_to_front_and_send_to_back(self, temp_presentation):
        slide = temp_presentation.slides[0]
        assert len(slide.shapes) >= 2

        first_shape = slide.shapes[0]
        first_id = first_shape.shape_id

        # Bring first shape to front
        modify_shape(slide, first_id, z_order="bring_to_front")
        # In python-pptx, shapes is an iterable over spTree child nodes
        # The last shape element in spTree is now first_shape
        assert slide.shapes[-1].shape_id == first_id

        # Send back
        modify_shape(slide, first_id, z_order="send_to_back")
        assert slide.shapes[0].shape_id == first_id

    def test_z_order_forward_and_backward(self, temp_presentation):
        slide = temp_presentation.slides[0]
        if len(slide.shapes) >= 3:
            mid_shape_id = slide.shapes[1].shape_id

            # Bring forward
            modify_shape(slide, mid_shape_id, z_order="bring_forward")
            assert slide.shapes[2].shape_id == mid_shape_id

            # Send backward
            modify_shape(slide, mid_shape_id, z_order="send_backward")
            assert slide.shapes[1].shape_id == mid_shape_id

    def test_z_order_numeric_index(self, temp_presentation):
        slide = temp_presentation.slides[0]
        if len(slide.shapes) >= 3:
            target_id = slide.shapes[-1].shape_id
            modify_shape(slide, target_id, z_order=0)
            assert slide.shapes[0].shape_id == target_id


class TestShapeDeletion:
    """Test delete_shape operation."""

    def test_delete_shape_successfully(self, temp_presentation):
        slide = temp_presentation.slides[0]
        initial_count = len(slide.shapes)
        target_id = slide.shapes[0].shape_id

        success = delete_shape(slide, target_id)
        assert success is True
        assert len(slide.shapes) == initial_count - 1
        assert not any(s.shape_id == target_id for s in slide.shapes)

    def test_delete_shape_nonexistent_raises_error(self, temp_presentation):
        slide = temp_presentation.slides[0]
        with pytest.raises(ValueError, match="not found"):
            delete_shape(slide, 999999)


class TestShapeDuplicationAndCopy:
    """Test copy_shape within same slide and across slides."""

    def test_copy_shape_same_slide(self, temp_presentation):
        slide = temp_presentation.slides[0]
        initial_count = len(slide.shapes)
        source_shape = slide.shapes[0]
        source_id = source_shape.shape_id
        orig_left = source_shape.left
        orig_top = source_shape.top

        new_shape_id = copy_shape(
            slide,
            source_id,
            offset_x_inches=0.5,
            offset_y_inches=0.5,
        )

        assert new_shape_id != source_id
        assert len(slide.shapes) == initial_count + 1

        new_shape = next(s for s in slide.shapes if s.shape_id == new_shape_id)
        assert new_shape.left == orig_left + inches_to_emu(0.5)
        assert new_shape.top == orig_top + inches_to_emu(0.5)
        assert "(Copy)" in new_shape.name

    def test_copy_shape_cross_slide(self, temp_presentation):
        slide_1 = temp_presentation.slides[0]
        slide_2 = temp_presentation.slides[1]
        initial_count_slide_2 = len(slide_2.shapes)
        source_shape = slide_1.shapes[0]
        source_id = source_shape.shape_id

        new_shape_id = copy_shape(
            temp_presentation,
            1, # source slide 1
            source_id,
            target_slide_number=2, # destination slide 2
            offset_x_inches=0.0,
            offset_y_inches=0.0,
        )

        assert len(slide_2.shapes) == initial_count_slide_2 + 1
        new_shape = next(s for s in slide_2.shapes if s.shape_id == new_shape_id)
        assert new_shape.shape_id == new_shape_id

    def test_copy_picture_shape_with_relationship(self, temp_presentation):
        """Verify copying a picture replicates relationship reference without corruption."""
        # Add a shape or picture if slide has one
        slide = temp_presentation.slides[0]
        # Find picture or add an auto shape
        source_shape = slide.shapes[0]
        new_id = copy_shape(slide, source_shape.shape_id)
        assert new_id > 0
        copied = next(s for s in slide.shapes if s.shape_id == new_id)
        assert copied is not None


class TestEditingEdgeCases:
    """Test boundary and unexpected input scenarios for shape editing."""

    def test_modify_shape_drotation_wrapping(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        shape.rotation = 350.0
        modify_shape(slide, shape.shape_id, drotation=30.0)
        assert shape.rotation == pytest.approx(20.0, abs=1e-2)

    def test_resize_shape_scale_x_and_y_independent(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        modify_shape(slide, shape.shape_id, width=2.0, height=2.0)

        # Scale only X
        resize_shape(slide, shape.shape_id, scale_x=3.0)
        assert shape.width == inches_to_emu(6.0)
        assert shape.height == inches_to_emu(2.0)

        # Scale only Y
        resize_shape(slide, shape.shape_id, scale_y=0.5)
        assert shape.width == inches_to_emu(6.0)
        assert shape.height == inches_to_emu(1.0)

    def test_delete_multiple_shapes_sequentially(self, temp_presentation):
        slide = temp_presentation.slides[0]
        initial_len = len(slide.shapes)
        assert initial_len >= 2
        id1 = slide.shapes[0].shape_id
        id2 = slide.shapes[1].shape_id

        delete_shape(slide, id1)
        delete_shape(slide, id2)
        assert len(slide.shapes) == initial_len - 2


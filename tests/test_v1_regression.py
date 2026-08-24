"""Regression test suite for PowerPoint MCP v1 fixes and integrity enhancements.

Verifies:
1. Target resolution with active session and explicit original presentation_path.
2. Target resolution with active session and omitted presentation_path.
3. Conflicting presentation_path during active session.
4. Integrity checks during save and save_as.
5. Save-as preserves original while writing to new output.
6. Concise inspection and validation summary modes.
"""

from pathlib import Path
import pytest
from pptx import Presentation

from powerpoint_mcp.server import app
from powerpoint_mcp.tools.editing import (
    ppt_modify_shape,
    ppt_modify_text,
    ppt_move_shape,
)
from powerpoint_mcp.tools.inspection import (
    ppt_inspect_presentation,
    ppt_inspect_slide,
    ppt_validate_slide,
)
from powerpoint_mcp.tools.versioning import (
    SessionManager,
    compute_file_hash,
    get_current_session,
    get_session_manager,
    open_presentation,
    resolve_active_target,
    revert_session,
    save_as,
    save_session,
)


def _extract_tool_result(res):
    """Helper to parse MCP call result."""
    import json
    if hasattr(res, "content") and res.content:
        first = res.content[0]
        if hasattr(first, "text"):
            try:
                return json.loads(first.text)
            except Exception:
                return {"text": first.text}
    if isinstance(res, dict):
        return res
    return {}


class TestActiveSessionPathResolution:
    """Tests for the primary v1 bug where passing original presentation_path bypassed working copy."""

    def test_explicit_original_path_targets_working_copy_and_persists_on_save(
        self, temp_deck_path: Path, clean_ppt_env
    ):
        """1. Open source presentation -> create session.
        2. Modify using explicit original presentation_path.
        3. Save.
        4. Verify that saved result contains modification and was not silently lost.
        """
        # Step 1: Open presentation
        session = open_presentation(temp_deck_path)
        assert session.session_id is not None
        orig_src_hash = compute_file_hash(temp_deck_path)

        # Step 2: Modify text passing explicit original presentation_path
        mod_res = ppt_modify_text(
            slide_number=1,
            shape_id=2,
            text="UPDATED TITLE VIA ORIGINAL PATH",
            presentation_path=str(temp_deck_path),
        )
        assert mod_res.get("success") is True

        # Verify mutation was applied to working copy, NOT directly to source file yet
        working_prs = Presentation(session.working_path)
        working_title = working_prs.slides[0].shapes[0].text_frame.text
        assert "UPDATED TITLE VIA ORIGINAL PATH" in working_title

        # Verify source file before save was NOT modified
        assert compute_file_hash(temp_deck_path) == orig_src_hash

        # Step 3: Save session
        save_res = save_session()
        assert save_res.get("success") is True
        assert save_res.get("backup_path") is not None
        assert Path(save_res["backup_path"]).exists()

        # Step 4: Verify saved file contains the edits
        saved_prs = Presentation(str(temp_deck_path))
        saved_title = saved_prs.slides[0].shapes[0].text_frame.text
        assert "UPDATED TITLE VIA ORIGINAL PATH" in saved_title

    def test_omitted_presentation_path_targets_working_copy_and_persists_on_save(
        self, temp_deck_path: Path, clean_ppt_env
    ):
        """1. Open source presentation -> create session.
        2. Modify while omitting presentation_path.
        3. Save.
        4. Verify modification exists in saved output.
        """
        session = open_presentation(temp_deck_path)

        # Modify omitting presentation_path
        mod_res = ppt_modify_shape(
            slide_number=1,
            shape_id=2,
            dx=0.5,
            dy=0.25,
            presentation_path=None,
        )
        assert mod_res.get("success") is True

        # Save session
        save_res = save_session()
        assert save_res.get("success") is True

        # Verify saved file has new position
        saved_prs = Presentation(str(temp_deck_path))
        title_shape = saved_prs.slides[0].shapes[0]
        assert title_shape.left > 1.0  # shifted by dx=0.5

    def test_conflicting_presentation_path_raises_clear_error(
        self, temp_deck_path: Path, tmp_path: Path, clean_ppt_env
    ):
        """Active session rejects mutation targeting an unrelated file path."""
        open_presentation(temp_deck_path)

        # Create a second dummy presentation
        other_path = tmp_path / "other_deck.pptx"
        prs2 = Presentation()
        prs2.slides.add_slide(prs2.slide_layouts[6])
        prs2.save(str(other_path))

        # Attempting to mutate other_path while session is active for temp_deck_path must fail
        with pytest.raises(ValueError, match="Conflicting presentation_path"):
            resolve_active_target(
                presentation_path=str(other_path),
                mutation=True,
            )


class TestSaveAsAndIntegrity:
    """Tests for save_as non-destructiveness and session state tracking."""

    def test_save_as_preserves_original_file(
        self, temp_deck_path: Path, tmp_path: Path, clean_ppt_env
    ):
        """save_as creates separate output file and leaves original intact."""
        orig_bytes = temp_deck_path.read_bytes()
        session = open_presentation(temp_deck_path)

        # Edit working copy
        ppt_modify_text(
            slide_number=1,
            shape_id=2,
            text="SAVED AS NEW DECK",
        )

        out_file = tmp_path / "exported_deck.pptx"
        res = save_as(output_path=out_file)
        assert res.get("success") is True
        assert out_file.exists()

        # Output file has changes
        out_prs = Presentation(str(out_file))
        assert "SAVED AS NEW DECK" in out_prs.slides[0].shapes[0].text_frame.text

        # Original source file is untouched
        assert temp_deck_path.read_bytes() == orig_bytes

    def test_session_tracks_mutations_and_hashes(
        self, temp_deck_path: Path, clean_ppt_env
    ):
        """Session object tracks mutation count and hashes."""
        session = open_presentation(temp_deck_path)
        assert session.mutation_count == 0
        assert session.initial_working_hash is not None

        # Perform 2 mutations
        ppt_modify_shape(slide_number=1, shape_id=2, dx=0.1)
        ppt_modify_shape(slide_number=1, shape_id=2, dy=0.1)

        curr = get_current_session()
        assert curr.mutation_count == 2
        assert curr.last_working_hash != curr.initial_working_hash


class TestConciseInspectionAndValidation:
    """Tests for detail='summary' vs detail='full' outputs."""

    def test_inspect_slide_summary_mode_is_concise(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        """ppt_inspect_slide defaults to concise summary representation."""
        open_presentation(synthetic_deck_path)

        summary_res = ppt_inspect_slide(slide_number=1, detail="summary")
        assert summary_res.get("success") is True
        assert summary_res.get("detail") == "summary"
        shapes = summary_res.get("shapes", [])
        assert len(shapes) > 0

        # Summary shapes contain high-level fields but omit raw paragraphs/runs collections
        first = shapes[0]
        assert "shape_id" in first
        assert "name" in first
        assert "semantic_role" in first
        assert "x" in first
        assert "y" in first
        assert "width" in first
        assert "height" in first
        assert "paragraphs" not in first  # omitted in summary

    def test_inspect_slide_full_mode_has_deep_data(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        """ppt_inspect_slide with detail='full' includes full text frames and runs."""
        open_presentation(synthetic_deck_path)

        full_res = ppt_inspect_slide(slide_number=1, detail="full")
        assert full_res.get("success") is True
        assert full_res.get("detail") == "full"
        first = full_res["shapes"][0]
        assert "text_frame" in first
        assert "paragraphs" in first["text_frame"]

    def test_validate_slide_summary_counts(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        """ppt_validate_slide returns structured summary counts."""
        open_presentation(synthetic_deck_path)

        # Slide 3 has intentional defects
        val_res = ppt_validate_slide(slide_number=3, detail="summary")
        assert val_res.get("success") is True
        assert "summary" in val_res
        assert "valid" in val_res
        assert val_res["valid"] is False
        summary = val_res["summary"]
        assert "overlaps" in summary
        assert "boundary_violations" in summary
        assert "text_overflow" in summary
        assert "tiny_fonts" in summary
        assert summary["overlaps"] >= 1


@pytest.mark.asyncio
class TestMCPToolEndToEndV1:
    """Async MCP protocol end-to-end tests for the fixed workflow."""

    async def test_mcp_full_edit_lifecycle_with_explicit_path(
        self, temp_deck_path: Path, clean_ppt_env
    ):
        """Test full MCP tool call cycle with explicit path, inspect, mutate, validate, save."""
        # 1. ppt_open
        open_res = await app.call_tool("ppt_open", {"presentation_path": str(temp_deck_path)})
        open_data = _extract_tool_result(open_res)
        assert open_data.get("success") is True

        # 2. ppt_inspect_slide
        insp_res = await app.call_tool("ppt_inspect_slide", {"slide_number": 1})
        insp_data = _extract_tool_result(insp_res)
        assert insp_data.get("success") is True
        assert insp_data.get("detail") == "summary"

        # 3. ppt_modify_text passing explicit presentation_path
        mod_res = await app.call_tool(
            "ppt_modify_text",
            {
                "slide_number": 1,
                "shape_id": 2,
                "text": "E2E VERIFIED TITLE FIX",
                "presentation_path": str(temp_deck_path),
            },
        )
        mod_data = _extract_tool_result(mod_res)
        assert mod_data.get("success") is True

        # 4. ppt_validate_slide
        val_res = await app.call_tool("ppt_validate_slide", {"slide_number": 1})
        val_data = _extract_tool_result(val_res)
        assert val_data.get("success") is True
        assert val_data.get("valid") is True

        # 5. ppt_save
        save_res = await app.call_tool("ppt_save", {})
        save_data = _extract_tool_result(save_res)
        assert save_data.get("success") is True

        # 6. Verify edits are in destination file
        final_prs = Presentation(str(temp_deck_path))
        assert "E2E VERIFIED TITLE FIX" in final_prs.slides[0].shapes[0].text_frame.text

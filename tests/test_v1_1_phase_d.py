"""Unit tests for PowerPoint MCP v1.1 Phase D: Style Inheritance & Presets."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.pptx.styles import extract_complete_shape_style
from powerpoint_mcp.tools.editing import ppt_apply_style
from powerpoint_mcp.tools.versioning import open_presentation


@pytest.fixture
def style_deck(tmp_path: Path) -> Path:
    """Create a test deck with shapes for style transfer and presets."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Shape 1: Styled Source Card
    s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(3.0), Inches(2.0))
    s1.name = "Source Card"
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(220, 240, 255)
    s1.line.color.rgb = RGBColor(50, 120, 220)
    s1.line.width = Pt(2.0)
    s1.text_frame.paragraphs[0].text = "Source Content"
    s1.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    s1.text_frame.paragraphs[0].runs[0].font.bold = True
    s1.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(20, 50, 100)

    # Shape 2: Plain Unstyled Card
    s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.0), Inches(3.0), Inches(2.0))
    s2.name = "Target Card"
    s2.text_frame.paragraphs[0].text = "Target Content"

    # Shape 3: Plain Badge Box
    s3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.0), Inches(1.5), Inches(0.4))
    s3.name = "Badge Target"
    s3.text_frame.paragraphs[0].text = "APPROVED"

    deck_path = tmp_path / "style_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_apply_style_preset_card_accent(style_deck: Path):
    """Test applying standard card_accent preset."""
    session = open_presentation(str(style_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    target_id = shapes[1].shape_id

    res = ppt_apply_style(slide_number=1, shape_id=target_id, preset="card_accent")
    assert res["success"] is True
    assert res["shapes_styled_count"] == 1

    # Reload and verify style applied
    prs = Presentation(session.working_path)
    styled_shape = next(s for s in prs.slides[0].shapes if s.shape_id == target_id)
    # Check text preserved
    assert styled_shape.text_frame.text == "Target Content"
    # Check fill
    style_snap = extract_complete_shape_style(styled_shape)
    assert style_snap["fill_color"] == "EFF6FF"
    assert style_snap["line_color"] == "3B82F6"


def test_apply_style_preset_badge_success(style_deck: Path):
    """Test applying badge_success preset."""
    session = open_presentation(str(style_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    badge_id = shapes[2].shape_id

    res = ppt_apply_style(slide_number=1, shape_id=badge_id, preset="badge_success")
    assert res["success"] is True

    prs = Presentation(session.working_path)
    styled_shape = next(s for s in prs.slides[0].shapes if s.shape_id == badge_id)
    assert styled_shape.text_frame.text == "APPROVED"
    style_snap = extract_complete_shape_style(styled_shape)
    assert style_snap["fill_color"] == "DEF7EC"
    assert style_snap["line_color"] == "31C48D"


def test_style_inheritance_from_source_shape(style_deck: Path):
    """Test transferring style from a source shape onto a target shape."""
    session = open_presentation(str(style_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    source_id, target_id = shapes[0].shape_id, shapes[1].shape_id

    res = ppt_apply_style(slide_number=1, shape_id=target_id, source_shape_id=source_id)
    assert res["success"] is True

    prs = Presentation(session.working_path)
    styled_shape = next(s for s in prs.slides[0].shapes if s.shape_id == target_id)
    assert styled_shape.text_frame.text == "Target Content"
    style_snap = extract_complete_shape_style(styled_shape)
    assert style_snap["fill_color"] == "DCF0FF"
    assert style_snap["line_color"] == "3278DC"


def test_batch_apply_style_with_overrides(style_deck: Path):
    """Test applying a preset to multiple shapes with an explicit color override."""
    session = open_presentation(str(style_deck))
    shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id = shapes[0].shape_id, shapes[1].shape_id

    res = ppt_apply_style(
        slide_number=1,
        shape_ids=[s1_id, s2_id],
        preset="card_default",
        fill_color="#FAFAFA",  # Override preset fill
    )
    assert res["success"] is True
    assert res["shapes_styled_count"] == 2

    prs = Presentation(session.working_path)
    for sid in (s1_id, s2_id):
        s = next(x for x in prs.slides[0].shapes if x.shape_id == sid)
        style_snap = extract_complete_shape_style(s)
        assert style_snap["fill_color"] == "FAFAFA"

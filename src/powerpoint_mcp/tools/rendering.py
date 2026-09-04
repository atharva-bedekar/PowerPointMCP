"""Rendering and visual comparison tools for PowerPoint MCP server."""

from datetime import datetime, timezone
from functools import wraps
import os
from pathlib import Path
import traceback
from typing import Any, Dict, List, Optional, Union

from PIL import Image, ImageDraw, ImageFont

from powerpoint_mcp.pptx.inspector import inspect_slide
from powerpoint_mcp.rendering.image_diff import VisualDiffResult, visual_diff
from powerpoint_mcp.rendering.renderer import (
    BaseRenderer,
    NullRenderer,
    PowerPointRenderer,
    get_available_renderer,
)
from powerpoint_mcp.tools.inspection import handle_tool_errors
from powerpoint_mcp.tools.versioning import get_session_manager, resolve_active_target
from powerpoint_mcp.utils.paths import get_session_renders_dir


def _iso_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _resolve_presentation_path(presentation_path: Optional[str] = None) -> str:
    """Resolve presentation path from argument or active session using resolve_active_target."""
    target_path, _ = resolve_active_target(
        presentation_path=presentation_path,
        require_session=False,
        mutation=False,
    )
    return target_path


def _get_presentation_dimensions(target_path: str, dpi: int) -> tuple[int, int]:
    """Dynamically derive slide dimensions in pixels from the presentation preserving exact aspect ratio.

    Supports 16:9 widescreen, 4:3 standard, portrait, and custom slide geometries.
    """
    from pptx import Presentation
    from powerpoint_mcp.models.shape import emu_to_inches

    w_in, h_in = 13.333, 7.5
    try:
        prs = Presentation(target_path)
        if prs.slide_width and prs.slide_height:
            w_in = emu_to_inches(prs.slide_width)
            h_in = emu_to_inches(prs.slide_height)
    except Exception:
        pass

    width_px = int(round(w_in * dpi))
    height_px = int(round(h_in * dpi))
    return width_px, height_px


def _render_pillow_fallback(
    presentation_path: str,
    slide_number: int,
    output_path: Path,
    width: int = 1920,
    height: int = 1080,
) -> str:
    """Fallback programmatic slide visualizer using Pillow when COM/LibreOffice are unavailable."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    try:
        slide_model = inspect_slide(presentation_path, slide_number)
        scale_x = width / max(slide_model.width_inches, 1.0)
        scale_y = height / max(slide_model.height_inches, 1.0)

        # Draw slide background frame
        draw.rectangle([0, 0, width - 1, height - 1], outline=(200, 200, 200), width=2)

        # Draw shapes
        for shape in slide_model.shapes:
            sx = int(round(shape.bounds.left_inches * scale_x))
            sy = int(round(shape.bounds.top_inches * scale_y))
            sw = int(round(shape.bounds.width_inches * scale_x))
            sh = int(round(shape.bounds.height_inches * scale_y))

            box = [sx, sy, sx + sw, sy + sh]

            # Background color
            fill_color = (235, 240, 248)
            if shape.fill and shape.fill.get("color_rgb"):
                hex_c = shape.fill["color_rgb"].lstrip("#")
                if len(hex_c) == 6:
                    try:
                        fill_color = (int(hex_c[:2], 16), int(hex_c[2:4], 16), int(hex_c[4:], 16))
                    except Exception:
                        pass

            draw.rectangle(box, fill=fill_color, outline=(100, 120, 150), width=2)

            # Draw label / text
            text = (shape.text_frame.text if shape.text_frame else shape.name or "").strip()
            if text:
                draw.text((sx + 8, sy + 8), text[:80], fill=(20, 20, 20))

    except Exception:
        # Fallback text banner
        draw.rectangle([50, 50, width - 50, height - 50], outline=(180, 180, 180), width=2)
        draw.text((100, 100), f"Slide {slide_number} (Pillow Synthetic Render)", fill=(50, 50, 50))

    img.save(str(output_path), format="PNG")
    return str(output_path)


@handle_tool_errors
def ppt_render_slide(
    slide_number: int,
    output_dir: Optional[str] = None,
    output_path: Optional[str] = None,
    renderer: str = "auto",
    dpi: int = 150,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Render a single slide to a high-resolution PNG image.

    Uses Windows PowerPoint COM automation when available, falling back to LibreOffice headless or Pillow.

    Args:
        slide_number: 1-indexed slide number.
        output_dir: Output directory path. If omitted, saves inside active session renders directory.
        output_path: Direct output PNG file path.
        renderer: Preferred renderer engine ('auto', 'powerpoint', 'libreoffice', 'none', 'mock').
        dpi: Render resolution DPI (default 150).
        presentation_path: Path to presentation file. If omitted, uses active session.

    Returns:
        Structured dictionary containing rendered image_path, slide_number, renderer name, and pixel dimensions.
    """
    if slide_number < 1:
        raise IndexError(f"Slide number must be >= 1, got {slide_number}")

    target_path = _resolve_presentation_path(presentation_path)
    mgr = get_session_manager()
    session = mgr.get_current_session()

    # Determine destination image path
    if output_path:
        out_file = Path(output_path).resolve()
    elif output_dir:
        out_file = Path(output_dir).resolve() / f"slide_{slide_number}.png"
    elif session:
        renders_dir = get_session_renders_dir(session.session_id, mgr.workspace_dir.parent)
        renders_dir.mkdir(parents=True, exist_ok=True)
        out_file = renders_dir / f"slide_{slide_number}.png"
    else:
        out_file = Path("./renders").resolve() / f"slide_{slide_number}.png"

    out_file.parent.mkdir(parents=True, exist_ok=True)

    # Dynamically derive dimensions from presentation (preserves 16:9, 4:3, portrait, etc.)
    width_px, height_px = _get_presentation_dimensions(target_path, dpi)

    import time
    t0 = time.perf_counter()

    # Cache check: if output file exists and is newer than target presentation
    is_cached = False
    if out_file.exists() and os.path.exists(target_path):
        target_mtime = os.path.getmtime(target_path)
        out_mtime = os.path.getmtime(str(out_file))
        if out_mtime >= target_mtime and out_file.stat().st_size > 0:
            is_cached = True
            render_time_ms = round((time.perf_counter() - t0) * 1000, 2)
            return {
                "success": True,
                "image_path": str(out_file),
                "slide_number": slide_number,
                "renderer": "cache",
                "cached": True,
                "render_time_ms": render_time_ms,
                "width_px": width_px,
                "height_px": height_px,
                "dimensions": {"width_px": width_px, "height_px": height_px},
            }

    renderer_used = "powerpoint_com"
    renderer_obj = get_available_renderer(preferred=renderer)

    if renderer.lower() in ("mock", "pillow"):
        _render_pillow_fallback(target_path, slide_number, out_file, width_px, height_px)
        renderer_used = "mock"
    elif renderer_obj.is_available:
        try:
            renderer_obj.render_slide(
                target_path,
                slide_number,
                out_file,
                width=width_px,
                height=height_px,
            )
            renderer_used = renderer_obj.renderer_name
        except Exception as exc:
            # Fallback to pillow render
            _render_pillow_fallback(target_path, slide_number, out_file, width_px, height_px)
            renderer_used = f"{renderer_obj.renderer_name}_fallback"
    else:
        _render_pillow_fallback(target_path, slide_number, out_file, width_px, height_px)
        renderer_used = "pillow_fallback"

    render_time_ms = round((time.perf_counter() - t0) * 1000, 2)

    # Update session renders record
    if session:
        session.renders.append({
            "slide_number": slide_number,
            "render_path": str(out_file),
            "renderer": renderer_used,
            "timestamp": _iso_now(),
        })
        session.save_metadata()

    return {
        "success": True,
        "image_path": str(out_file),
        "slide_number": slide_number,
        "renderer": renderer_used,
        "cached": False,
        "render_time_ms": render_time_ms,
        "width_px": width_px,
        "height_px": height_px,
        "dimensions": {"width_px": width_px, "height_px": height_px},
    }


@handle_tool_errors
def ppt_render_slides(
    slide_numbers: List[int],
    output_dir: Optional[str] = None,
    renderer: str = "auto",
    dpi: int = 150,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Render a specified list of slides in batch to PNG images.

    Args:
        slide_numbers: List of 1-indexed slide numbers to render (e.g. [3, 4, 5, 6]).
        output_dir: Target output directory for PNG slide renders.
        renderer: Renderer preference ('auto', 'powerpoint', 'libreoffice', 'mock').
        dpi: Output resolution DPI (default 150).
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured dictionary listing rendered slide numbers, paths, and images map.
    """
    if not slide_numbers:
        raise ValueError("slide_numbers must contain at least one slide number")

    target_path = _resolve_presentation_path(presentation_path)
    mgr = get_session_manager()
    session = mgr.get_current_session()

    if output_dir:
        out_dir_path = Path(output_dir).resolve()
    elif session:
        out_dir_path = get_session_renders_dir(session.session_id, mgr.workspace_dir.parent)
    else:
        out_dir_path = Path("./renders").resolve()

    out_dir_path.mkdir(parents=True, exist_ok=True)

    rendered_slides: List[Dict[str, Any]] = []
    images_map: Dict[int, str] = {}
    renderer_name = "powerpoint_com"

    for s_num in slide_numbers:
        if s_num < 1:
            raise IndexError(f"Slide number must be >= 1, got {s_num}")

    renderer_inst = get_available_renderer(renderer)
    width_px, height_px = _get_presentation_dimensions(target_path, dpi)

    rendered_map: Dict[int, str] = {}
    if renderer.lower() in ("mock", "pillow"):
        for s_num in slide_numbers:
            out_file = out_dir_path / f"slide_{s_num}.png"
            _render_pillow_fallback(target_path, s_num, out_file, width_px, height_px)
            rendered_map[s_num] = str(out_file)
        renderer_name = "mock"
    elif renderer_inst.is_available and hasattr(renderer_inst, "render_slides"):
        try:
            rendered_map = renderer_inst.render_slides(
                presentation_path=target_path,
                slide_numbers=slide_numbers,
                output_dir=out_dir_path,
                width=width_px,
                height=height_px,
            )
            renderer_name = renderer_inst.renderer_name
        except Exception:
            rendered_map = {}

    if not rendered_map:
        for s_num in slide_numbers:
            res = ppt_render_slide(
                slide_number=s_num,
                output_dir=str(out_dir_path),
                renderer=renderer,
                dpi=dpi,
                presentation_path=target_path,
            )
            if res.get("success"):
                img_path = res.get("image_path")
                rendered_map[s_num] = img_path
                renderer_name = res.get("renderer", renderer_name)

    for s_num in slide_numbers:
        img_p = rendered_map.get(s_num)
        if img_p:
            rendered_slides.append({"slide_number": s_num, "image_path": img_p})
            images_map[s_num] = img_p

    return {
        "success": True,
        "slide_count": len(rendered_slides),
        "slide_numbers": slide_numbers,
        "rendered_slides": rendered_slides,
        "images": images_map,
        "renderer": renderer_name,
    }


@handle_tool_errors
def ppt_render_presentation(
    output_dir: Optional[str] = None,
    renderer: str = "auto",
    dpi: int = 150,
    presentation_path: Optional[str] = None,
) -> Dict[str, Any]:
    """Render all slides in a presentation to PNG images.

    Args:
        output_dir: Target output directory for PNG slide renders.
        renderer: Renderer preference ('auto', 'powerpoint', 'libreoffice', 'mock').
        dpi: Output resolution DPI (default 150).
        presentation_path: Path to presentation. If omitted, uses active session.

    Returns:
        Structured dictionary listing rendered slide paths and total slide count.
    """
    target_path = _resolve_presentation_path(presentation_path)
    mgr = get_session_manager()
    session = mgr.get_current_session()

    if output_dir:
        out_dir_path = Path(output_dir).resolve()
    elif session:
        out_dir_path = get_session_renders_dir(session.session_id, mgr.workspace_dir.parent)
    else:
        out_dir_path = Path("./renders").resolve()

    out_dir_path.mkdir(parents=True, exist_ok=True)

    from pptx import Presentation
    prs = Presentation(target_path)
    slide_count = len(prs.slides)

    rendered_slides: List[Dict[str, Any]] = []
    renderer_name = "powerpoint_com"

    for i in range(1, slide_count + 1):
        res = ppt_render_slide(
            slide_number=i,
            output_dir=str(out_dir_path),
            renderer=renderer,
            dpi=dpi,
            presentation_path=target_path,
        )
        if res.get("success"):
            rendered_slides.append({
                "slide_number": i,
                "image_path": res.get("image_path"),
            })
            renderer_name = res.get("renderer", renderer_name)

    return {
        "success": True,
        "slide_count": len(rendered_slides),
        "rendered_slides": rendered_slides,
        "renderer": renderer_name,
    }


@handle_tool_errors
def ppt_visual_diff(
    before_image: str,
    after_image: str,
    output_diff_path: Optional[str] = None,
    threshold: float = 0.1,
) -> Dict[str, Any]:
    """Perform deterministic pixel-level comparison between two rendered slide images.

    Generates a difference heat map overlay image, computes bounding boxes of change, and calculates similarity percentage.

    Args:
        before_image: Path to baseline / before PNG image.
        after_image: Path to modified / after PNG image.
        output_diff_path: Optional destination path for difference overlay image.
        threshold: Pixel difference sensitivity threshold (0.0 to 1.0 or 0 to 255).

    Returns:
        Structured dictionary detailing similarity percentage, changed pixel count, and changed bounding boxes.
    """
    p_before = Path(before_image).resolve()
    p_after = Path(after_image).resolve()

    if not p_before.exists():
        raise FileNotFoundError(f"Before image not found: {p_before}")
    if not p_after.exists():
        raise FileNotFoundError(f"After image not found: {p_after}")

    thresh_int = int(round(threshold * 255.0)) if threshold <= 1.0 else int(threshold)

    diff_result: VisualDiffResult = visual_diff(
        str(p_before),
        str(p_after),
        diff_output_path=output_diff_path,
        threshold=thresh_int,
    )

    return {
        "success": True,
        **diff_result.to_dict(),
    }

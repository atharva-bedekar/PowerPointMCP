"""Unit tests for PowerPoint MCP v1.1 Phase A: Relative Typography & Advanced Text Validation."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import (
    BoundingBox,
    ParagraphModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    TextFrameModel,
    TextRunModel,
    TextStyle,
)
from powerpoint_mcp.models.slide import SlideModel
from powerpoint_mcp.pptx.editor import modify_text, scale_slide_typography
from powerpoint_mcp.tools.editing import (
    ppt_batch_modify_text,
    ppt_modify_text,
    ppt_scale_slide_typography,
)
from powerpoint_mcp.tools.versioning import open_presentation
from powerpoint_mcp.utils.validation import IssueSeverity, validate_slide


@pytest.fixture
def typography_deck(tmp_path: Path) -> Path:
    """Create a test deck with multiple text shapes and diverse font sizes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Shape 1: Title (28pt)
    tb1 = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.0), Inches(1.0))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "Presentation Main Title"
    p1.runs[0].font.size = Pt(28)

    # Shape 2: Subtitle (18pt)
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(10.0), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Strategic Platform Overview"
    p2.runs[0].font.size = Pt(18)

    # Shape 3: Card Body (12pt)
    tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(4.0), Inches(2.5))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "First Bullet Point"
    p3.runs[0].font.size = Pt(12)
    p3_2 = tb3.text_frame.add_paragraph()
    p3_2.text = "Second Bullet Point"
    p3_2.runs[0].font.size = Pt(10)

    # Shape 4: Badge/Pill (7pt)
    tb4 = slide.shapes.add_textbox(Inches(10.5), Inches(0.5), Inches(1.8), Inches(0.35))
    tb4.name = "Status Badge"
    p4 = tb4.text_frame.paragraphs[0]
    p4.text = "ACTIVE RELEASE"
    p4.runs[0].font.size = Pt(7)

    deck_path = tmp_path / "typography_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_relative_font_size_delta(typography_deck: Path):
    """Test modifying text font size using positive and negative deltas."""
    session = open_presentation(str(typography_deck))
    prs = Presentation(session.working_path)
    slide = prs.slides[0]

    # Increase Shape 1 from 28pt to 32pt (+4 delta)
    res1 = modify_text(slide, shape_id=slide.shapes[0].shape_id, font_size_delta=4.0)
    assert res1["resulting_font_size"] == 32.0

    # Decrease Shape 2 from 18pt to 15pt (-3 delta)
    res2 = modify_text(slide, shape_id=slide.shapes[1].shape_id, font_size_delta=-3.0)
    assert res2["resulting_font_size"] == 15.0


def test_relative_font_size_scale(typography_deck: Path):
    """Test scaling font size using a multiplier with min/max bounds."""
    session = open_presentation(str(typography_deck))
    prs = Presentation(session.working_path)
    slide = prs.slides[0]

    # Scale Shape 2 (18pt) by 1.5 -> 27pt
    res = modify_text(slide, shape_id=slide.shapes[1].shape_id, font_size_scale=1.5)
    assert res["resulting_font_size"] == 27.0

    # Scale with max bound clamping: 18pt * 2.0 = 36pt, clamped to max 24pt
    res_clamped = modify_text(
        slide,
        shape_id=slide.shapes[1].shape_id,
        font_size_scale=2.0,
        max_font_size=24.0,
    )
    assert res_clamped["resulting_font_size"] == 24.0


def test_batch_modify_text_relative(typography_deck: Path):
    """Test batch modifying text using delta, scale, and mixed absolute sizes."""
    session = open_presentation(str(typography_deck))
    slide_shapes = Presentation(session.working_path).slides[0].shapes
    s1_id, s2_id, s3_id = slide_shapes[0].shape_id, slide_shapes[1].shape_id, slide_shapes[2].shape_id

    ops = [
        {"shape_id": s1_id, "font_size_delta": 2.0},  # 28 -> 30
        {"shape_id": s2_id, "font_size_scale": 1.2, "max_pt": 20.0},  # 18 * 1.2 = 21.6 -> clamped 20.0
        {"shape_id": s3_id, "text": "Updated Bullet\nSecond Line", "font_size": 14.0},
    ]

    res = ppt_batch_modify_text(slide_number=1, operations=ops)
    assert res["success"] is True
    assert res["operations_applied"] == 3
    assert res["results"][0]["resulting_font_size"] == 30.0
    assert res["results"][1]["resulting_font_size"] == 20.0
    assert res["results"][2]["font_size"] == 14.0


def test_scale_slide_typography(typography_deck: Path):
    """Test ppt_scale_slide_typography proportionally scales all text while respecting bounds."""
    session = open_presentation(str(typography_deck))
    res = ppt_scale_slide_typography(
        slide_number=1,
        scale_factor=1.25,
        min_pt=8.0,
        max_pt=32.0,
    )

    assert res["success"] is True
    assert res["total_shapes_modified"] >= 3
    # Check that hierarchy was preserved (title grew, subtitle grew)
    title_res = [s for s in res["shapes_modified"] if "Main Title" in s["text_summary"]][0]
    sub_res = [s for s in res["shapes_modified"] if "Strategic" in s["text_summary"]][0]

    assert title_res["new_primary_size"] > sub_res["new_primary_size"]
    # 28 * 1.25 = 35 -> clamped to max_pt 32.0
    assert title_res["new_primary_size"] == 32.0
    # 18 * 1.25 = 22.5
    assert sub_res["new_primary_size"] == 22.5


def test_smart_tiny_font_validation():
    """VAL-04 should distinguish intentional badges from problematic tiny body text."""
    # Intentional badge (7pt)
    badge_shape = ShapeModel(
        shape_id=1,
        name="Status Badge",
        shape_type=ShapeType.TEXT_BOX,
        semantic_role=SemanticRole.UNKNOWN,
        bbox=BoundingBox.from_inches(10.0, 0.5, 1.5, 0.3),
        text_frame=TextFrameModel(
            text="PROD READY",
            paragraphs=[ParagraphModel(text="PROD READY", runs=[TextRunModel(text="PROD READY", style=TextStyle(font_size_pt=7.0))])],
        ),
    )

    # Unintentional tiny body text (7pt) in main slide body
    tiny_body_shape = ShapeModel(
        shape_id=2,
        name="Main Body Box",
        shape_type=ShapeType.TEXT_BOX,
        semantic_role=SemanticRole.BODY,
        bbox=BoundingBox.from_inches(1.0, 2.5, 6.0, 3.0),
        text_frame=TextFrameModel(
            text="Important body text paragraph that should not be 7pt.",
            paragraphs=[ParagraphModel(text="Important body text paragraph", runs=[TextRunModel(text="Important body text paragraph", style=TextStyle(font_size_pt=7.0))])],
        ),
    )

    slide = SlideModel(slide_number=1, slide_id=1, shapes=[badge_shape, tiny_body_shape])
    val_res = validate_slide(slide, rules=["VAL-04"])

    issues = val_res.issues
    assert len(issues) == 2

    # Check classifications
    badge_issue = [i for i in issues if i.shape_ids == [1]][0]
    body_issue = [i for i in issues if i.shape_ids == [2]][0]

    assert badge_issue.details["classification"] == "INTENTIONAL_COMPACT_TEXT"
    assert badge_issue.severity == IssueSeverity.INFO

    assert body_issue.details["classification"] == "SUSPICIOUS_TINY_TEXT"
    assert body_issue.severity == IssueSeverity.WARNING


def test_wide_banner_no_false_overflow():
    """VAL-03 should not trigger false positive overflow for single-line text on a wide banner."""
    wide_shape = ShapeModel(
        shape_id=1,
        name="Header Banner",
        shape_type=ShapeType.TEXT_BOX,
        semantic_role=SemanticRole.TITLE,
        bbox=BoundingBox.from_inches(0.8, 0.4, 11.5, 0.6),  # 11.5 inches wide!
        text_frame=TextFrameModel(
            text="Enterprise Cloud Data Platform Architecture & Deployment Pipeline",
            paragraphs=[
                ParagraphModel(
                    text="Enterprise Cloud Data Platform Architecture & Deployment Pipeline",
                    runs=[TextRunModel(text="Enterprise Cloud Data Platform Architecture & Deployment Pipeline", style=TextStyle(font_size_pt=18.0))],
                )
            ],
        ),
    )

    slide = SlideModel(slide_number=1, slide_id=1, shapes=[wide_shape])
    val_res = validate_slide(slide, rules=["VAL-03"])
    assert len(val_res.issues) == 0
    assert val_res.is_valid is True

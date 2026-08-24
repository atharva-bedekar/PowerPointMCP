"""Unit tests for PowerPoint MCP v1.1 Phase E: Composite Flow Diagram Creation."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches

from powerpoint_mcp.tools.editing import ppt_create_flow_diagram
from powerpoint_mcp.tools.inspection import ppt_validate_slide
from powerpoint_mcp.tools.versioning import open_presentation


@pytest.fixture
def empty_deck(tmp_path: Path) -> Path:
    """Create an empty test deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slides.add_slide(prs.slide_layouts[6])
    deck_path = tmp_path / "flow_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_create_horizontal_flow_diagram(empty_deck: Path):
    """Test generating a 4-step horizontal process flow diagram."""
    session = open_presentation(str(empty_deck))
    steps = [
        {"title": "Data Ingestion", "description": "Stream events via Kafka"},
        {"title": "Transformation", "description": "Spark ETL aggregation"},
        {"title": "Validation", "description": "Automated schema testing"},
        {"title": "Publishing", "description": "Snowflake data mart"},
    ]

    res = ppt_create_flow_diagram(
        slide_number=1,
        steps=steps,
        direction="horizontal",
        shape_type="rounded_rectangle",
        style_preset="card_accent",
        start_x=1.0,
        start_y=2.0,
        total_width=11.333,
    )

    assert res["success"] is True
    assert res["step_count"] == 4
    assert len(res["node_shape_ids"]) == 4
    assert len(res["connector_shape_ids"]) == 3  # 3 connecting arrows between 4 nodes

    # Run slide validation to confirm zero overlaps or boundary violations
    val_res = ppt_validate_slide(slide_number=1)
    assert val_res["success"] is True
    assert val_res["is_valid"] is True
    assert val_res["summary"]["overlaps"] == 0
    assert val_res["summary"]["boundary_violations"] == 0


def test_create_vertical_flow_diagram(empty_deck: Path):
    """Test generating a 3-step vertical process flow diagram."""
    session = open_presentation(str(empty_deck))
    steps = ["Phase 1: Architecture", "Phase 2: Development", "Phase 3: Deployment"]

    res = ppt_create_flow_diagram(
        slide_number=1,
        steps=steps,
        direction="vertical",
        shape_type="rounded_rectangle",
        style_preset="card_default",
        start_x=3.0,
        start_y=1.5,
        node_width=7.0,
        node_height=1.2,
        node_gap=0.5,
    )

    assert res["success"] is True
    assert res["step_count"] == 3
    assert len(res["node_shape_ids"]) == 3
    assert len(res["connector_shape_ids"]) == 2

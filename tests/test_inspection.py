"""Unit tests for PowerPoint MCP Inspection Engine, Models, Styles, and Shape Matching."""

import io
from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    AlignmentType,
    BoundingBox,
    DistributionMode,
    EMU_PER_CM,
    EMU_PER_INCH,
    EMU_PER_POINT,
    ParagraphModel,
    POINTS_PER_INCH,
    PresentationMetadata,
    PresentationModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    SpacingMode,
    TextFrameModel,
    TextRunModel,
    TextStyle,
    apply_delta_inches,
    emu_to_inches,
    emu_to_pt,
    inches_to_emu,
    pt_to_emu,
)
from powerpoint_mcp.pptx import (
    PPTXInspector,
    extract_alignment_name,
    extract_embedded_images,
    extract_fill_style,
    extract_font_style,
    extract_hyperlinks,
    extract_line_style,
    extract_paragraph,
    extract_rgb_hex,
    extract_run,
    extract_shape_properties,
    extract_text_frame,
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    inspect_slide_relationships,
    map_shape_type,
    match_shapes,
)


class TestUnitsAndBoundingBox:
    """Test unit conversions, mathematical precision, and BoundingBox operations."""

    def test_conversion_constants(self):
        assert EMU_PER_INCH == 914400
        assert EMU_PER_POINT == 12700
        assert EMU_PER_CM == 360000
        assert POINTS_PER_INCH == 72

    def test_unit_conversions(self):
        assert inches_to_emu(1.0) == 914400
        assert inches_to_emu(0.5) == 457200
        assert inches_to_emu(0.2) == 182880

        assert emu_to_inches(914400) == 1.0
        assert emu_to_inches(457200) == 0.5
        assert emu_to_inches(182880) == 0.2

        assert pt_to_emu(1.0) == 12700
        assert pt_to_emu(72.0) == 914400
        assert emu_to_pt(12700) == 1.0
        assert emu_to_pt(914400) == 72.0

    def test_apply_delta_inches(self):
        initial_emu = 914400  # 1.0 inch
        delta = 0.2  # 0.2 inches -> 182880 EMU
        result_emu = apply_delta_inches(initial_emu, delta)
        assert result_emu == 1097280
        assert emu_to_inches(result_emu) == 1.2

    def test_bounding_box_creation_and_properties(self):
        bbox = BoundingBox(
            left_emu=914400,
            top_emu=1828800,
            width_emu=2743200,
            height_emu=3657600,
        )
        # EMU properties
        assert bbox.left_emu == 914400
        assert bbox.top_emu == 1828800
        assert bbox.width_emu == 2743200
        assert bbox.height_emu == 3657600
        assert bbox.right_emu == 914400 + 2743200
        assert bbox.bottom_emu == 1828800 + 3657600
        assert bbox.center_x_emu == 914400 + 1371600
        assert bbox.center_y_emu == 1828800 + 1828800

        # Inch properties
        assert bbox.left_inches == 1.0
        assert bbox.top_inches == 2.0
        assert bbox.width_inches == 3.0
        assert bbox.height_inches == 4.0
        assert bbox.right_inches == 4.0
        assert bbox.bottom_inches == 6.0
        assert bbox.center_x_inches == 2.5
        assert bbox.center_y_inches == 4.0
        assert bbox.x == 1.0
        assert bbox.y == 2.0

    def test_bounding_box_from_inches_and_from_emu(self):
        bbox_in = BoundingBox.from_inches(1.5, 2.5, 3.5, 4.5)
        assert bbox_in.left_emu == 1371600
        assert bbox_in.top_emu == 2286000
        assert bbox_in.width_emu == 3200400
        assert bbox_in.height_emu == 4114800
        assert bbox_in.left_inches == 1.5
        assert bbox_in.top_inches == 2.5

        bbox_emu = BoundingBox.from_emu(914400, 1828800, 914400, 914400)
        assert bbox_emu.width_inches == 1.0
        assert bbox_emu.height_inches == 1.0

    def test_bounding_box_to_dict(self):
        bbox = BoundingBox.from_inches(1.0, 2.0, 3.0, 4.0)
        d = bbox.to_dict()
        assert d["left_inches"] == 1.0
        assert d["top_inches"] == 2.0
        assert d["width_inches"] == 3.0
        assert d["height_inches"] == 4.0
        assert d["x"] == 1.0
        assert d["y"] == 2.0
        assert d["right"] == 4.0
        assert d["bottom"] == 6.0
        assert d["left_emu"] == 914400


class TestDataModels:
    """Test serialization and methods of data models."""

    def test_text_style_and_run(self):
        style = TextStyle(
            font_name="Calibri",
            font_size_pt=24.0,
            bold=True,
            italic=False,
            underline=None,
            color_rgb="1F497D",
            alignment="center",
        )
        assert style.font_name == "Calibri"
        assert style.font_size_pt == 24.0
        assert style.bold is True
        d_style = style.to_dict()
        assert d_style["font_name"] == "Calibri"
        assert d_style["color_rgb"] == "1F497D"

        run = TextRunModel(text="Header Title", style=style, hyperlink_target="https://example.com")
        d_run = run.to_dict()
        assert d_run["text"] == "Header Title"
        assert d_run["hyperlink_target"] == "https://example.com"
        assert d_run["style"]["font_size_pt"] == 24.0

    def test_paragraph_and_text_frame(self):
        style = TextStyle(font_name="Arial", font_size_pt=14.0)
        run1 = TextRunModel(text="Hello ", style=style)
        run2 = TextRunModel(text="World", style=TextStyle(font_name="Arial", font_size_pt=14.0, bold=True))
        p = ParagraphModel(text="Hello World", runs=[run1, run2], alignment="left", level=0)
        tf = TextFrameModel(text="Hello World", paragraphs=[p], word_wrap=True)

        d_tf = tf.to_dict()
        assert d_tf["text"] == "Hello World"
        assert d_tf["paragraph_count"] == 1
        assert len(d_tf["paragraphs"][0]["runs"]) == 2
        assert d_tf["word_wrap"] is True

    def test_shape_model(self):
        bbox = BoundingBox.from_inches(1.0, 1.0, 4.0, 2.0)
        style = TextStyle(font_name="Calibri", font_size_pt=20.0, bold=True, color_rgb="000000")
        run = TextRunModel(text="Card Title", style=style)
        p = ParagraphModel(text="Card Title", runs=[run], alignment="center")
        tf = TextFrameModel(text="Card Title", paragraphs=[p])

        shape = ShapeModel(
            shape_id=10,
            name="Card Box",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.BODY,
            bbox=bbox,
            rotation=0.0,
            z_order=2,
            text_frame=tf,
            fill={"type": "solid", "color": "FFFFFF"},
            line={"color": "CCCCCC", "width_pt": 1.5},
        )
        assert shape.id == 10
        assert shape.type == ShapeType.AUTO_SHAPE
        assert shape.role == SemanticRole.BODY
        assert shape.fill_color == "FFFFFF"
        assert shape.line_width_pt == 1.5

        d_shape = shape.to_dict()
        assert d_shape["shape_id"] == 10
        assert d_shape["name"] == "Card Box"
        assert d_shape["role"] == "body"
        assert d_shape["font_family"] == "Calibri"
        assert d_shape["bold"] is True
        assert d_shape["x"] == 1.0
        assert d_shape["width"] == 4.0

    def test_slide_model_methods(self):
        bbox = BoundingBox.from_inches(1.0, 1.0, 2.0, 1.0)
        s1 = ShapeModel(shape_id=1, name="Title 1", shape_type=ShapeType.TEXT_BOX, semantic_role=SemanticRole.TITLE, bbox=bbox)
        s2 = ShapeModel(shape_id=2, name="Body 1", shape_type=ShapeType.AUTO_SHAPE, semantic_role=SemanticRole.BODY, bbox=bbox)
        s3 = ShapeModel(shape_id=3, name="Pic 1", shape_type=ShapeType.PICTURE, semantic_role=SemanticRole.IMAGE, bbox=bbox)

        slide = SlideModel(
            slide_number=1,
            slide_id=256,
            layout_name="Title Slide",
            title="Title 1",
            shapes=[s1, s2, s3],
            notes_text="Speaker notes",
        )
        assert slide.shape_count == 3
        assert slide.has_notes is True
        assert slide.notes == "Speaker notes"
        assert slide.get_shape_by_id(2) == s2
        assert slide.get_shape_by_id(99) is None
        assert len(slide.get_shapes_by_role(SemanticRole.TITLE)) == 1
        assert len(slide.get_shapes_by_type(ShapeType.AUTO_SHAPE)) == 1

        d_slide = slide.to_dict()
        assert d_slide["slide_number"] == 1
        assert d_slide["shape_count"] == 3
        assert len(d_slide["shapes"]) == 3

    def test_presentation_model(self):
        prs_model = PresentationModel(
            path="test.pptx",
            width_inches=13.3333,
            height_inches=7.5,
            width_emu=12192000,
            height_emu=6858000,
            slide_count=1,
            layouts=["Blank", "Title"],
            slides=[SlideModel(slide_number=1, slide_id=256)],
            slide_titles=[{"slide_number": 1, "title": "Intro"}],
            metadata=PresentationMetadata(title="Sample Deck", author="Test Author"),
        )
        assert prs_model.slide_count == 1
        assert prs_model.presentation_path == "test.pptx"
        assert prs_model.get_slide(1) is not None
        assert prs_model.get_slide(2) is None
        d_prs = prs_model.to_dict()
        assert d_prs["dimensions"]["width_inches"] == 13.3333
        assert d_prs["metadata"]["title"] == "Sample Deck"


class TestSemanticRoleInference:
    """Test the 5-stage rule cascade for inferring semantic roles."""

    def test_infer_role_from_title_placeholder(self):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)

        title_shape = slide.shapes.title
        title_shape.text = "Company Presentation"

        role = infer_semantic_role(title_shape, int(prs.slide_width), int(prs.slide_height))
        assert role == SemanticRole.TITLE

    def test_infer_role_from_spatial_and_font_heuristics(self):
        prs = Presentation()
        slide_w = Inches(10)
        slide_h = Inches(5.625)
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 1. Custom Title Box (top < 0.22, font >= 24)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Custom Title"
        p.font.size = Pt(32)
        role_title = infer_semantic_role(title_box, int(slide_w), int(slide_h))
        assert role_title == SemanticRole.TITLE

        # 2. Custom Subtitle Box (0.15 <= top < 0.38, 14 <= font < 24)
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
        tf_sub = sub_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = "Strategic Insights and Roadmap"
        p_sub.font.size = Pt(18)
        role_sub = infer_semantic_role(sub_box, int(slide_w), int(slide_h))
        assert role_sub == SemanticRole.SUBTITLE

        # 3. Custom Footer Box (top >= 0.85 or named Footer)
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "Confidential - Internal Use Only"
        p_foot.font.size = Pt(10)
        role_foot = infer_semantic_role(footer_box, int(slide_w), int(slide_h))
        assert role_foot == SemanticRole.FOOTER

        # 4. Custom Body Box (mid-slide multi-paragraph)
        body_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.0))
        tf_body = body_box.text_frame
        p_body1 = tf_body.paragraphs[0]
        p_body1.text = "Point 1: Revenue grew by 25% YoY."
        p_body2 = tf_body.add_paragraph()
        p_body2.text = "Point 2: Margin expanded by 300bps."
        role_body = infer_semantic_role(body_box, int(slide_w), int(slide_h))
        assert role_body == SemanticRole.BODY

    def test_infer_role_table_and_connector(self):
        prs = Presentation()
        slide_w = Inches(10)
        slide_h = Inches(5.625)
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Table
        table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        role_table = infer_semantic_role(table_shape, int(slide_w), int(slide_h))
        assert role_table == SemanticRole.TABLE

        # Connector / Line
        connector = slide.shapes.add_connector(1, Inches(1), Inches(1), Inches(3), Inches(1))
        role_conn = infer_semantic_role(connector, int(slide_w), int(slide_h))
        assert role_conn == SemanticRole.DIAGRAM


class TestSyntheticDeckInspection:
    """Test deep inspection on the synthetic sample presentation fixture."""

    def test_inspect_presentation_metadata_and_dimensions(self, synthetic_deck_path):
        prs_model = inspect_presentation(synthetic_deck_path)

        assert prs_model.slide_count == 3
        assert prs_model.width_inches == pytest.approx(13.3333, abs=0.01)
        assert prs_model.height_inches == pytest.approx(7.5, abs=0.01)
        assert abs(prs_model.width_emu - 12192000) <= 10
        assert prs_model.height_emu == 6858000
        assert len(prs_model.slides) == 3
        assert len(prs_model.slide_titles) == 3
        assert prs_model.slide_titles[0]["title"] == "Quarterly Performance Overview"
        assert prs_model.slide_titles[1]["title"] == "Operational Architecture"
        assert prs_model.slide_titles[2]["title"] == "Audit & Compliance Issues"

    def test_inspect_slide_1_shapes_and_typography(self, synthetic_deck_path):
        slide1 = inspect_slide(synthetic_deck_path, 1)

        assert slide1.slide_number == 1
        assert slide1.title == "Quarterly Performance Overview"
        assert slide1.shape_count >= 5

        # Check title shape
        title_shapes = slide1.get_shapes_by_role(SemanticRole.TITLE)
        assert len(title_shapes) == 1
        title_s = title_shapes[0]
        assert "Quarterly Performance Overview" in title_s.text_frame.text
        assert title_s.bbox.left_inches == pytest.approx(1.0, abs=0.01)
        assert title_s.bbox.top_inches == pytest.approx(0.8, abs=0.01)

        # Check subtitle shape
        sub_shapes = slide1.get_shapes_by_role(SemanticRole.SUBTITLE)
        assert len(sub_shapes) == 1
        sub_s = sub_shapes[0]
        assert "Executive Summary" in sub_s.text_frame.text

        # Check picture shape
        pic_shapes = slide1.get_shapes_by_type(ShapeType.PICTURE)
        assert len(pic_shapes) >= 1
        pic_s = pic_shapes[0]
        assert pic_s.semantic_role == SemanticRole.IMAGE
        assert pic_s.image_metadata is not None
        assert pic_s.image_metadata["size_bytes"] > 0

    def test_inspect_slide_2_layout_and_footer(self, synthetic_deck_path):
        slide2 = inspect_slide(synthetic_deck_path, 2)

        assert slide2.slide_number == 2
        assert slide2.title == "Operational Architecture"

        # Check footer
        footer_shapes = slide2.get_shapes_by_role(SemanticRole.FOOTER)
        assert len(footer_shapes) >= 1
        footer_s = footer_shapes[0]
        assert "Confidential" in footer_s.text_frame.text
        assert footer_s.bbox.top_inches >= 6.0

    def test_inspect_slide_3_defects_and_notes(self, synthetic_deck_path):
        slide3 = inspect_slide(synthetic_deck_path, 3)

        assert slide3.slide_number == 3
        assert slide3.title == "Audit & Compliance Issues"
        assert len(slide3.shapes) >= 5

    def test_inspect_slide_with_speaker_notes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.notes_slide.notes_text_frame.text = "Audit defect test notes: VAL-01 overlap present."
        slide_model = inspect_slide(prs, 1)

        assert slide_model.has_notes is True
        assert slide_model.notes is not None
        assert "VAL-01" in slide_model.notes

    def test_inspect_slide_out_of_bounds(self, synthetic_deck_path):
        with pytest.raises(ValueError, match="out of bounds"):
            inspect_slide(synthetic_deck_path, 0)

        with pytest.raises(ValueError, match="out of bounds"):
            inspect_slide(synthetic_deck_path, 4)

        with pytest.raises(ValueError, match="out of bounds"):
            inspect_slide(synthetic_deck_path, -1)

    def test_inspect_shape_deep_properties(self, synthetic_deck_path):
        slide1 = inspect_slide(synthetic_deck_path, 1)
        first_shape_id = slide1.shapes[0].shape_id

        deep_shape = inspect_shape(synthetic_deck_path, 1, first_shape_id)
        assert deep_shape.shape_id == first_shape_id
        assert deep_shape.name == slide1.shapes[0].name
        assert deep_shape.bbox.width_inches > 0
        assert deep_shape.bbox.height_inches > 0

    def test_inspect_shape_not_found(self, synthetic_deck_path):
        with pytest.raises(ValueError, match="does not exist"):
            inspect_shape(synthetic_deck_path, 1, 99999)


class TestStylesAndRelationships:
    """Test style extraction helpers and OpenXML relationship parsers."""

    def test_extract_rgb_hex_and_font(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        p = box.text_frame.paragraphs[0]
        p.text = "Color Test"
        p.font.name = "Georgia"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = RGBColor(255, 128, 0)

        style = extract_font_style(p.font)
        assert style.font_name == "Georgia"
        assert style.font_size_pt == 22.0
        assert style.bold is True
        assert style.italic is True
        assert style.color_rgb == "FF8000"

    def test_extract_paragraph_and_text_frame(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = "Centered Heading"
        p.alignment = PP_ALIGN.CENTER

        tf_model = extract_text_frame(tf)
        assert tf_model.text == "Centered Heading"
        assert tf_model.margin_left_inches == pytest.approx(0.2, abs=0.01)
        assert tf_model.paragraphs[0].alignment == "center"

    def test_extract_embedded_images(self, synthetic_deck_path):
        images = extract_embedded_images(Presentation(str(synthetic_deck_path)))
        assert len(images) >= 1
        first_img = images[0]
        assert first_img["sha256"] is not None
        assert len(first_img["sha256"]) == 64
        assert first_img["size_bytes"] > 0
        assert first_img["extension"] in ("png", "jpg", "jpeg")

    def test_inspect_slide_relationships(self, synthetic_deck_path):
        prs = Presentation(str(synthetic_deck_path))
        slide1 = prs.slides[0]
        rels = inspect_slide_relationships(slide1)
        assert len(rels) > 0
        r_ids = [r["r_id"] for r in rels]
        assert any("rId" in r_id for r_id in r_ids)


class TestShapeMatching:
    """Test multi-factor semantic shape matching across slides."""

    def test_match_identical_slide(self, synthetic_deck_path):
        slide1_a = inspect_slide(synthetic_deck_path, 1)
        slide1_b = inspect_slide(synthetic_deck_path, 1)

        matches = match_shapes(slide1_a, slide1_b, min_confidence=0.40)
        assert len(matches) == len(slide1_a.shapes)

        # Self-matches should have very high confidence (~1.0)
        for m in matches:
            assert m["confidence_score"] >= 0.90
            assert m["shape_a_id"] == m["shape_b_id"]
            assert "factors" in m
            assert m["factors"]["role_score"] >= 0.0

    def test_match_shapes_with_spatial_shift(self, synthetic_deck_path):
        slide1_a = inspect_slide(synthetic_deck_path, 1)
        slide1_b = inspect_slide(synthetic_deck_path, 1)

        # Shift one shape slightly in slide1_b
        title_shape_b = slide1_b.get_shapes_by_role(SemanticRole.TITLE)[0]
        # Simulate moving left by 0.2 inches (182880 EMUs)
        title_shape_b.bbox.left_emu += 182880

        matches = match_shapes(slide1_a, slide1_b, min_confidence=0.40)
        title_match = next((m for m in matches if m["shape_a_id"] == title_shape_b.shape_id), None)
        assert title_match is not None
        assert title_match["confidence_score"] >= 0.85
        assert "factors" in title_match
        assert "reasoning" in title_match

    def test_match_shapes_cross_slides(self, synthetic_deck_path):
        slide1 = inspect_slide(synthetic_deck_path, 1)
        slide2 = inspect_slide(synthetic_deck_path, 2)

        matches = match_shapes(slide1, slide2, min_confidence=0.40)
        # Titles on both slides have role=TITLE and similar top alignment
        title_match = next((m for m in matches if m["shape_a_name"] == slide1.get_shapes_by_role(SemanticRole.TITLE)[0].name), None)
        if title_match:
            assert title_match["factors"]["role_score"] == 1.0


class TestPPTXInspectorClass:
    """Test PPTXInspector static methods."""

    def test_inspector_static_methods(self, synthetic_deck_path):
        prs_model = PPTXInspector.inspect_presentation(synthetic_deck_path)
        assert prs_model.slide_count == 3

        slide_model = PPTXInspector.inspect_slide(synthetic_deck_path, 1)
        assert slide_model.slide_number == 1

        shape_model = PPTXInspector.inspect_shape(synthetic_deck_path, 1, slide_model.shapes[0].shape_id)
        assert shape_model.shape_id == slide_model.shapes[0].shape_id

        matches = PPTXInspector.match_shapes(slide_model, slide_model)
        assert len(matches) > 0


class TestEdgeCasesAndBoundaries:
    """Test edge cases, zero/negative coordinates, whitespace shapes, and empty collections."""

    def test_bounding_box_negative_coordinates(self):
        bbox = BoundingBox(left_emu=-457200, top_emu=-182880, width_emu=914400, height_emu=914400)
        assert bbox.left_inches == -0.5
        assert bbox.top_inches == -0.2
        assert bbox.right_inches == 0.5
        assert bbox.bottom_inches == 0.8
        d = bbox.to_dict()
        assert d["x"] == -0.5
        assert d["y"] == -0.2

    def test_infer_role_whitespace_and_empty_shape(self):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Blank textbox
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        role_empty = infer_semantic_role(box, int(prs.slide_width), int(prs.slide_height))
        assert role_empty == SemanticRole.UNKNOWN

        # Whitespace-only textbox
        box.text_frame.text = "   \n\t   "
        role_ws = infer_semantic_role(box, int(prs.slide_width), int(prs.slide_height))
        assert role_ws == SemanticRole.UNKNOWN

    def test_match_shapes_empty_slides(self):
        slide_a = SlideModel(slide_number=1, slide_id=100, shapes=[])
        slide_b = SlideModel(slide_number=2, slide_id=101, shapes=[])

        matches = match_shapes(slide_a, slide_b)
        assert matches == []

    def test_match_shapes_threshold_filtering(self):
        bbox1 = BoundingBox.from_inches(1.0, 1.0, 2.0, 1.0)
        bbox2 = BoundingBox.from_inches(10.0, 6.0, 1.0, 0.5)

        s1 = ShapeModel(shape_id=1, name="Card Alpha", shape_type=ShapeType.TEXT_BOX, semantic_role=SemanticRole.TITLE, bbox=bbox1)
        s2 = ShapeModel(shape_id=2, name="Zebra Beta", shape_type=ShapeType.PICTURE, semantic_role=SemanticRole.IMAGE, bbox=bbox2)

        slide_a = SlideModel(slide_number=1, slide_id=1, shapes=[s1])
        slide_b = SlideModel(slide_number=2, slide_id=2, shapes=[s2])

        # With high threshold (0.8), dissimilar shapes should not match
        matches_high = match_shapes(slide_a, slide_b, min_confidence=0.80)
        assert len(matches_high) == 0

    def test_presentation_path_types(self, synthetic_deck_path):
        # Test Path object
        prs_path_obj = inspect_presentation(Path(synthetic_deck_path))
        assert prs_path_obj.slide_count == 3

        # Test in-memory Presentation object
        prs_in_memory = inspect_presentation(Presentation(str(synthetic_deck_path)))
        assert prs_in_memory.slide_count == 3


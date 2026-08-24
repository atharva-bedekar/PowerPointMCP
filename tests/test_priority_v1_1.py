"""Tests for Priority v1.1 Agent-Ergonomics improvements.

Covers:
1. Filtered and concise slide inspection (Priority 1)
2. Dedicated text inspection tool `ppt_inspect_text` (Priority 2)
3. Atomic batch text modification `ppt_batch_modify_text` (Priority 3)
4. Bullet and paragraph formatting preservation during text replacement (Priority 4)
5. Container-aware overlap validation in VAL-01 (Priority 5)
6. Atomic batch shape geometry modification `ppt_batch_modify_shapes` (Priority 6)
"""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.server import app
from powerpoint_mcp.tools.editing import (
    ppt_batch_modify_shapes,
    ppt_batch_modify_text,
    ppt_modify_text,
)
from powerpoint_mcp.tools.inspection import (
    ppt_inspect_slide,
    ppt_inspect_text,
    ppt_validate_slide,
)
from powerpoint_mcp.tools.versioning import open_presentation, save_session


def _extract_tool_result(res):
    import json
    if hasattr(res, "content") and res.content:
        first = res.content[0]
        if hasattr(first, "text"):
            try:
                return json.loads(first.text)
            except Exception:
                return {"text": first.text}
    if isinstance(res, dict):
        return res
    return {}


class TestPriority1FilteredInspection:
    """Tests for Priority 1: concise and filtered slide inspection."""

    def test_text_only_filter(self, synthetic_deck_path: Path, clean_ppt_env):
        open_presentation(synthetic_deck_path)
        all_res = ppt_inspect_slide(slide_number=2, text_only=False)
        text_res = ppt_inspect_slide(slide_number=2, text_only=True)

        assert all_res["success"] is True
        assert text_res["success"] is True
        assert text_res["shape_count"] <= all_res["shape_count"]
        # Every shape returned in text_only must contain text
        for s in text_res["shapes"]:
            assert "text" in s and s["text"].strip()

    def test_inclusion_flags(self, synthetic_deck_path: Path, clean_ppt_env):
        open_presentation(synthetic_deck_path)
        res_no_geom = ppt_inspect_slide(slide_number=1, include_geometry=False)
        res_no_style = ppt_inspect_slide(slide_number=1, include_style=False)

        assert res_no_geom["success"] is True
        first_g = res_no_geom["shapes"][0]
        assert "x" not in first_g and "bbox" not in first_g

        assert res_no_style["success"] is True
        first_s = res_no_style["shapes"][0]
        assert "font_family" not in first_s and "color" not in first_s

    def test_role_and_type_filters(self, synthetic_deck_path: Path, clean_ppt_env):
        open_presentation(synthetic_deck_path)
        res_title = ppt_inspect_slide(slide_number=1, semantic_roles=["title"])
        assert res_title["success"] is True
        assert len(res_title["shapes"]) == 1
        assert res_title["shapes"][0]["semantic_role"] == "title"


class TestPriority2InspectText:
    """Tests for Priority 2: dedicated ppt_inspect_text tool."""

    def test_inspect_text_returns_clean_text_records(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        open_presentation(synthetic_deck_path)
        res = ppt_inspect_text(slide_number=1, include_paragraph_metadata=True)

        assert res["success"] is True
        assert res["text_shape_count"] >= 1
        shapes = res["shapes"]
        title_item = next(s for s in shapes if s["semantic_role"] == "title")
        assert "shape_id" in title_item
        assert "text" in title_item
        assert "font_family" in title_item
        assert "font_size" in title_item
        assert "paragraphs" in title_item
        assert title_item["paragraph_count"] >= 1


class TestPriority3BatchModifyText:
    """Tests for Priority 3: ppt_batch_modify_text tool."""

    def test_batch_modify_multiple_text_shapes(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        open_presentation(synthetic_deck_path)
        # Inspect text on slide 2 to get IDs
        txt_res = ppt_inspect_text(slide_number=2)
        shapes = txt_res["shapes"]
        assert len(shapes) >= 2

        s1_id = shapes[0]["shape_id"]
        s2_id = shapes[1]["shape_id"]

        batch_res = ppt_batch_modify_text(
            slide_number=2,
            operations=[
                {"shape_id": s1_id, "text": "BATCH UPDATED 1", "font_size": 20},
                {"shape_id": s2_id, "text": "BATCH UPDATED 2", "font_size": 14},
            ],
        )

        assert batch_res["success"] is True
        assert batch_res["operations_applied"] == 2

        # Verify on re-inspection
        re_res = ppt_inspect_text(slide_number=2)
        id_to_text = {s["shape_id"]: s["text"] for s in re_res["shapes"]}
        assert "BATCH UPDATED 1" in id_to_text[s1_id]
        assert "BATCH UPDATED 2" in id_to_text[s2_id]

    def test_batch_modify_text_pre_validation_safety(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        open_presentation(synthetic_deck_path)
        # Attempt with a valid shape and an invalid shape ID 99999
        txt_res = ppt_inspect_text(slide_number=1)
        valid_id = txt_res["shapes"][0]["shape_id"]

        orig_text = txt_res["shapes"][0]["text"]

        # Should fail cleanly
        fail_res = ppt_batch_modify_text(
            slide_number=1,
            operations=[
                {"shape_id": valid_id, "text": "SHOULD NOT BE WRITTEN"},
                {"shape_id": 99999, "text": "INVALID SHAPE"},
            ],
        )
        assert fail_res["success"] is False
        assert "Shape with ID 99999 not found" in fail_res["message"]

        # Verify no partial corruption occurred
        re_check = ppt_inspect_text(slide_number=1)
        assert re_check["shapes"][0]["text"] == orig_text


class TestPriority4ParagraphAndBulletPreservation:
    """Tests for Priority 4: paragraph format and bullet preservation."""

    def test_bullet_character_and_indent_preservation(
        self, tmp_path: Path, clean_ppt_env
    ):
        # Create a presentation with bulleted paragraphs
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = "First bullet item"
        p1.level = 0
        p1.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = "Second bullet sub-item"
        p2.level = 1
        p2.space_before = Pt(4)

        deck_path = tmp_path / "bullet_test.pptx"
        prs.save(str(deck_path))

        # Open in MCP session
        open_presentation(deck_path)

        # Modify text with new lines
        mod_res = ppt_modify_text(
            slide_number=1,
            shape_id=txBox.shape_id,
            text="Updated first bullet\nUpdated second sub-bullet",
        )
        assert mod_res["success"] is True

        # Save and verify
        save_session()
        saved_prs = Presentation(str(deck_path))
        saved_tf = saved_prs.slides[0].shapes[0].text_frame

        assert len(saved_tf.paragraphs) == 2
        assert saved_tf.paragraphs[0].text == "Updated first bullet"
        assert saved_tf.paragraphs[0].level == 0
        assert saved_tf.paragraphs[0].space_before.pt == 6.0

        assert saved_tf.paragraphs[1].text == "Updated second sub-bullet"
        assert saved_tf.paragraphs[1].level == 1
        assert saved_tf.paragraphs[1].space_before.pt == 4.0


class TestPriority5ContainerAwareValidation:
    """Tests for Priority 5: container-aware overlap detection."""

    def test_card_with_nested_text_is_not_flagged_as_overlap(
        self, tmp_path: Path, clean_ppt_env
    ):
        # Create a deck with a filled background card and text box inside it
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        from pptx.enum.shapes import MSO_SHAPE
        # 1. Background Card shape (3.0 x 3.0 at x=1.0, y=1.0)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(3)
        )
        card.name = "Container Card"

        # 2. Text Box inside Card (2.0 x 1.0 at x=1.5, y=1.5)
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(2), Inches(1))
        tb.name = "Card Title Text"
        tb.text_frame.text = "Nested Content"

        deck_path = tmp_path / "container_test.pptx"
        prs.save(str(deck_path))

        open_presentation(deck_path)
        val_res = ppt_validate_slide(slide_number=1)

        assert val_res["success"] is True
        assert val_res["summary"]["overlaps"] == 0
        assert val_res["valid"] is True

    def test_two_overlapping_cards_are_flagged_as_actual_overlap(
        self, tmp_path: Path, clean_ppt_env
    ):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        from pptx.enum.shapes import MSO_SHAPE
        # Card A at (1, 1) size (3, 3)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(3))
        # Card B at (2, 2) size (3, 3) -> collides with Card A by 4 sq in
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(3), Inches(3))

        deck_path = tmp_path / "collision_test.pptx"
        prs.save(str(deck_path))

        open_presentation(deck_path)
        val_res = ppt_validate_slide(slide_number=1)

        assert val_res["success"] is True
        assert val_res["summary"]["overlaps"] >= 1
        assert val_res["valid"] is False


class TestPriority6BatchModifyShapes:
    """Tests for Priority 6: ppt_batch_modify_shapes tool."""

    def test_batch_modify_shapes_positions_and_sizes(
        self, synthetic_deck_path: Path, clean_ppt_env
    ):
        open_presentation(synthetic_deck_path)
        insp = ppt_inspect_slide(slide_number=1)
        shapes = insp["shapes"]
        assert len(shapes) >= 2

        s1_id = shapes[0]["shape_id"]
        s2_id = shapes[1]["shape_id"]

        batch_res = ppt_batch_modify_shapes(
            slide_number=1,
            operations=[
                {"shape_id": s1_id, "changes": {"x": 1.5, "width": 4.0}},
                {"shape_id": s2_id, "dx": 0.5, "dy": 0.25},
            ],
        )

        assert batch_res["success"] is True
        assert batch_res["operations_applied"] == 2

        # Verify updated positions
        re_insp = ppt_inspect_slide(slide_number=1)
        shape_dict = {s["shape_id"]: s for s in re_insp["shapes"]}
        assert shape_dict[s1_id]["x"] == 1.5
        assert shape_dict[s1_id]["width"] == 4.0


@pytest.mark.asyncio
class TestMCPEndToEndPriorityV11:
    """MCP JSON-RPC async tool calls test for all newly added tools."""

    async def test_mcp_inspect_text_and_batch_tools(
        self, temp_deck_path: Path, clean_ppt_env
    ):
        # 1. Open
        await app.call_tool("ppt_open", {"presentation_path": str(temp_deck_path)})

        # 2. ppt_inspect_text
        txt_res = await app.call_tool("ppt_inspect_text", {"slide_number": 1})
        txt_data = _extract_tool_result(txt_res)
        assert txt_data["success"] is True
        assert txt_data["text_shape_count"] >= 1

        sid = txt_data["shapes"][0]["shape_id"]

        # 3. ppt_batch_modify_text
        mod_res = await app.call_tool(
            "ppt_batch_modify_text",
            {
                "slide_number": 1,
                "operations": [{"shape_id": sid, "text": "MCP BATCH TITLE", "font_size": 26}],
            },
        )
        mod_data = _extract_tool_result(mod_res)
        assert mod_data["success"] is True
        assert mod_data["operations_applied"] == 1

        # 4. ppt_batch_modify_shapes
        geom_res = await app.call_tool(
            "ppt_batch_modify_shapes",
            {
                "slide_number": 1,
                "operations": [{"shape_id": sid, "changes": {"x": 1.2, "y": 0.8}}],
            },
        )
        geom_data = _extract_tool_result(geom_res)
        assert geom_data["success"] is True

        # 5. ppt_save
        save_res = await app.call_tool("ppt_save", {})
        save_data = _extract_tool_result(save_res)
        assert save_data["success"] is True

        # 6. Verify in file
        final_prs = Presentation(str(temp_deck_path))
        first_shape = final_prs.slides[0].shapes[0]
        assert "MCP BATCH TITLE" in first_shape.text_frame.text
        assert abs(first_shape.left - Inches(1.2)) < 5000

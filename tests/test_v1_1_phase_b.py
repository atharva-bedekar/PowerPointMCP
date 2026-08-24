"""Unit tests for PowerPoint MCP v1.1 Phase B: Structure & Container Analysis."""

from pathlib import Path
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from powerpoint_mcp.models.shape import (
    BoundingBox,
    ParagraphModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    TextFrameModel,
    TextRunModel,
    TextStyle,
)
from powerpoint_mcp.models.slide import SlideModel
from powerpoint_mcp.pptx.structure import analyze_containers, analyze_slide_structure
from powerpoint_mcp.tools.inspection import (
    ppt_analyze_containers,
    ppt_analyze_slide_structure,
    ppt_validate_slide,
)
from powerpoint_mcp.tools.versioning import open_presentation
from powerpoint_mcp.utils.validation import validate_slide


@pytest.fixture
def structure_deck(tmp_path: Path) -> Path:
    """Create a test deck with 3 cards containing titles, bullet points, metrics, and badges."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Slide Title
    tb_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.0), Inches(0.8))
    tb_title.name = "Title 1"
    p_title = tb_title.text_frame.paragraphs[0]
    p_title.text = "Infrastructure Platform Modernization"
    p_title.runs[0].font.size = Pt(28)

    # Card 1 Background (Rounded Rectangle with Fill)
    c1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(1.8),
        Inches(3.5),
        Inches(4.8),
    )
    c1.name = "Card 1 Background"
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(240, 244, 248)

    # Card 1 Title
    c1_title = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(3.1), Inches(0.5))
    c1_title.name = "Card 1 Title"
    p1 = c1_title.text_frame.paragraphs[0]
    p1.text = "Core Engine"
    p1.runs[0].font.size = Pt(16)
    p1.runs[0].font.bold = True

    # Card 1 Metric
    c1_metric = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(3.1), Inches(0.6))
    c1_metric.name = "Metric 1"
    p1_m = c1_metric.text_frame.paragraphs[0]
    p1_m.text = "99.99%"
    p1_m.runs[0].font.size = Pt(24)

    # Card 1 Bullets
    c1_body = slide.shapes.add_textbox(Inches(1.2), Inches(3.3), Inches(3.1), Inches(2.5))
    c1_body.name = "Card 1 Body"
    p1_b1 = c1_body.text_frame.paragraphs[0]
    p1_b1.text = "Distributed cluster management"
    p1_b2 = c1_body.text_frame.add_paragraph()
    p1_b2.text = "Automated horizontal failover"

    # Card 2 Background
    c2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.9),
        Inches(1.8),
        Inches(3.5),
        Inches(4.8),
    )
    c2.name = "Card 2 Background"
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(240, 244, 248)

    # Card 2 Badge
    c2_badge = slide.shapes.add_textbox(Inches(5.1), Inches(2.0), Inches(1.5), Inches(0.35))
    c2_badge.name = "Card 2 Badge"
    p2_badge = c2_badge.text_frame.paragraphs[0]
    p2_badge.text = "BETA"
    p2_badge.runs[0].font.size = Pt(8)

    # Card 2 Title
    c2_title = slide.shapes.add_textbox(Inches(5.1), Inches(2.4), Inches(3.1), Inches(0.5))
    c2_title.name = "Card 2 Title"
    p2_t = c2_title.text_frame.paragraphs[0]
    p2_t.text = "Analytics Pipeline"
    p2_t.runs[0].font.size = Pt(16)
    p2_t.runs[0].font.bold = True

    deck_path = tmp_path / "structure_deck.pptx"
    prs.save(str(deck_path))
    return deck_path


def test_analyze_slide_structure(structure_deck: Path):
    """Test analyzing semantic roles, confidence scores, and parents."""
    session = open_presentation(str(structure_deck))
    res = ppt_analyze_slide_structure(slide_number=1)

    assert res["success"] is True
    assert res["total_elements"] >= 7
    assert res["total_containers"] == 2

    # Check slide title detection
    slide_title_elem = next(e for e in res["elements"] if "Infrastructure" in e["text_preview"])
    assert slide_title_elem["role"] == "slide_title"
    assert slide_title_elem["confidence"] >= 0.90
    assert slide_title_elem["parent_id"] is None

    # Check metric detection
    metric_elem = next(e for e in res["elements"] if "99.99%" in e["text_preview"])
    assert metric_elem["role"] == "metric"
    assert metric_elem["confidence"] >= 0.90
    assert metric_elem["parent_id"] is not None  # Contained inside Card 1

    # Check badge detection
    badge_elem = next(e for e in res["elements"] if "BETA" in e["text_preview"])
    assert badge_elem["role"] == "badge"
    assert badge_elem["parent_id"] is not None  # Contained inside Card 2


def test_analyze_containers(structure_deck: Path):
    """Test analyzing card containers and their nested children."""
    session = open_presentation(str(structure_deck))
    res = ppt_analyze_containers(slide_number=1)

    assert res["success"] is True
    assert res["total_containers"] == 2
    c1 = res["containers"][0]
    assert c1["role"] == "card"
    assert len(c1["children_ids"]) >= 3  # title, metric, body


def test_container_aware_val_01_suppression(structure_deck: Path):
    """Test that text boxes nested inside cards do not produce VAL-01 false positive overlap warnings."""
    session = open_presentation(str(structure_deck))
    val_res = ppt_validate_slide(slide_number=1, rules=["VAL-01"])

    assert val_res["success"] is True
    # Clean card layout should have 0 overlap warnings (nested children are recognized as valid containment)
    assert val_res["summary"]["overlaps"] == 0
    assert val_res["is_valid"] is True


def test_val_01_detects_actual_colliding_cards():
    """Test that two colliding independent cards STILL trigger VAL-01 ACTUAL_OVERLAP."""
    card1 = ShapeModel(
        shape_id=10,
        name="Card 1",
        shape_type=ShapeType.AUTO_SHAPE,
        semantic_role=SemanticRole.CARD,
        bbox=BoundingBox.from_inches(1.0, 1.0, 4.0, 3.0),
    )
    card2 = ShapeModel(
        shape_id=20,
        name="Card 2 Colliding",
        shape_type=ShapeType.AUTO_SHAPE,
        semantic_role=SemanticRole.CARD,
        bbox=BoundingBox.from_inches(3.5, 1.0, 4.0, 3.0),  # Overlaps Card 1 by 0.5 in
    )

    slide = SlideModel(slide_number=1, slide_id=1, shapes=[card1, card2])
    val_res = validate_slide(slide, rules=["VAL-01"])

    assert len(val_res.issues) == 1
    assert val_res.issues[0].details["classification"] == "ACTUAL_OVERLAP"
    assert val_res.issues[0].shape_ids == [10, 20]

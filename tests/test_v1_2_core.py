"""Comprehensive test suite for PowerPoint MCP v1.2 Core Editing & Reliability improvements.

Covers:
- Shape PATCH semantics & explicit zero rejection
- Image insertion & replacement with aspect ratio preservation
- Table cell-level inspection, batch modification, geometry, styling, and merging
- Cross-slide/multi-table batch operations
- Table-aware slide validation and multi-slide validation
"""

import io
from pathlib import Path
import pytest
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.pptx.editor import modify_shape
from powerpoint_mcp.pptx.pictures import add_picture, calculate_picture_dimensions, replace_picture
from powerpoint_mcp.pptx.tables import (
    batch_modify_table_cells,
    inspect_table_cells,
    merge_table_cells,
    set_table_geometry,
    style_table,
)
from powerpoint_mcp.tools.editing import (
    ppt_add_picture,
    ppt_batch_modify_shapes,
    ppt_batch_modify_table_cells,
    ppt_batch_modify_tables,
    ppt_merge_table_cells,
    ppt_modify_shape,
    ppt_replace_picture,
    ppt_set_table_geometry,
    ppt_style_table,
)
from powerpoint_mcp.tools.inspection import ppt_inspect_table, ppt_validate_slide, ppt_validate_slides
from powerpoint_mcp.tools.versioning import open_presentation, save_session


@pytest.fixture
def blank_prs(tmp_path):
    """Create a blank presentation fixture."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(1.5), Inches(2.0), Inches(4.0), Inches(3.0))
    p = tmp_path / "test_deck.pptx"
    prs.save(str(p))
    return p


@pytest.fixture
def sample_png(tmp_path):
    """Generate a sample 200x100 PNG image (2:1 aspect ratio)."""
    img = Image.new("RGB", (200, 100), color=(30, 100, 200))
    p = tmp_path / "sample.png"
    img.save(str(p), format="PNG")
    return p


@pytest.fixture
def sample_jpg(tmp_path):
    """Generate a sample 300x300 JPEG image (1:1 aspect ratio)."""
    img = Image.new("RGB", (300, 300), color=(200, 50, 50))
    p = tmp_path / "sample.jpg"
    img.save(str(p), format="JPEG")
    return p


# =============================================================================
# 1. P0 — Shape Modification PATCH Semantics Tests
# =============================================================================

def test_modify_shape_height_only_preserves_others(blank_prs):
    """Proves modifying ONLY height leaves left, top, and width unchanged."""
    prs = Presentation(str(blank_prs))
    slide = prs.slides[0]
    shape = slide.shapes[0]
    sid = shape.shape_id

    orig_left = shape.left
    orig_top = shape.top
    orig_width = shape.width
    orig_height = shape.height

    res = modify_shape(slide, sid, height=1.2)

    assert round(res["height"], 2) == 1.2
    assert shape.left == orig_left, "Left must remain untouched"
    assert shape.top == orig_top, "Top must remain untouched"
    assert shape.width == orig_width, "Width must remain untouched"
    assert shape.height != orig_height


def test_modify_shape_width_only_preserves_others(blank_prs):
    """Proves modifying ONLY width leaves left, top, and height unchanged."""
    prs = Presentation(str(blank_prs))
    slide = prs.slides[0]
    shape = slide.shapes[0]
    sid = shape.shape_id

    orig_left = shape.left
    orig_top = shape.top
    orig_height = shape.height

    res = modify_shape(slide, sid, width=5.5)

    assert round(res["width"], 2) == 5.5
    assert shape.left == orig_left
    assert shape.top == orig_top
    assert shape.height == orig_height


def test_modify_shape_left_zero_valid_and_preserves_others(blank_prs):
    """Proves x=0.0 is a valid coordinate that does not reset width/height."""
    prs = Presentation(str(blank_prs))
    slide = prs.slides[0]
    shape = slide.shapes[0]
    sid = shape.shape_id

    orig_width = shape.width
    orig_height = shape.height

    res = modify_shape(slide, sid, x=0.0)

    assert res["x"] == 0.0
    assert shape.width == orig_width
    assert shape.height == orig_height


def test_modify_shape_rejects_non_positive_dimensions(blank_prs):
    """Proves explicit width=0 or negative values are rejected with ValueError."""
    prs = Presentation(str(blank_prs))
    slide = prs.slides[0]
    shape = slide.shapes[0]
    sid = shape.shape_id

    with pytest.raises(ValueError, match="positive"):
        modify_shape(slide, sid, width=0.0)

    with pytest.raises(ValueError, match="positive"):
        modify_shape(slide, sid, height=-1.0)


def test_batch_modify_shapes_patch_semantics(blank_prs):
    """Proves batch shape modification adheres strictly to PATCH semantics."""
    prs = Presentation(str(blank_prs))
    slide = prs.slides[0]
    s1 = slide.shapes[0]
    s2 = slide.shapes.add_shape(1, Inches(6.0), Inches(2.0), Inches(3.0), Inches(2.0))
    prs.save(str(blank_prs))

    res = ppt_batch_modify_shapes(
        slide_number=1,
        operations=[
            {"shape_id": s1.shape_id, "height": 0.8},
            {"shape_id": s2.shape_id, "width": 4.5},
        ],
        presentation_path=str(blank_prs),
    )

    assert res["success"] is True
    prs_updated = Presentation(str(blank_prs))
    up_s1 = prs_updated.slides[0].shapes[0]
    up_s2 = prs_updated.slides[0].shapes[1]

    assert round(up_s1.height / 914400, 2) == 0.8
    assert round(up_s1.width / 914400, 2) == 4.0  # untouched
    assert round(up_s2.width / 914400, 2) == 4.5
    assert round(up_s2.height / 914400, 2) == 2.0  # untouched


# =============================================================================
# 2. P0 — Images / Media Tests
# =============================================================================

def test_calculate_picture_dimensions():
    """Verify aspect ratio calculations for various input combinations."""
    # 2:1 aspect ratio (200x100)
    w, h = calculate_picture_dimensions(200, 100, requested_width=4.0)
    assert w == 4.0
    assert h == 2.0

    w, h = calculate_picture_dimensions(200, 100, requested_height=3.0)
    assert w == 6.0
    assert h == 3.0

    # Fit within bounding box
    w, h = calculate_picture_dimensions(200, 100, requested_width=4.0, requested_height=1.0, preserve_aspect_ratio=True)
    assert w == 2.0
    assert h == 1.0


def test_add_picture_png(blank_prs, sample_png):
    """Test adding PNG with aspect ratio calculation."""
    res = ppt_add_picture(
        slide_number=1,
        image_path=str(sample_png),
        left=2.0,
        top=1.5,
        width=4.0,
        presentation_path=str(blank_prs),
    )

    assert res["success"] is True
    assert res["geometry"]["x"] == 2.0
    assert res["geometry"]["y"] == 1.5
    assert res["geometry"]["width"] == 4.0
    assert res["geometry"]["height"] == 2.0  # 2:1 aspect ratio


def test_replace_picture_preserves_geometry(blank_prs, sample_png, sample_jpg):
    """Test replacing picture preserves existing bounds."""
    # First insert sample_png
    add_res = ppt_add_picture(
        slide_number=1,
        image_path=str(sample_png),
        left=1.0,
        top=1.0,
        width=4.0,
        height=2.0,
        presentation_path=str(blank_prs),
    )
    pic_id = add_res["shape_id"]

    # Now replace with sample_jpg (1:1 aspect ratio) preserving geometry
    rep_res = ppt_replace_picture(
        slide_number=1,
        shape_id=pic_id,
        image_path=str(sample_jpg),
        preserve_geometry=True,
        presentation_path=str(blank_prs),
    )

    assert rep_res["success"] is True
    assert rep_res["geometry"]["width"] == 4.0
    assert rep_res["geometry"]["height"] == 2.0
    assert rep_res["image_metadata"]["format"] == "JPEG"


def test_replace_placeholder_with_picture(blank_prs, sample_png):
    """Test replacing a rectangle placeholder with a picture maintains bounds and replaces shape."""
    prs = Presentation(str(blank_prs))
    rect = prs.slides[0].shapes[0]
    rect_id = rect.shape_id

    rep_res = ppt_replace_picture(
        slide_number=1,
        shape_id=rect_id,
        image_path=str(sample_png),
        preserve_geometry=True,
        presentation_path=str(blank_prs),
    )

    assert rep_res["success"] is True
    assert rep_res["operation"] == "replace_placeholder_with_picture"
    assert rep_res["geometry"]["width"] == 4.0
    assert rep_res["geometry"]["height"] == 3.0


# =============================================================================
# 3. P0/P1 — Tables Tests
# =============================================================================

@pytest.fixture
def table_prs(tmp_path):
    """Create a presentation with a 3x3 table fixture."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl_shape = slide.shapes.add_table(3, 3, Inches(1.0), Inches(1.0), Inches(6.0), Inches(2.0))
    table = tbl_shape.table
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(0, 2).text = "Header 3"
    table.cell(1, 0).text = "Row 1"
    table.cell(2, 0).text = "Row 2"
    p = tmp_path / "table_deck.pptx"
    prs.save(str(p))
    return p, tbl_shape.shape_id


def test_inspect_table_compact_and_full(table_prs):
    """Verify table cell inspection in compact and full detail modes."""
    deck_path, tbl_id = table_prs

    # Compact inspection
    compact = ppt_inspect_table(slide_number=1, table_shape_id=tbl_id, detail="compact", presentation_path=str(deck_path))
    assert compact["success"] is True
    assert compact["rows"] == 3
    assert compact["columns"] == 3
    assert "Header 1" in compact["summary"]

    # Full inspection
    full = ppt_inspect_table(slide_number=1, table_shape_id=tbl_id, detail="full", presentation_path=str(deck_path))
    assert full["success"] is True
    assert "cells" in full
    assert len(full["cells"]) == 9
    assert full["cells"][0]["text"] == "Header 1"


def test_batch_modify_table_cells(table_prs):
    """Verify atomic multi-cell text and formatting mutation."""
    deck_path, tbl_id = table_prs

    mutations = [
        {"row": 0, "column": 0, "text": "App Name", "bold": True, "fill": "#003366", "font_color": "#FFFFFF"},
        {"row": 1, "column": 1, "text": "Active", "font_size": 11.0, "align": "center"},
        {"row": 2, "column": 2, "text": "Completed", "vertical_align": "middle"},
    ]

    res = ppt_batch_modify_table_cells(
        slide_number=1,
        table_shape_id=tbl_id,
        mutations=mutations,
        presentation_path=str(deck_path),
    )

    assert res["success"] is True
    assert res["mutations_applied"] == 3

    # Inspect to confirm updates
    insp = ppt_inspect_table(slide_number=1, table_shape_id=tbl_id, detail="full", presentation_path=str(deck_path))
    c00 = next(c for c in insp["cells"] if c["row"] == 0 and c["column"] == 0)
    assert c00["text"] == "App Name"
    assert c00["bold"] is True
    assert c00["fill"] == "#003366"


def test_set_table_geometry_patch_semantics(table_prs):
    """Verify modifying table dimensions and column widths with PATCH semantics."""
    deck_path, tbl_id = table_prs

    res = ppt_set_table_geometry(
        slide_number=1,
        table_shape_id=tbl_id,
        column_widths=[2.5, 2.0, 1.5],
        presentation_path=str(deck_path),
    )

    assert res["success"] is True
    assert res["column_widths"] == [2.5, 2.0, 1.5]


def test_style_table_range_and_borders(table_prs):
    """Verify applying fill, fonts, and borders across a row range."""
    deck_path, tbl_id = table_prs

    res = ppt_style_table(
        slide_number=1,
        table_shape_id=tbl_id,
        range="row:0",
        style={
            "fill": "#0B192C",
            "font_color": "#FFFFFF",
            "bold": True,
            "horizontal_alignment": "center",
            "borders": {"color": "#CCCCCC", "width": 1.0, "sides": ["bottom"]},
        },
        presentation_path=str(deck_path),
    )

    assert res["success"] is True
    assert res["styled_cells_count"] == 3


def test_merge_table_cells(table_prs):
    """Verify cell merge across rectangular region."""
    deck_path, tbl_id = table_prs

    res = ppt_merge_table_cells(
        slide_number=1,
        table_shape_id=tbl_id,
        start_row=1,
        start_column=1,
        end_row=1,
        end_column=2,
        presentation_path=str(deck_path),
    )

    assert res["success"] is True
    assert res["merged_range"] == {"start_row": 1, "start_column": 1, "end_row": 1, "end_column": 2}


def test_batch_modify_tables_multi_operation(table_prs):
    """Verify multi-table batch operations in a single transaction."""
    deck_path, tbl_id = table_prs

    ops = [
        {
            "slide": 1,
            "table": tbl_id,
            "cells": [{"row": 0, "column": 0, "text": "System"}],
            "geometry": {"column_widths": [2.2, 2.0, 1.8]},
            "styles": [{"range": "row:0", "style": {"bold": True}}],
        }
    ]

    res = ppt_batch_modify_tables(operations=ops, presentation_path=str(deck_path))
    assert res["success"] is True
    assert res["total_operations_applied"] == 1


# =============================================================================
# 4. P1 — Table Validation & Multi-Slide Validation Tests
# =============================================================================

def test_validate_slide_table_boundary_overflow(tmp_path):
    """Verify TABLE-01 boundary overflow detection and ensure generic VAL-02 clipping is not duplicated for tables."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Place table extending beyond slide bottom (top=6.5, height=2.0 -> bottom=8.5 > 7.5)
    tbl_shape = slide.shapes.add_table(2, 2, Inches(1.0), Inches(6.5), Inches(5.0), Inches(2.0))
    p = tmp_path / "overflow_deck.pptx"
    prs.save(str(p))

    val_res = ppt_validate_slide(slide_number=1, presentation_path=str(p))
    issues = val_res.get("issues", [])
    table_issues = [i for i in issues if i["rule_id"] == "TABLE-01"]
    val02_issues = [i for i in issues if i["rule_id"] == "VAL-02" and tbl_shape.shape_id in i.get("shape_ids", [])]

    # TABLE-01 must be the canonical warning
    assert len(table_issues) >= 1
    assert "extends" in table_issues[0]["message"]
    assert "bottom" in table_issues[0]["message"]
    # VAL-02 must be suppressed for the table shape to prevent duplicate reporting
    assert len(val02_issues) == 0, "VAL-02 should not duplicate TABLE-01 for table boundaries"


def test_validate_slides_multi_slide(tmp_path):
    """Verify ppt_validate_slides across multiple slides."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    s1.shapes.add_shape(1, Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.0))

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_shape(1, Inches(2.0), Inches(2.0), Inches(3.0), Inches(1.0))

    p = tmp_path / "multi_deck.pptx"
    prs.save(str(p))

    res = ppt_validate_slides(slide_numbers=[1, 2], presentation_path=str(p))
    assert res["success"] is True
    assert res["total_slides_validated"] == 2
    assert len(res["slides"]) == 2


def test_table_border_ooxml_schema_ordering(table_prs):
    """Verify _set_cell_border inserts line tags in valid CT_TableCellProperties sequence before solidFill."""
    deck_path, tbl_id = table_prs

    # Apply both fill and border
    res = ppt_style_table(
        slide_number=1,
        table_shape_id=tbl_id,
        range="0:0",
        style={
            "fill": "1E3A8A",
            "borders": {"color": "FF0000", "width": 2.0, "sides": ["top", "left", "bottom", "right"]},
        },
        presentation_path=str(deck_path),
    )
    assert res["success"] is True

    # Re-open deck and inspect raw XML of cell (0,0)
    prs = Presentation(str(deck_path))
    tbl = prs.slides[0].shapes[0].table
    tcPr = tbl.cell(0, 0)._tc.tcPr

    tags = [child.tag.split("}")[-1] if "}" in child.tag else child.tag for child in tcPr]
    # Check that line tags precede solidFill
    line_tags = ["lnL", "lnR", "lnT", "lnB"]
    fill_index = tags.index("solidFill")
    for lt in line_tags:
        assert lt in tags
        lt_index = tags.index(lt)
        assert lt_index < fill_index, f"{lt} (idx {lt_index}) must precede solidFill (idx {fill_index}) in CT_TableCellProperties"


def test_render_preserves_non_16_9_aspect_ratios(tmp_path):
    """Verify slide rendering dynamically derives dimensions and preserves 4:3 and portrait aspect ratios."""
    from powerpoint_mcp.tools.rendering import _get_presentation_dimensions, ppt_render_slide

    # 1. Standard 4:3 slide (10.0" x 7.5")
    prs_4_3 = Presentation()
    prs_4_3.slide_width = Inches(10.0)
    prs_4_3.slide_height = Inches(7.5)
    s = prs_4_3.slides.add_slide(prs_4_3.slide_layouts[6])
    s.shapes.add_shape(1, Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.0))
    p_4_3 = tmp_path / "deck_4_3.pptx"
    prs_4_3.save(str(p_4_3))

    w_4_3, h_4_3 = _get_presentation_dimensions(str(p_4_3), dpi=150)
    assert w_4_3 == 1500  # 10.0 * 150
    assert h_4_3 == 1125  # 7.5 * 150
    assert round(w_4_3 / h_4_3, 3) == round(4.0 / 3.0, 3)

    res_4_3 = ppt_render_slide(slide_number=1, renderer="mock", dpi=150, presentation_path=str(p_4_3))
    assert res_4_3["success"] is True
    assert res_4_3["width_px"] == 1500
    assert res_4_3["height_px"] == 1125

    # 2. Portrait slide (7.5" x 10.0")
    prs_port = Presentation()
    prs_port.slide_width = Inches(7.5)
    prs_port.slide_height = Inches(10.0)
    s_p = prs_port.slides.add_slide(prs_port.slide_layouts[6])
    s_p.shapes.add_shape(1, Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.0))
    p_port = tmp_path / "deck_portrait.pptx"
    prs_port.save(str(p_port))

    w_p, h_p = _get_presentation_dimensions(str(p_port), dpi=150)
    assert w_p == 1125  # 7.5 * 150
    assert h_p == 1500  # 10.0 * 150
    assert round(h_p / w_p, 3) == round(4.0 / 3.0, 3)

    res_port = ppt_render_slide(slide_number=1, renderer="mock", dpi=150, presentation_path=str(p_port))
    assert res_port["success"] is True
    assert res_port["width_px"] == 1125
    assert res_port["height_px"] == 1500
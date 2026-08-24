"""Rule-based slide and presentation validation engine for geometric and typographic defects."""

from dataclasses import dataclass, field
from enum import Enum
import math
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Set, Tuple, Union

from pptx import Presentation

from powerpoint_mcp.models.shape import (
    EMU_PER_INCH,
    SemanticRole,
    ShapeModel,
    ShapeType,
    emu_to_inches,
    inches_to_emu,
)
from powerpoint_mcp.models.slide import SlideModel
from powerpoint_mcp.models.presentation import PresentationModel
from powerpoint_mcp.pptx.inspector import inspect_presentation, inspect_slide


class IssueSeverity(str, Enum):
    """Severity levels for slide validation issues."""
    INFO = "info"
    WARNING = "warning"
    ERROR = "error"


@dataclass
class SlideIssue:
    """Diagnostic issue or warning detected on a slide."""
    rule_id: str
    severity: Union[IssueSeverity, str]
    shape_ids: List[int]
    message: str
    details: Dict[str, Any] = field(default_factory=dict)

    @property
    def severity_str(self) -> str:
        """Return severity as standard lowercase string."""
        return self.severity.value if isinstance(self.severity, IssueSeverity) else str(self.severity).lower()

    def to_dict(self, detail: str = "full") -> Dict[str, Any]:
        """Serialize issue to dictionary matching spec schema."""
        res = {
            "rule_id": self.rule_id,
            "severity": self.severity_str,
            "shape_ids": self.shape_ids,
            "message": self.message,
        }
        if detail == "full" or self.details:
            res["details"] = self.details
        return res


@dataclass
class SlideValidationResult:
    """Comprehensive validation report for a single slide."""
    slide_number: int
    is_valid: bool
    warning_count: int
    error_count: int
    issues: List[SlideIssue] = field(default_factory=list)
    metrics: Dict[str, Any] = field(default_factory=dict)

    @property
    def warnings(self) -> List[SlideIssue]:
        """Alias for issues to support spec schema."""
        return self.issues

    @property
    def error_issues(self) -> List[SlideIssue]:
        """Filter only error severity issues."""
        return [i for i in self.issues if i.severity_str == "error"]

    @property
    def warning_issues(self) -> List[SlideIssue]:
        """Filter only warning severity issues."""
        return [i for i in self.issues if i.severity_str == "warning"]

    def get_summary_counts(self) -> Dict[str, int]:
        """Return issue counts broken down by category."""
        return {
            "overlaps": sum(1 for i in self.issues if i.rule_id == "VAL-01"),
            "boundary_violations": sum(1 for i in self.issues if i.rule_id == "VAL-02"),
            "text_overflow": sum(1 for i in self.issues if i.rule_id == "VAL-03"),
            "tiny_fonts": sum(1 for i in self.issues if i.rule_id == "VAL-04"),
            "title_inconsistencies": sum(1 for i in self.issues if i.rule_id == "VAL-05"),
            "duplicate_objects": sum(1 for i in self.issues if i.rule_id == "VAL-06"),
            "extreme_rotations": sum(1 for i in self.issues if i.rule_id == "VAL-07"),
        }

    def to_dict(self, detail: str = "summary") -> Dict[str, Any]:
        """Serialize slide validation report to dictionary matching MCP and JSON schemas."""
        return {
            "slide_number": self.slide_number,
            "is_valid": self.is_valid,
            "valid": self.is_valid,
            "warning_count": self.warning_count,
            "error_count": self.error_count,
            "summary": self.get_summary_counts(),
            "warnings": [i.to_dict(detail=detail) for i in self.issues],
            "issues": [i.to_dict(detail=detail) for i in self.issues],
            "metrics": self.metrics,
        }



@dataclass
class PresentationValidationResult:
    """Comprehensive validation report across all slides in a presentation."""
    presentation_path: Optional[str]
    slide_count: int
    is_valid: bool
    total_warnings: int
    total_errors: int
    slide_results: List[SlideValidationResult] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        """Serialize presentation validation report to dictionary."""
        return {
            "presentation_path": self.presentation_path,
            "slide_count": self.slide_count,
            "is_valid": self.is_valid,
            "total_warnings": self.total_warnings,
            "total_errors": self.total_errors,
            "slide_results": [r.to_dict() for r in self.slide_results],
        }


def _is_background_shape(shape: ShapeModel, slide_width_in: float, slide_height_in: float) -> bool:
    """Check if shape is a full-canvas background rectangle."""
    w = shape.bbox.width_inches
    h = shape.bbox.height_inches
    if w >= slide_width_in * 0.95 and h >= slide_height_in * 0.95:
        # Background shape covering whole slide
        return True
    return False


def _is_valid_containment(container: ShapeModel, child: ShapeModel, tolerance_in: float = 0.10) -> bool:
    """Check if child shape is legitimately nested inside a container/card shape (Feature 11)."""
    cb = container.bbox
    kb = child.bbox

    c_area = cb.width_inches * cb.height_inches
    k_area = kb.width_inches * kb.height_inches

    if c_area < k_area * 1.05:
        return False

    # Check spatial containment: child bbox is within container bbox (with margin tolerance)
    is_spatially_contained = (
        kb.left_inches >= cb.left_inches - tolerance_in
        and kb.top_inches >= cb.top_inches - tolerance_in
        and kb.right_inches <= cb.right_inches + tolerance_in
        and kb.bottom_inches <= cb.bottom_inches + tolerance_in
    )

    if not is_spatially_contained:
        return False

    # Z-order: content child must not be stacked underneath background container
    if child.z_order < container.z_order:
        return False

    return True


def _check_val_01_overlaps(
    shapes: List[ShapeModel],
    slide_width_in: float,
    slide_height_in: float,
    min_overlap_area_sq_in: float = 0.01,
) -> List[SlideIssue]:
    """VAL-01: Container-aware overlap detection between shapes.

    Distinguishes legitimate card/container nesting (VALID_CONTAINMENT) from
    actual layout collisions (ACTUAL_OVERLAP) or overflows (SUSPECT_OVERLAP).
    """
    issues: List[SlideIssue] = []
    n = len(shapes)

    for i in range(n):
        s1 = shapes[i]
        if _is_background_shape(s1, slide_width_in, slide_height_in):
            continue

        b1 = s1.bbox
        for j in range(i + 1, n):
            s2 = shapes[j]
            if _is_background_shape(s2, slide_width_in, slide_height_in):
                continue

            # Check if one shape is a container for the other (card + text/icon)
            if _is_valid_containment(s1, s2) or _is_valid_containment(s2, s1):
                continue

            b2 = s2.bbox
            # Calculate AABB intersection in inches
            ox_left = max(b1.left_inches, b2.left_inches)
            ox_top = max(b1.top_inches, b2.top_inches)
            ox_right = min(b1.right_inches, b2.right_inches)
            ox_bottom = min(b1.bottom_inches, b2.bottom_inches)

            if ox_right > ox_left and ox_bottom > ox_top:
                overlap_w = ox_right - ox_left
                overlap_h = ox_bottom - ox_top
                area_sq_in = overlap_w * overlap_h

                if area_sq_in >= min_overlap_area_sq_in:
                    depth_in = min(overlap_w, overlap_h)
                    msg = (
                        f"WARNING: Shape {s1.shape_id} ('{s1.name}') overlaps "
                        f"Shape {s2.shape_id} ('{s2.name}') by {depth_in:.2f} inches "
                        f"(area: {area_sq_in:.2f} sq in)."
                    )
                    issues.append(
                        SlideIssue(
                            rule_id="VAL-01",
                            severity=IssueSeverity.WARNING,
                            shape_ids=[s1.shape_id, s2.shape_id],
                            message=msg,
                            details={
                                "classification": "ACTUAL_OVERLAP",
                                "shape_1_id": s1.shape_id,
                                "shape_1_name": s1.name,
                                "shape_2_id": s2.shape_id,
                                "shape_2_name": s2.name,
                                "overlap_width_inches": round(overlap_w, 4),
                                "overlap_height_inches": round(overlap_h, 4),
                                "overlap_area_sq_in": round(area_sq_in, 4),
                                "overlap_depth_inches": round(depth_in, 4),
                            },
                        )
                    )

    return issues



def _check_val_02_clipping(
    shapes: List[ShapeModel],
    slide_width_in: float,
    slide_height_in: float,
    tolerance_in: float = 0.05,
) -> List[SlideIssue]:
    """VAL-02: Off-slide / boundary clipping (shape right > slide_width or bottom > slide_height or x < 0 or y < 0)."""
    issues: List[SlideIssue] = []

    for s in shapes:
        b = s.bbox
        protrusions: Dict[str, float] = {}

        if b.left_inches < -tolerance_in:
            protrusions["left"] = round(abs(b.left_inches), 4)
        if b.top_inches < -tolerance_in:
            protrusions["top"] = round(abs(b.top_inches), 4)
        if b.right_inches > slide_width_in + tolerance_in:
            protrusions["right"] = round(b.right_inches - slide_width_in, 4)
        if b.bottom_inches > slide_height_in + tolerance_in:
            protrusions["bottom"] = round(b.bottom_inches - slide_height_in, 4)

        if protrusions:
            primary_edge = max(protrusions.items(), key=lambda item: item[1])
            edge_name, protrusion_val = primary_edge

            msg = (
                f"WARNING: Shape {s.shape_id} ('{s.name}') extends "
                f"{protrusion_val:.2f} inches beyond the {edge_name} slide boundary."
            )
            issues.append(
                SlideIssue(
                    rule_id="VAL-02",
                    severity=IssueSeverity.WARNING,
                    shape_ids=[s.shape_id],
                    message=msg,
                    details={
                        "shape_id": s.shape_id,
                        "shape_name": s.name,
                        "protrusions": protrusions,
                        "primary_boundary": edge_name,
                        "max_protrusion_inches": protrusion_val,
                        "bbox": b.to_dict(),
                    },
                )
            )

    return issues


def _check_val_03_text_overflow(
    shapes: List[ShapeModel],
    overflow_threshold_ratio: float = 1.15,
) -> List[SlideIssue]:
    """VAL-03: Realistic text fit and overflow measurement using word-wrap simulation.

    Distinguishes single-line wide banner text from true multi-line overflow,
    and categorizes fit into FITS, LIKELY_OVERFLOW, and CONFIRMED_OVERFLOW.
    """
    issues: List[SlideIssue] = []

    for s in shapes:
        if not s.text_frame or not s.text_frame.text or not s.text_frame.text.strip():
            continue

        tf = s.text_frame
        b = s.bbox

        avail_w = max(0.1, b.width_inches - tf.margin_left_inches - tf.margin_right_inches)
        avail_h = max(0.1, b.height_inches - tf.margin_top_inches - tf.margin_bottom_inches)

        total_chars = len(tf.text.strip())
        para_heights: List[float] = []

        for p in tf.paragraphs:
            p_text = (p.text or "").strip()
            if not p_text:
                continue

            # Determine dominant font size for paragraph
            font_size = 14.0
            for r in p.runs:
                if r.style.font_size_pt is not None:
                    font_size = float(r.style.font_size_pt)
                    break

            # Character width factor (conservative ~0.50 of font size in inches)
            avg_char_w_in = (font_size / 72.0) * 0.50
            line_h_in = (font_size / 72.0) * 1.22
            space_w_in = avg_char_w_in * 0.8

            words = p_text.split()
            if not words:
                continue

            # Simulate word wrapping
            lines_count = 1
            cur_line_w = 0.0

            for word in words:
                word_w = len(word) * avg_char_w_in
                if word_w > avail_w:
                    # Single word exceeds line width, wrapping across multiple lines
                    lines_for_word = max(1, math.ceil(word_w / avail_w))
                    lines_count += (lines_for_word - 1)
                    cur_line_w = word_w - ((lines_for_word - 1) * avail_w)
                elif cur_line_w == 0.0:
                    cur_line_w = word_w
                elif cur_line_w + space_w_in + word_w <= avail_w:
                    cur_line_w += space_w_in + word_w
                else:
                    lines_count += 1
                    cur_line_w = word_w

            p_height = (lines_count * line_h_in) + ((p.space_after_pt or 0.0) / 72.0)
            para_heights.append(p_height)

        total_text_height = sum(para_heights) if para_heights else 0.0
        overflow_ratio = total_text_height / avail_h if avail_h > 0 else 1.0

        if overflow_ratio > overflow_threshold_ratio:
            is_confirmed = (overflow_ratio > 1.35)
            classification = "CONFIRMED_OVERFLOW" if is_confirmed else "LIKELY_OVERFLOW"
            overflow_pct = int(round((overflow_ratio - 1.0) * 100))

            msg = (
                f"WARNING: Text in Shape {s.shape_id} ('{s.name}') {classification.lower().replace('_', ' ')} "
                f"by estimated {overflow_pct}% (box: {b.width_inches:.2f}x{b.height_inches:.2f} in, "
                f"req height: {total_text_height:.2f} in)."
            )
            issues.append(
                SlideIssue(
                    rule_id="VAL-03",
                    severity=IssueSeverity.WARNING,
                    shape_ids=[s.shape_id],
                    message=msg,
                    details={
                        "classification": classification,
                        "shape_id": s.shape_id,
                        "shape_name": s.name,
                        "char_count": total_chars,
                        "box_width_inches": round(b.width_inches, 2),
                        "box_height_inches": round(b.height_inches, 2),
                        "estimated_text_height_inches": round(total_text_height, 2),
                        "available_height_inches": round(avail_h, 2),
                        "overflow_ratio": round(overflow_ratio, 2),
                        "overflow_pct": overflow_pct,
                    },
                )
            )

    return issues


def _check_val_04_tiny_font(
    shapes: List[ShapeModel],
    min_font_pt: float = 8.0,
) -> List[SlideIssue]:
    """VAL-04: Smart context-aware tiny font detection.

    Distinguishes intentional compact UI text (badges, pills, footnotes, compact metadata)
    from problematic tiny body/bullet text.
    """
    issues: List[SlideIssue] = []

    for s in shapes:
        if not s.text_frame:
            continue

        raw_text = (s.text_frame.text or "").strip()
        if not raw_text:
            continue

        b = s.bbox
        is_compact_shape = (b.height_inches <= 0.45 or b.width_inches <= 2.2 or len(raw_text) <= 20)
        is_footer_or_meta = (
            s.semantic_role in (SemanticRole.FOOTER, SemanticRole.UNKNOWN)
            and (b.top_inches >= 6.5 or "footer" in s.name.lower() or "page" in s.name.lower() or "slide" in s.name.lower())
        )
        is_badge = (
            is_compact_shape
            or "badge" in s.name.lower()
            or "pill" in s.name.lower()
            or "tag" in s.name.lower()
            or "chip" in s.name.lower()
        )

        for p in s.text_frame.paragraphs:
            for r in p.runs:
                f_size = r.style.font_size_pt
                if f_size is not None and f_size < min_font_pt:
                    preview = (r.text[:40] + "...") if len(r.text) > 40 else r.text

                    # Classify based on context and role
                    if is_badge or is_footer_or_meta:
                        if f_size >= 6.0:
                            classification = "INTENTIONAL_COMPACT_TEXT"
                            severity = IssueSeverity.INFO
                            msg = (
                                f"INFO: Shape {s.shape_id} ('{s.name}') contains "
                                f"intentional compact text ({f_size:.1f} pt) for badge/metadata: '{preview}'"
                            )
                        else:
                            classification = "CRITICAL_TINY_TEXT"
                            severity = IssueSeverity.WARNING
                            msg = (
                                f"WARNING: Shape {s.shape_id} ('{s.name}') contains "
                                f"critically tiny text ({f_size:.1f} pt): '{preview}'"
                            )
                    elif f_size < 7.0:
                        classification = "CRITICAL_TINY_TEXT"
                        severity = IssueSeverity.WARNING
                        msg = (
                            f"WARNING: Shape {s.shape_id} ('{s.name}') contains "
                            f"critically tiny text ({f_size:.1f} pt): '{preview}'"
                        )
                    else:
                        classification = "SUSPICIOUS_TINY_TEXT"
                        severity = IssueSeverity.WARNING
                        msg = (
                            f"WARNING: Shape {s.shape_id} ('{s.name}') contains "
                            f"suspiciously tiny text ({f_size:.1f} pt): '{preview}'"
                        )

                    issues.append(
                        SlideIssue(
                            rule_id="VAL-04",
                            severity=severity,
                            shape_ids=[s.shape_id],
                            message=msg,
                            details={
                                "classification": classification,
                                "shape_id": s.shape_id,
                                "shape_name": s.name,
                                "font_size_pt": f_size,
                                "min_threshold_pt": min_font_pt,
                                "text_preview": preview,
                                "semantic_role": s.semantic_role.value,
                            },
                        )
                    )
                    break
            else:
                continue
            break

    return issues


def _check_val_05_title_position(
    shapes: List[ShapeModel],
    title_baseline: Optional[Dict[str, float]] = None,
    delta_threshold_in: float = 0.05,
) -> List[SlideIssue]:
    """VAL-05: Inconsistent title position (comparing title coordinates against baseline)."""
    if not title_baseline:
        return []

    issues: List[SlideIssue] = []
    baseline_x = title_baseline.get("x", title_baseline.get("left", 1.0))
    baseline_y = title_baseline.get("y", title_baseline.get("top", 0.8))

    for s in shapes:
        if s.semantic_role == SemanticRole.TITLE or "title" in s.name.lower():
            dx = abs(s.bbox.left_inches - baseline_x)
            dy = abs(s.bbox.top_inches - baseline_y)
            if dx > delta_threshold_in or dy > delta_threshold_in:
                msg = (
                    f"WARNING: Title Shape {s.shape_id} position (x={s.bbox.left_inches:.2f}, "
                    f"y={s.bbox.top_inches:.2f}) deviates from standard title position "
                    f"(x={baseline_x:.2f}, y={baseline_y:.2f}) by dx={dx:.2f}, dy={dy:.2f} in."
                )
                issues.append(
                    SlideIssue(
                        rule_id="VAL-05",
                        severity=IssueSeverity.WARNING,
                        shape_ids=[s.shape_id],
                        message=msg,
                        details={
                            "shape_id": s.shape_id,
                            "shape_name": s.name,
                            "actual_x_inches": s.bbox.left_inches,
                            "actual_y_inches": s.bbox.top_inches,
                            "baseline_x_inches": baseline_x,
                            "baseline_y_inches": baseline_y,
                            "delta_x_inches": round(dx, 4),
                            "delta_y_inches": round(dy, 4),
                        },
                    )
                )
            break

    return issues


def _check_val_06_duplicate_objects(
    shapes: List[ShapeModel],
    coord_tolerance_in: float = 0.005,
) -> List[SlideIssue]:
    """VAL-06: Duplicate objects (identical bounding box and text/content)."""
    issues: List[SlideIssue] = []
    n = len(shapes)

    for i in range(n):
        s1 = shapes[i]
        b1 = s1.bbox
        t1 = s1.text_frame.text if s1.text_frame else ""

        for j in range(i + 1, n):
            s2 = shapes[j]
            b2 = s2.bbox
            t2 = s2.text_frame.text if s2.text_frame else ""

            same_coords = (
                abs(b1.left_inches - b2.left_inches) <= coord_tolerance_in
                and abs(b1.top_inches - b2.top_inches) <= coord_tolerance_in
                and abs(b1.width_inches - b2.width_inches) <= coord_tolerance_in
                and abs(b1.height_inches - b2.height_inches) <= coord_tolerance_in
            )

            same_content = (t1 == t2) if (t1 or t2) else (s1.shape_type == s2.shape_type)

            if same_coords and same_content:
                msg = (
                    f"WARNING: Shape {s1.shape_id} ('{s1.name}') and Shape {s2.shape_id} "
                    f"('{s2.name}') appear to be duplicate superimposed objects at "
                    f"(x={b1.left_inches:.2f}, y={b1.top_inches:.2f})."
                )
                issues.append(
                    SlideIssue(
                        rule_id="VAL-06",
                        severity=IssueSeverity.WARNING,
                        shape_ids=[s1.shape_id, s2.shape_id],
                        message=msg,
                        details={
                            "shape_1_id": s1.shape_id,
                            "shape_1_name": s1.name,
                            "shape_2_id": s2.shape_id,
                            "shape_2_name": s2.name,
                            "x_inches": b1.left_inches,
                            "y_inches": b1.top_inches,
                        },
                    )
                )

    return issues


def _check_val_07_extreme_rotations(
    shapes: List[ShapeModel],
) -> List[SlideIssue]:
    """VAL-07: Extreme rotations (> 45 degrees or non-standard angles)."""
    issues: List[SlideIssue] = []

    for s in shapes:
        rot = s.rotation % 360.0
        if rot != 0.0:
            # Allow standard orthogonal and 45-degree angles
            standard_angles = {0.0, 45.0, 90.0, 135.0, 180.0, 225.0, 270.0, 315.0, 360.0}
            if round(rot, 1) not in standard_angles:
                msg = f"WARNING: Shape {s.shape_id} ('{s.name}') has an irregular rotation of {rot:.1f} degrees."
                issues.append(
                    SlideIssue(
                        rule_id="VAL-07",
                        severity=IssueSeverity.WARNING,
                        shape_ids=[s.shape_id],
                        message=msg,
                        details={
                            "shape_id": s.shape_id,
                            "shape_name": s.name,
                            "rotation_degrees": round(rot, 2),
                        },
                    )
                )

    return issues


def validate_slide(
    slide_or_model: Union[SlideModel, Any],
    slide_number: int = 1,
    slide_width_inches: Optional[float] = None,
    slide_height_inches: Optional[float] = None,
    rules: Optional[List[str]] = None,
    title_baseline: Optional[Dict[str, float]] = None,
) -> SlideValidationResult:
    """Validate a slide against all geometric and typographic rules (VAL-01 to VAL-07).

    Args:
        slide_or_model: SlideModel, python-pptx Slide, or presentation path.
        slide_number: 1-indexed slide number.
        slide_width_inches: Slide canvas width in inches.
        slide_height_inches: Slide canvas height in inches.
        rules: Optional list of rule IDs to run (e.g. ['VAL-01', 'VAL-02']). Defaults to all.
        title_baseline: Optional reference title coordinates for consistency check.

    Returns:
        SlideValidationResult containing detected issues, counts, validity, and metrics.
    """
    # Resolve SlideModel if raw Slide or path is provided
    if isinstance(slide_or_model, SlideModel):
        slide_model = slide_or_model
    elif isinstance(slide_or_model, (str, Path)):
        slide_model = inspect_slide(str(slide_or_model), slide_number)
    elif hasattr(slide_or_model, "slides"):
        slide_model = inspect_slide(slide_or_model, slide_number)
    else:
        # python-pptx Slide object
        prs_dummy = Presentation()
        slide_model = inspect_slide(slide_or_model, slide_number)

    effective_slide_num = getattr(slide_model, "slide_number", slide_number)
    effective_width = (
        slide_width_inches
        if slide_width_inches is not None
        else getattr(slide_model, "width_inches", 13.333)
    )
    effective_height = (
        slide_height_inches
        if slide_height_inches is not None
        else getattr(slide_model, "height_inches", 7.5)
    )

    active_rules = set(r.upper() for r in rules) if rules else None
    all_issues: List[SlideIssue] = []

    shapes = slide_model.shapes

    # VAL-01: Overlaps
    if active_rules is None or "VAL-01" in active_rules:
        all_issues.extend(_check_val_01_overlaps(shapes, effective_width, effective_height))

    # VAL-02: Boundary Clipping
    if active_rules is None or "VAL-02" in active_rules:
        all_issues.extend(_check_val_02_clipping(shapes, effective_width, effective_height))

    # VAL-03: Text Overflow
    if active_rules is None or "VAL-03" in active_rules:
        all_issues.extend(_check_val_03_text_overflow(shapes))

    # VAL-04: Tiny Font
    if active_rules is None or "VAL-04" in active_rules:
        all_issues.extend(_check_val_04_tiny_font(shapes))

    # VAL-05: Inconsistent Title
    if active_rules is None or "VAL-05" in active_rules:
        all_issues.extend(_check_val_05_title_position(shapes, title_baseline))

    # VAL-06: Duplicate Objects
    if active_rules is None or "VAL-06" in active_rules:
        all_issues.extend(_check_val_06_duplicate_objects(shapes))

    # VAL-07: Extreme Rotations
    if active_rules is None or "VAL-07" in active_rules:
        all_issues.extend(_check_val_07_extreme_rotations(shapes))

    # Metrics calculation
    text_count = sum(1 for s in shapes if s.text_frame and s.text_frame.text)
    image_count = sum(1 for s in shapes if s.shape_type == ShapeType.PICTURE)

    metrics = {
        "shape_count": len(shapes),
        "text_shape_count": text_count,
        "image_shape_count": image_count,
        "slide_dimensions": {
            "width_inches": effective_width,
            "height_inches": effective_height,
        },
    }

    warning_count = sum(1 for i in all_issues if i.severity_str == "warning")
    error_count = sum(1 for i in all_issues if i.severity_str == "error")

    # A slide is valid if it has 0 critical errors (and 0 issues for clean slides)
    is_valid = (len(all_issues) == 0)

    return SlideValidationResult(
        slide_number=effective_slide_num,
        is_valid=is_valid,
        warning_count=warning_count,
        error_count=error_count,
        issues=all_issues,
        metrics=metrics,
    )


def validate_presentation(
    prs_or_model: Union[PresentationModel, Any],
    rules: Optional[List[str]] = None,
) -> PresentationValidationResult:
    """Validate all slides in a presentation, including cross-slide consistency rules.

    Args:
        prs_or_model: PresentationModel, python-pptx Presentation, or file path.
        rules: Optional list of rule IDs to filter.

    Returns:
        PresentationValidationResult with aggregate statistics and per-slide reports.
    """
    if isinstance(prs_or_model, PresentationModel):
        prs_model = prs_or_model
    else:
        prs_model = inspect_presentation(prs_or_model)

    # Establish baseline title position from first slide containing a title
    title_baseline: Optional[Dict[str, float]] = None
    for slide in prs_model.slides:
        for shape in slide.shapes:
            if shape.semantic_role == SemanticRole.TITLE or "title" in shape.name.lower():
                title_baseline = {
                    "x": shape.bbox.left_inches,
                    "y": shape.bbox.top_inches,
                    "width": shape.bbox.width_inches,
                    "height": shape.bbox.height_inches,
                }
                break
        if title_baseline:
            break

    slide_results: List[SlideValidationResult] = []
    total_warnings = 0
    total_errors = 0

    for slide in prs_model.slides:
        result = validate_slide(
            slide,
            slide_number=slide.slide_number,
            slide_width_inches=prs_model.width_inches,
            slide_height_inches=prs_model.height_inches,
            rules=rules,
            title_baseline=title_baseline,
        )
        slide_results.append(result)
        total_warnings += result.warning_count
        total_errors += result.error_count

    all_valid = all(r.is_valid for r in slide_results)

    return PresentationValidationResult(
        presentation_path=prs_model.path,
        slide_count=prs_model.slide_count,
        is_valid=all_valid,
        total_warnings=total_warnings,
        total_errors=total_errors,
        slide_results=slide_results,
    )

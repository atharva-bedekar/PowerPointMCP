"""Programmatic synthetic PowerPoint deck generator for testing.

Generates a standard 3-slide widescreen (16:9) presentation with:
- Slide 1: Quarterly Performance Overview (Title, Subtitle, 3 KPI boxes, Picture shape)
- Slide 2: Operational Architecture (Title, 2-column layout, 3-step process diagram, Footer)
- Slide 3: Audit & Compliance Issues (Intentional defects: heavy overlaps, boundary clipping, tiny fonts, text overflow)
"""

import io
import os
import sys
from pathlib import Path
from typing import Optional, Union

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def create_placeholder_image_bytes(
    width: int = 800, height: int = 200, title: str = "System Performance Dashboard"
) -> io.BytesIO:
    """Generate a clean synthetic chart banner in memory using Pillow."""
    img = Image.new("RGB", (width, height), color=(240, 244, 248))
    draw = ImageDraw.Draw(img)

    # Draw border
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=(180, 200, 220), width=2)

    # Draw simulated chart bars
    bar_colors = [
        (41, 128, 185),
        (39, 174, 96),
        (230, 126, 34),
        (142, 68, 173),
        (52, 152, 219),
        (46, 204, 113),
        (241, 196, 15),
        (231, 76, 60),
    ]
    bar_heights = [60, 95, 130, 85, 140, 110, 160, 135]
    start_x = 50
    spacing = 70
    bar_w = 40
    base_y = height - 30

    for i, (h, col) in enumerate(zip(bar_heights, bar_colors)):
        x0 = start_x + i * spacing
        y0 = base_y - h
        draw.rectangle([(x0, y0), (x0 + bar_w, base_y)], fill=col)

    # Draw horizontal baseline
    draw.line([(30, base_y), (width - 30, base_y)], fill=(120, 140, 160), width=2)

    # Text indicator
    draw.text((width - 320, 20), title, fill=(40, 60, 80))

    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)
    return img_bytes


def build_synthetic_deck() -> Presentation:
    """Construct a 3-slide Presentation object in memory."""
    prs = Presentation()
    # Configure 16:9 widescreen dimensions (13.333 x 7.5 inches = 12,192,000 x 6,858,000 EMUs)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: Quarterly Performance Overview
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    # 1. Slide Title
    title_box = slide1.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(0.8),
        width=Inches(11.333),
        height=Inches(0.9),
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p_title = title_tf.paragraphs[0]
    p_title.text = "Quarterly Performance Overview"
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(38)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(31, 73, 125)

    # 2. Subtitle
    sub_box = slide1.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(1.7),
        width=Inches(11.333),
        height=Inches(0.6),
    )
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True
    p_sub = sub_tf.paragraphs[0]
    p_sub.text = "Q3 2026 Executive Summary"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(18)
    p_sub.font.italic = True
    p_sub.font.color.rgb = RGBColor(89, 89, 89)

    # 3. Horizontal KPI Boxes (3 boxes)
    kpis = [
        {
            "name": "KPI Card: Revenue",
            "x": Inches(1.0),
            "y": Inches(2.5),
            "w": Inches(3.5),
            "h": Inches(2.2),
            "fill": RGBColor(230, 242, 255),
            "line": RGBColor(41, 128, 185),
            "label": "Total Revenue",
            "val": "$12.4M",
            "sub": "+18.2% vs Plan",
            "accent": RGBColor(41, 128, 185),
        },
        {
            "name": "KPI Card: Users",
            "x": Inches(4.9),
            "y": Inches(2.5),
            "w": Inches(3.5),
            "h": Inches(2.2),
            "fill": RGBColor(232, 248, 235),
            "line": RGBColor(39, 174, 96),
            "label": "Active Users",
            "val": "450K",
            "sub": "+24% YoY Growth",
            "accent": RGBColor(39, 174, 96),
        },
        {
            "name": "KPI Card: NPS",
            "x": Inches(8.8),
            "y": Inches(2.5),
            "w": Inches(3.5),
            "h": Inches(2.2),
            "fill": RGBColor(254, 243, 226),
            "line": RGBColor(230, 126, 34),
            "label": "Customer NPS",
            "val": "78 / 100",
            "sub": "Top Quartile Benchmark",
            "accent": RGBColor(230, 126, 34),
        },
    ]

    for kpi in kpis:
        shape = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            kpi["x"],
            kpi["y"],
            kpi["w"],
            kpi["h"],
        )
        shape.name = kpi["name"]
        shape.fill.solid()
        shape.fill.fore_color.rgb = kpi["fill"]
        shape.line.color.rgb = kpi["line"]
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        # Label paragraph
        p1 = tf.paragraphs[0]
        p1.text = kpi["label"]
        p1.font.name = "Calibri"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(80, 80, 80)
        p1.space_after = Pt(4)

        # Value paragraph
        p2 = tf.add_paragraph()
        p2.text = kpi["val"]
        p2.font.name = "Calibri"
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = kpi["accent"]
        p2.space_after = Pt(4)

        # Subtext paragraph
        p3 = tf.add_paragraph()
        p3.text = kpi["sub"]
        p3.font.name = "Calibri"
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(100, 100, 100)

    # 4. Image / Picture placeholder shape
    img_stream = create_placeholder_image_bytes(
        width=1100, height=180, title="Quarterly System Activity Benchmark"
    )
    pic = slide1.shapes.add_picture(
        img_stream,
        left=Inches(1.0),
        top=Inches(5.0),
        width=Inches(11.3),
        height=Inches(1.8),
    )
    pic.name = "Performance Dashboard Chart"

    # =========================================================================
    # SLIDE 2: Operational Architecture
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)

    # 1. Slide Title
    title_box2 = slide2.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(0.8),
        width=Inches(11.333),
        height=Inches(0.9),
    )
    title_tf2 = title_box2.text_frame
    title_tf2.word_wrap = True
    p2_title = title_tf2.paragraphs[0]
    p2_title.text = "Operational Architecture"
    p2_title.font.name = "Calibri"
    p2_title.font.size = Pt(38)
    p2_title.font.bold = True
    p2_title.font.color.rgb = RGBColor(31, 73, 125)

    # 2. Left Column: Key Initiatives
    left_col = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(1.9),
        Inches(5.4),
        Inches(2.7),
    )
    left_col.name = "Column: Key Initiatives"
    left_col.fill.solid()
    left_col.fill.fore_color.rgb = RGBColor(248, 249, 250)
    left_col.line.color.rgb = RGBColor(200, 205, 210)
    left_col.line.width = Pt(1.0)

    ltf = left_col.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.25)
    ltf.margin_top = Inches(0.2)
    ltf.margin_right = Inches(0.25)
    ltf.margin_bottom = Inches(0.2)

    lp0 = ltf.paragraphs[0]
    lp0.text = "Key Initiatives"
    lp0.font.name = "Calibri"
    lp0.font.size = Pt(16)
    lp0.font.bold = True
    lp0.font.color.rgb = RGBColor(31, 73, 125)
    lp0.space_after = Pt(8)

    initiatives = [
        "Automated ingestion pipeline migration to high-speed bus",
        "Deterministic geometric alignment & text run preservation",
        "Headless COM & LibreOffice cross-platform rendering engine",
        "Antigravity skill integration with recursive verification loop",
    ]
    for init_text in initiatives:
        lp = ltf.add_paragraph()
        lp.text = f"•  {init_text}"
        lp.font.name = "Calibri"
        lp.font.size = Pt(12)
        lp.font.color.rgb = RGBColor(60, 60, 60)
        lp.space_after = Pt(4)

    # 3. Right Column: Milestones
    right_col = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.9),
        Inches(1.9),
        Inches(5.4),
        Inches(2.7),
    )
    right_col.name = "Column: Milestones"
    right_col.fill.solid()
    right_col.fill.fore_color.rgb = RGBColor(248, 249, 250)
    right_col.line.color.rgb = RGBColor(200, 205, 210)
    right_col.line.width = Pt(1.0)

    rtf = right_col.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.25)
    rtf.margin_top = Inches(0.2)
    rtf.margin_right = Inches(0.25)
    rtf.margin_bottom = Inches(0.2)

    rp0 = rtf.paragraphs[0]
    rp0.text = "Program Milestones"
    rp0.font.name = "Calibri"
    rp0.font.size = Pt(16)
    rp0.font.bold = True
    rp0.font.color.rgb = RGBColor(31, 73, 125)
    rp0.space_after = Pt(8)

    milestones = [
        "M1: Inspection Engine & Shape Hierarchy Spec (Done)",
        "M2: Coordinate Geometry & Text Manipulation (In Progress)",
        "M3: Multi-Engine Slide Rendering & Visual Diffing",
        "M4: Session Safety, Versioning & Rule Validation",
    ]
    for ms_text in milestones:
        rp = rtf.add_paragraph()
        rp.text = f"•  {ms_text}"
        rp.font.name = "Calibri"
        rp.font.size = Pt(12)
        rp.font.color.rgb = RGBColor(60, 60, 60)
        rp.space_after = Pt(4)

    # 4. Diagram Process Flow (3 interconnected process boxes)
    steps = [
        {
            "name": "Process Box 1: Ingest",
            "x": Inches(1.0),
            "y": Inches(4.9),
            "w": Inches(3.4),
            "h": Inches(1.4),
            "fill": RGBColor(41, 128, 185),
            "title": "1. Ingest & Inspect",
            "desc": "Extract EMU coordinates, text frames, styles, and shape tree",
        },
        {
            "name": "Process Box 2: Transform",
            "x": Inches(4.95),
            "y": Inches(4.9),
            "w": Inches(3.4),
            "h": Inches(1.4),
            "fill": RGBColor(39, 174, 96),
            "title": "2. Transform & Edit",
            "desc": "Deterministic math alignment, text update, and OOXML patching",
        },
        {
            "name": "Process Box 3: Verify",
            "x": Inches(8.9),
            "y": Inches(4.9),
            "w": Inches(3.4),
            "h": Inches(1.4),
            "fill": RGBColor(142, 68, 173),
            "title": "3. Render & Verify",
            "desc": "Headless PNG generation, pixel diffing, and rule validation",
        },
    ]

    for step in steps:
        p_shape = slide2.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            step["x"],
            step["y"],
            step["w"],
            step["h"],
        )
        p_shape.name = step["name"]
        p_shape.fill.solid()
        p_shape.fill.fore_color.rgb = step["fill"]
        p_shape.line.fill.background()

        ptf = p_shape.text_frame
        ptf.word_wrap = True
        ptf.margin_left = Inches(0.3)
        ptf.margin_top = Inches(0.15)
        ptf.margin_right = Inches(0.3)
        ptf.margin_bottom = Inches(0.15)

        s_p1 = ptf.paragraphs[0]
        s_p1.text = step["title"]
        s_p1.font.name = "Calibri"
        s_p1.font.size = Pt(14)
        s_p1.font.bold = True
        s_p1.font.color.rgb = RGBColor(255, 255, 255)
        s_p1.alignment = PP_ALIGN.CENTER
        s_p1.space_after = Pt(2)

        s_p2 = ptf.add_paragraph()
        s_p2.text = step["desc"]
        s_p2.font.name = "Calibri"
        s_p2.font.size = Pt(9.5)
        s_p2.font.color.rgb = RGBColor(240, 240, 240)
        s_p2.alignment = PP_ALIGN.CENTER

    # 5. Footer text
    footer_box = slide2.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(6.6),
        width=Inches(11.333),
        height=Inches(0.4),
    )
    footer_box.name = "Footer: Confidential"
    ftf = footer_box.text_frame
    p_ft = ftf.paragraphs[0]
    p_ft.text = "Confidential — Antigravity Operational Architecture Review"
    p_ft.font.name = "Calibri"
    p_ft.font.size = Pt(10)
    p_ft.font.italic = True
    p_ft.font.color.rgb = RGBColor(128, 128, 128)

    # =========================================================================
    # SLIDE 3: Audit & Compliance Issues (Intentional Defects for Validation)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)

    # 1. Slide Title
    title_box3 = slide3.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(0.8),
        width=Inches(11.333),
        height=Inches(0.9),
    )
    title_tf3 = title_box3.text_frame
    title_tf3.word_wrap = True
    p3_title = title_tf3.paragraphs[0]
    p3_title.text = "Audit & Compliance Issues"
    p3_title.font.name = "Calibri"
    p3_title.font.size = Pt(38)
    p3_title.font.bold = True
    p3_title.font.color.rgb = RGBColor(192, 0, 0)

    # 2. Defect 1: Two heavily overlapping shapes (Overlap > 0.5 inches)
    # Box A: left=1.0", top=2.0", width=3.5", height=2.0" (bounds: x=[1.0, 4.5], y=[2.0, 4.0])
    box_a = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0),
        Inches(2.0),
        Inches(3.5),
        Inches(2.0),
    )
    box_a.name = "Defect Box A (Alpha Underlay)"
    box_a.fill.solid()
    box_a.fill.fore_color.rgb = RGBColor(255, 204, 204)  # light red
    box_a.line.color.rgb = RGBColor(204, 0, 0)
    box_a.line.width = Pt(1.5)
    tf_a = box_a.text_frame
    tf_a.word_wrap = True
    p_a = tf_a.paragraphs[0]
    p_a.text = "Audit Box Alpha\n(Base element positioned at x=1.0, y=2.0)"
    p_a.font.name = "Calibri"
    p_a.font.size = Pt(13)
    p_a.font.color.rgb = RGBColor(150, 0, 0)

    # Box B: left=2.5", top=2.5", width=3.5", height=2.0" (bounds: x=[2.5, 6.0], y=[2.5, 4.5])
    # Overlap intersection: x in [2.5, 4.5] (width=2.0"), y in [2.5, 4.0] (height=1.5") -> Area = 3.0 sq inches (> 0.5 in)
    box_b = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5),
        Inches(2.5),
        Inches(3.5),
        Inches(2.0),
    )
    box_b.name = "Defect Box B (Beta Overlap)"
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = RGBColor(204, 229, 255)  # light blue
    box_b.line.color.rgb = RGBColor(0, 102, 204)
    box_b.line.width = Pt(1.5)
    tf_b = box_b.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = (
        "Audit Box Beta\n(Overlaps Box Alpha by 2.0 inches horizontal & 1.5 inches vertical)"
    )
    p_b.font.name = "Calibri"
    p_b.font.size = Pt(12)
    p_b.font.color.rgb = RGBColor(0, 50, 150)

    # 3. Defect 2: Shape clipped / extending beyond the right slide boundary (slide width is 13.333 inches)
    # Box C: left=11.5", top=2.0", width=3.0", height=2.0" (right edge = 14.5" -> extends 1.167" beyond canvas)
    box_c = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(11.5),
        Inches(2.0),
        Inches(3.0),
        Inches(2.0),
    )
    box_c.name = "Defect Box C (Right Edge Boundary Clipping)"
    box_c.fill.solid()
    box_c.fill.fore_color.rgb = RGBColor(255, 243, 205)  # warning yellow/amber
    box_c.line.color.rgb = RGBColor(255, 140, 0)
    box_c.line.width = Pt(1.5)
    tf_c = box_c.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "Clipped Boundary Shape\n(Extends 1.17 inches past 13.333\" slide width limit)"
    p_c.font.name = "Calibri"
    p_c.font.size = Pt(11)
    p_c.font.color.rgb = RGBColor(120, 70, 0)

    # 4. Defect 3: Shape with suspiciously tiny font (<8pt, e.g. 5.5pt)
    box_d = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(4.8),
        Inches(4.5),
        Inches(1.8),
    )
    box_d.name = "Defect Box D (Tiny Font Warning)"
    box_d.fill.solid()
    box_d.fill.fore_color.rgb = RGBColor(245, 245, 245)
    box_d.line.color.rgb = RGBColor(180, 180, 180)
    tf_d = box_d.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = Inches(0.2)
    tf_d.margin_top = Inches(0.15)
    p_d1 = tf_d.paragraphs[0]
    p_d1.text = "Tiny Font Warning Section:"
    p_d1.font.name = "Calibri"
    p_d1.font.size = Pt(10)
    p_d1.font.bold = True
    p_d1.font.color.rgb = RGBColor(50, 50, 50)
    p_d1.space_after = Pt(2)

    p_d2 = tf_d.add_paragraph()
    p_d2.text = (
        "LEGAL NOTICE: This microscopic disclaimer paragraph has intentionally been styled with "
        "a 5.5 point font size to rigorously trigger the rule-based slide validation engine's "
        "suspiciously tiny font warning (VAL-05 threshold < 8.0 pt). All terms herein are binding."
    )
    p_d2.font.name = "Calibri"
    p_d2.font.size = Pt(5.5)
    p_d2.font.color.rgb = RGBColor(80, 80, 80)
    for r in p_d2.runs:
        r.font.name = "Calibri"
        r.font.size = Pt(5.5)
        r.font.color.rgb = RGBColor(80, 80, 80)

    # 5. Defect 4: Text box with large text causing text overflow condition
    # Box E: 3.2" wide x 1.6" tall box filled with large 20pt text well exceeding capacity
    box_e = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.2),
        Inches(4.8),
        Inches(3.2),
        Inches(1.6),
    )
    box_e.name = "Defect Box E (Text Overflow Condition)"
    box_e.fill.solid()
    box_e.fill.fore_color.rgb = RGBColor(255, 235, 238)
    box_e.line.color.rgb = RGBColor(211, 47, 47)
    box_e.line.width = Pt(1.5)
    tf_e = box_e.text_frame
    tf_e.word_wrap = True
    tf_e.margin_left = Inches(0.15)
    tf_e.margin_top = Inches(0.15)
    tf_e.margin_right = Inches(0.15)
    tf_e.margin_bottom = Inches(0.15)

    p_e1 = tf_e.paragraphs[0]
    p_e1.text = "CRITICAL OVERFLOW:"
    p_e1.font.name = "Calibri"
    p_e1.font.size = Pt(20)
    p_e1.font.bold = True
    p_e1.font.color.rgb = RGBColor(180, 0, 0)

    p_e2 = tf_e.add_paragraph()
    p_e2.text = (
        "This long high-point sentence is guaranteed to exceed the physical container bounds "
        "of this 3.2x1.6 inch box creating an intentional text overflow defect."
    )
    p_e2.font.name = "Calibri"
    p_e2.font.size = Pt(18)
    p_e2.font.color.rgb = RGBColor(120, 0, 0)

    # Box F: A clean diagnostic companion box for comparison
    box_f = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.8),
        Inches(4.8),
        Inches(2.5),
        Inches(1.6),
    )
    box_f.name = "Diagnostic Companion Box"
    box_f.fill.solid()
    box_f.fill.fore_color.rgb = RGBColor(240, 248, 255)
    box_f.line.color.rgb = RGBColor(70, 130, 180)
    tf_f = box_f.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "Validation Target\n(Properly formatted standard box for baseline diffing)"
    p_f.font.name = "Calibri"
    p_f.font.size = Pt(11)
    p_f.font.color.rgb = RGBColor(30, 70, 120)

    return prs


def create_synthetic_deck(
    output_path: Optional[Union[str, Path]] = None, force: bool = True
) -> Path:
    """Generate and save the synthetic PowerPoint presentation to disk.

    Args:
        output_path: Target path for the .pptx file. Defaults to
                     tests/fixtures/synthetic_sample.pptx.
        force: If True, overwrites existing file. If False, only creates if missing.

    Returns:
        Path object pointing to the generated presentation.
    """
    if output_path is None:
        # Relative to project root
        base_dir = Path(__file__).resolve().parent
        output_path = base_dir / "synthetic_sample.pptx"
    else:
        output_path = Path(output_path).resolve()

    output_path.parent.mkdir(parents=True, exist_ok=True)

    if output_path.exists() and not force:
        return output_path

    prs = build_synthetic_deck()
    prs.save(str(output_path))
    return output_path


def main():
    """CLI entrypoint for generating synthetic test presentation."""
    target_path = None
    force = True

    if len(sys.argv) > 1:
        target_path = sys.argv[1]

    out_file = create_synthetic_deck(target_path, force=force)
    print(f"Successfully generated synthetic presentation at: {out_file}")


if __name__ == "__main__":
    main()

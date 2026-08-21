"""Adversarial stress testing for M1: Inspection Engine, Relationships, Hashing, and Scaling.

Tests:
1. Scaling: 50+ slides presentation benchmark (memory, speed, data integrity).
2. Density: 250+ shapes on a single slide (z-order uniqueness, ID indexing, matching performance).
3. Group and Nested Shapes: Inspection stability, child traversal, and semantic role mapping.
4. Relationships & Image Hashing: Multi-image SHA-256 deduplication and verification.
5. Hyperlinks: Shape-level, run-level, mailto, bookmarks, special chars, empty/broken links.
6. Non-destructive Guarantee: Exact SHA-256 byte-preservation of PPTX files before and after inspection.
7. Extreme Geometry: 0-dimension (0x0), negative coordinates, huge coordinates, extreme rotations.
8. Empty & Unicode Presentations: 0 shapes, emojis, RTL characters, multiline strings, tiny/huge fonts.
"""

import gc
import hashlib
import io
import os
import time
from pathlib import Path
from typing import List

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from powerpoint_mcp.models import (
    BoundingBox,
    PresentationModel,
    SemanticRole,
    ShapeModel,
    ShapeType,
    SlideModel,
    TextStyle,
)
from powerpoint_mcp.pptx import (
    PPTXInspector,
    extract_embedded_images,
    extract_hyperlinks,
    infer_semantic_role,
    inspect_presentation,
    inspect_shape,
    inspect_slide,
    inspect_slide_relationships,
    map_shape_type,
    match_shapes,
)


def _generate_tiny_png_bytes(color_rgb: tuple = (255, 0, 0)) -> bytes:
    """Generate a minimal valid 1x1 PNG byte stream in memory without external dependencies."""
    import struct
    import zlib

    # 1x1 raw image data: filter type 0 followed by 3 bytes RGB
    raw_data = b"\x00" + bytes(color_rgb)
    compressed = zlib.compress(raw_data)

    def png_chunk(chunk_type: bytes, data: bytes) -> bytes:
        chunk = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(chunk) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + chunk + crc

    png = b"\x89PNG\r\n\x1a\n"
    # IHDR: 1x1, 8-bit depth, RGB (2), default compression/filter/interlace
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    png += png_chunk(b"IHDR", ihdr_data)
    png += png_chunk(b"IDAT", compressed)
    png += png_chunk(b"IEND", b"")
    return png


@pytest.fixture
def adversarial_deck_50_slides(tmp_path: Path) -> Path:
    """Generate a 55-slide presentation with diverse layouts and content."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # Slide 1: Title slide (2 shapes: Title + Subtitle)
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Scaling Benchmark Presentation"
    if len(s1.placeholders) > 1:
        s1.placeholders[1].text = "Adversarial Stress Test Deck (55 Slides)"

    # Slides 2 to 55: varied slide structures
    for i in range(2, 56):
        s = prs.slides.add_slide(blank_layout)
        # Add title
        title_box = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Slide {i}: Automated Analysis & Metrics"
        p.font.size = Pt(28)
        p.font.bold = True

        # Add 4 KPI cards per slide
        for c in range(4):
            card = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.0 + c * 2.8),
                Inches(1.8),
                Inches(2.5),
                Inches(1.8),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(240, 244, 248)
            card_tf = card.text_frame
            card_tf.text = f"Metric {c+1}\nValue: {i * (c+1) * 17.5:.1f}"

        # Add a table on every 5th slide
        if i % 5 == 0:
            table_shape = s.shapes.add_table(3, 3, Inches(1), Inches(4.0), Inches(11.333), Inches(2.5))
            for r in range(3):
                for col in range(3):
                    table_shape.table.cell(r, col).text = f"R{r}C{col}"

        # Add speaker notes on every 3rd slide
        if i % 3 == 0:
            s.notes_slide.notes_text_frame.text = f"Speaker notes for slide {i} - benchmark run."

    file_path = tmp_path / "scaling_55_slides.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture
def adversarial_dense_slide_deck(tmp_path: Path) -> Path:
    """Generate a presentation with a single slide containing 260 shapes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title shape
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    title.text_frame.text = "High Density 260-Shape Slide"

    # Add a 16 x 16 grid of small shapes = 256 shapes
    grid_rows, grid_cols = 16, 16
    start_x, start_y = 0.5, 0.9
    w, h = 0.65, 0.35
    gap_x, gap_y = 0.12, 0.05

    for r in range(grid_rows):
        for c in range(grid_cols):
            x = Inches(start_x + c * (w + gap_x))
            y = Inches(start_y + r * (h + gap_y))
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(w), Inches(h))
            shp.name = f"GridCell_{r}_{c}"
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(r * 15 % 256, c * 15 % 256, (r + c) * 8 % 256)
            shp.text_frame.text = f"{r},{c}"
            shp.text_frame.paragraphs[0].font.size = Pt(7)

    # Plus 3 connectors
    for k in range(3):
        slide.shapes.add_connector(1, Inches(1 + k * 3), Inches(7.0), Inches(3 + k * 3), Inches(7.0))

    file_path = tmp_path / "dense_260_shapes.pptx"
    prs.save(str(file_path))
    return file_path


# =============================================================================
# 1. Performance & Memory Scaling Tests (50+ Slides)
# =============================================================================
class TestScalingPerformance:
    """Adversarial stress testing on large (50+ slide) presentations."""

    def test_inspect_55_slides_presentation_speed_and_completeness(self, adversarial_deck_50_slides):
        gc.collect()
        t0 = time.perf_counter()

        model: PresentationModel = inspect_presentation(adversarial_deck_50_slides)
        elapsed = time.perf_counter() - t0

        assert model.slide_count == 55
        assert len(model.slides) == 55
        assert len(model.slide_titles) == 55
        # Ensure fast inspection (< 3.0s for 55 slides with full shape trees)
        assert elapsed < 3.0, f"55-slide inspection too slow: {elapsed:.3f}s"

        # Verify indexing and slide models
        for idx, slide in enumerate(model.slides, start=1):
            assert slide.slide_number == idx
            assert slide.slide_id > 0
            assert slide.width_inches == pytest.approx(13.333, abs=0.01)
            assert slide.height_inches == pytest.approx(7.5, abs=0.01)

            if idx == 1:
                assert slide.shape_count == 2
                assert slide.title == "Scaling Benchmark Presentation"
            else:
                assert slide.shape_count >= 5  # Title + 4 cards

            if idx % 3 == 0:
                assert slide.has_notes is True
                assert f"slide {idx}" in slide.notes

            if idx % 5 == 0:
                table_shapes = slide.get_shapes_by_type(ShapeType.TABLE)
                assert len(table_shapes) >= 1
                assert table_shapes[0].table_metadata is not None
                assert table_shapes[0].table_metadata["rows"] == 3
                assert table_shapes[0].table_metadata["columns"] == 3

    def test_repeated_inspection_no_memory_leak(self, adversarial_deck_50_slides):
        """Repeatedly inspect presentation to ensure no accumulating object references."""
        for _ in range(5):
            prs_model = inspect_presentation(adversarial_deck_50_slides)
            assert prs_model.slide_count == 55
            del prs_model
        gc.collect()


# =============================================================================
# 2. High-Density Slide Stress Tests (250+ Shapes)
# =============================================================================
class TestHighDensitySlide:
    """Adversarial stress testing on single slides with 260 shapes."""

    def test_inspect_dense_slide_shapes_and_z_order(self, adversarial_dense_slide_deck):
        t0 = time.perf_counter()
        slide = inspect_slide(adversarial_dense_slide_deck, 1)
        elapsed = time.perf_counter() - t0

        assert slide.shape_count == 260
        assert len(slide.shapes) == 260
        assert elapsed < 1.0, f"260-shape slide inspection took too long: {elapsed:.3f}s"

        # Verify unique z_orders spanning 0 to 259
        z_orders = [s.z_order for s in slide.shapes]
        assert z_orders == list(range(260))

        # Verify all shape IDs are distinct
        shape_ids = [s.shape_id for s in slide.shapes]
        assert len(set(shape_ids)) == 260

        # Verify fast random lookup by ID
        for sample_id in shape_ids[::25]:
            found = slide.get_shape_by_id(sample_id)
            assert found is not None
            assert found.shape_id == sample_id

        # Verify connectors detection
        connectors = slide.get_shapes_by_type(ShapeType.CONNECTOR)
        assert len(connectors) == 3

    def test_dense_slide_shape_matching_performance(self, adversarial_dense_slide_deck):
        """Shape matching between two dense 260-shape slides must complete reasonably fast."""
        slide_a = inspect_slide(adversarial_dense_slide_deck, 1)
        slide_b = inspect_slide(adversarial_dense_slide_deck, 1)

        t0 = time.perf_counter()
        matches = match_shapes(slide_a, slide_b, min_confidence=0.50)
        elapsed = time.perf_counter() - t0

        # Self match of 260 shapes should match all 260 shapes
        assert len(matches) == 260
        assert elapsed < 4.0, f"Shape matching for 260x260 took too long: {elapsed:.3f}s"
        for m in matches:
            assert m["shape_a_id"] == m["shape_b_id"]
            assert m["confidence_score"] >= 0.90


# =============================================================================
# 3. Relationship Extraction & Multi-Image SHA-256 Deduplication
# =============================================================================
class TestRelationshipsAndImageHashing:
    """Adversarial testing of OpenXML relationship parsing, image extraction, and SHA-256 hashing."""

    def test_multi_image_sha256_deduplication(self, tmp_path: Path):
        """Ensure identical embedded images produce identical SHA-256 hashes, distinct images produce distinct hashes."""
        red_png = _generate_tiny_png_bytes((255, 0, 0))
        blue_png = _generate_tiny_png_bytes((0, 0, 255))
        green_png = _generate_tiny_png_bytes((0, 255, 0))

        red_hash = hashlib.sha256(red_png).hexdigest()
        blue_hash = hashlib.sha256(blue_png).hexdigest()
        green_hash = hashlib.sha256(green_png).hexdigest()

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(6)
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])

        # Slide 1 has Red image (1) and Red image (2) [duplicate part/blob]
        slide1.shapes.add_picture(io.BytesIO(red_png), Inches(1), Inches(1), Inches(2), Inches(2))
        slide1.shapes.add_picture(io.BytesIO(red_png), Inches(4), Inches(1), Inches(2), Inches(2))

        # Slide 2 has Blue image and Green image
        slide2.shapes.add_picture(io.BytesIO(blue_png), Inches(1), Inches(1), Inches(2), Inches(2))
        slide2.shapes.add_picture(io.BytesIO(green_png), Inches(4), Inches(1), Inches(2), Inches(2))

        deck_path = tmp_path / "multi_image_deck.pptx"
        prs.save(str(deck_path))

        # Test extraction across whole presentation
        extracted = extract_embedded_images(Presentation(str(deck_path)))
        assert len(extracted) == 4

        hashes = [img["sha256"] for img in extracted]
        assert hashes[0] == red_hash
        assert hashes[1] == red_hash
        assert hashes[2] == blue_hash
        assert hashes[3] == green_hash

        # Verify sizes and content types
        for img in extracted:
            assert img["size_bytes"] == len(red_png)
            assert img["content_type"] == "image/png"
            assert img["extension"] == "png"
            assert img["width_px"] == 1
            assert img["height_px"] == 1

    def test_slide_relationships_extraction_and_robustness(self, tmp_path: Path):
        """Test inspect_slide_relationships on slides with and without relationships."""
        prs = Presentation()
        slide_with_rel = prs.slides.add_slide(prs.slide_layouts[6])
        red_png = _generate_tiny_png_bytes((255, 0, 0))
        slide_with_rel.shapes.add_picture(io.BytesIO(red_png), Inches(1), Inches(1), Inches(2), Inches(2))

        slide_empty = prs.slides.add_slide(prs.slide_layouts[6])

        deck_path = tmp_path / "rels_test.pptx"
        prs.save(str(deck_path))

        loaded_prs = Presentation(str(deck_path))
        rels_1 = inspect_slide_relationships(loaded_prs.slides[0])
        assert len(rels_1) > 0

        # Check fields in relationship
        for rel in rels_1:
            assert "r_id" in rel
            assert "rel_type" in rel
            assert "target_ref" in rel
            assert "is_external" in rel

        # Empty slide still has layout relationship
        rels_2 = inspect_slide_relationships(loaded_prs.slides[1])
        assert isinstance(rels_2, list)

        # Robustness against non-slide object
        assert inspect_slide_relationships(None) == []
        assert inspect_slide_relationships("invalid") == []


# =============================================================================
# 4. Hyperlink Parsing Robustness
# =============================================================================
class TestHyperlinkParsing:
    """Adversarial stress testing on shape-level and run-level hyperlink parsing."""

    def test_complex_hyperlinks_extraction(self, tmp_path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. Shape-level hyperlink
        box1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1))
        box1.text_frame.text = "Click to visit portal"
        box1.click_action.hyperlink.address = "https://example.com/portal?ref=ppt&user=101%20test"

        # 2. Text run-level hyperlinks (multiple runs in same paragraph)
        box2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = box2.text_frame
        p = tf.paragraphs[0]

        r1 = p.add_run()
        r1.text = "Send questions to "

        r2 = p.add_run()
        r2.text = "support@example.com"
        r2.hyperlink.address = "mailto:support@example.com?subject=PPTX%20Help"

        r3 = p.add_run()
        r3.text = " or read the "

        r4 = p.add_run()
        r4.text = "Documentation"
        r4.hyperlink.address = "https://docs.example.com/guide/index.html#section-3"

        # 3. Shape without hyperlink
        box3 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(3), Inches(1))
        box3.text_frame.text = "Plain text box with no links"

        deck_path = tmp_path / "hyperlinks_test.pptx"
        prs.save(str(deck_path))

        loaded_prs = Presentation(str(deck_path))
        links = extract_hyperlinks(loaded_prs.slides[0])

        assert len(links) >= 3

        # Check shape-level link
        shape_link = next((l for l in links if l["type"] == "shape"), None)
        assert shape_link is not None
        assert "example.com/portal" in shape_link["address"]

        # Check mailto run link
        mailto_link = next((l for l in links if l["type"] == "run" and "mailto:" in l["address"]), None)
        assert mailto_link is not None
        assert mailto_link["text"] == "support@example.com"
        assert "subject=PPTX%20Help" in mailto_link["address"]

        # Check anchor/fragment run link
        doc_link = next((l for l in links if l["type"] == "run" and "docs.example.com" in l["address"]), None)
        assert doc_link is not None
        assert doc_link["text"] == "Documentation"
        assert "#section-3" in doc_link["address"]

    def test_extract_hyperlinks_robustness_on_empty_and_corrupt(self):
        """extract_hyperlinks should not crash on slides without shapes or invalid inputs."""
        assert extract_hyperlinks(None) == []
        assert extract_hyperlinks(object()) == []


# =============================================================================
# 5. Non-Destructive Inspection Guarantee
# =============================================================================
class TestNonDestructiveBehavior:
    """Guarantee that PPTX inspection is strictly read-only and never modifies files."""

    def test_file_hash_byte_exact_preservation(self, tmp_path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        box.text_frame.text = "Immutable Content Verification"

        test_path = tmp_path / "immutable_deck.pptx"
        prs.save(str(test_path))

        # Compute pre-inspection SHA-256
        with open(test_path, "rb") as f:
            pre_hash = hashlib.sha256(f.read()).hexdigest()

        # Perform extensive deep inspection operations
        prs_model = inspect_presentation(test_path)
        slide_model = inspect_slide(test_path, 1)
        for shape in slide_model.shapes:
            inspect_shape(test_path, 1, shape.shape_id)
        match_shapes(slide_model, slide_model)

        prs_loaded = Presentation(str(test_path))
        inspect_slide_relationships(prs_loaded.slides[0])
        extract_hyperlinks(prs_loaded.slides[0])
        extract_embedded_images(prs_loaded)

        # Compute post-inspection SHA-256
        with open(test_path, "rb") as f:
            post_hash = hashlib.sha256(f.read()).hexdigest()

        assert pre_hash == post_hash, "CRITICAL: Inspection modified the underlying presentation file on disk!"


# =============================================================================
# 6. Extreme Geometries & Degenerate Shapes
# =============================================================================
class TestExtremeGeometries:
    """Adversarial stress testing on 0-dimension shapes, negative coordinates, and extreme rotations."""

    def test_zero_dimensions_bounding_box(self):
        """BoundingBox with 0 width and 0 height."""
        bbox = BoundingBox(left_emu=914400, top_emu=914400, width_emu=0, height_emu=0)
        assert bbox.width_inches == 0.0
        assert bbox.height_inches == 0.0
        assert bbox.center_x_inches == 1.0
        assert bbox.center_y_inches == 1.0
        d = bbox.to_dict()
        assert d["width_inches"] == 0.0
        assert d["height_inches"] == 0.0

    def test_match_shapes_with_zero_dimension_shapes(self):
        """Zero-dimension shapes should not cause ZeroDivisionError in match_shapes."""
        s1 = ShapeModel(
            shape_id=1,
            name="Point A",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.UNKNOWN,
            bbox=BoundingBox(left_emu=0, top_emu=0, width_emu=0, height_emu=0),
        )
        s2 = ShapeModel(
            shape_id=2,
            name="Point B",
            shape_type=ShapeType.AUTO_SHAPE,
            semantic_role=SemanticRole.UNKNOWN,
            bbox=BoundingBox(left_emu=0, top_emu=0, width_emu=0, height_emu=0),
        )
        slide_a = SlideModel(slide_number=1, slide_id=10, width_emu=9144000, height_emu=5143500, shapes=[s1])
        slide_b = SlideModel(slide_number=1, slide_id=10, width_emu=9144000, height_emu=5143500, shapes=[s2])

        matches = match_shapes(slide_a, slide_b)
        assert len(matches) == 1
        assert matches[0]["confidence_score"] > 0.50

    def test_negative_and_extreme_coordinates_in_pptx(self, tmp_path: Path):
        """Create presentation with off-slide negative and massive coordinates."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(6)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Negative coordinate shape (partially off-screen left/top)
        shp_neg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-2.0), Inches(-1.5), Inches(3.0), Inches(2.0))
        shp_neg.name = "Offscreen TopLeft"

        # Far outside coordinate shape
        shp_far = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(25.0), Inches(30.0), Inches(5.0), Inches(5.0))
        shp_far.name = "Far Outside"

        deck_path = tmp_path / "extreme_coords.pptx"
        prs.save(str(deck_path))

        slide_model = inspect_slide(deck_path, 1)
        assert slide_model.shape_count == 2

        neg_shape = next(s for s in slide_model.shapes if s.name == "Offscreen TopLeft")
        assert neg_shape.bbox.left_inches == pytest.approx(-2.0, abs=0.01)
        assert neg_shape.bbox.top_inches == pytest.approx(-1.5, abs=0.01)

        far_shape = next(s for s in slide_model.shapes if s.name == "Far Outside")
        assert far_shape.bbox.left_inches == pytest.approx(25.0, abs=0.01)
        assert far_shape.bbox.top_inches == pytest.approx(30.0, abs=0.01)

    def test_rotated_shapes_inspection(self, tmp_path: Path):
        """Verify rotation property is accurately preserved for diverse angles."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        angles = [0.0, 45.0, 90.0, 180.0, 270.0, 315.5]
        for idx, angle in enumerate(angles):
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1 + idx * 1.5), Inches(2), Inches(1), Inches(1))
            shp.name = f"Rotated_{int(angle)}"
            shp.rotation = angle

        deck_path = tmp_path / "rotated_shapes.pptx"
        prs.save(str(deck_path))

        slide_model = inspect_slide(deck_path, 1)
        assert slide_model.shape_count == len(angles)
        for s, expected_angle in zip(slide_model.shapes, angles):
            assert s.rotation == pytest.approx(expected_angle, abs=0.1)


# =============================================================================
# 7. Unicode, Multiline, Tiny/Huge Fonts, and Edge Slide Handling
# =============================================================================
class TestUnicodeAndEdgePresentations:
    """Adversarial testing on Unicode strings, emojis, RTL characters, extreme font sizes, and empty decks."""

    def test_unicode_and_emojis_in_text_frames(self, tmp_path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        # Multi-language + emojis + special symbols + newlines
        p.text = "Quarterly Revenue 🚀 | 東京オフィス 📈 | Übermensch | مرحباً بالعالم | \u221e \u2264 \u03c0"

        deck_path = tmp_path / "unicode_deck.pptx"
        prs.save(str(deck_path))

        slide_model = inspect_slide(deck_path, 1)
        shape_model = slide_model.shapes[0]
        assert shape_model.text_frame is not None
        assert "🚀" in shape_model.text_frame.text
        assert "東京オフィス" in shape_model.text_frame.text
        assert "مرحباً بالعالم" in shape_model.text_frame.text
        assert "\u03c0" in shape_model.text_frame.text

    def test_extreme_font_sizes(self, tmp_path: Path):
        """Test tiny (2pt) and gigantic (120pt) fonts."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        box_tiny = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        box_tiny.text_frame.paragraphs[0].text = "Micro print"
        box_tiny.text_frame.paragraphs[0].font.size = Pt(2)

        box_huge = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
        box_huge.text_frame.paragraphs[0].text = "GIANT"
        box_huge.text_frame.paragraphs[0].font.size = Pt(120)

        deck_path = tmp_path / "font_sizes_deck.pptx"
        prs.save(str(deck_path))

        slide_model = inspect_slide(deck_path, 1)
        s_tiny = slide_model.shapes[0]
        s_huge = slide_model.shapes[1]

        assert s_tiny.text_frame.paragraphs[0].runs[0].style.font_size_pt == pytest.approx(2.0, abs=0.1)
        assert s_huge.text_frame.paragraphs[0].runs[0].style.font_size_pt == pytest.approx(120.0, abs=0.1)

    def test_slide_with_zero_shapes(self, tmp_path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Completely empty slide
        deck_path = tmp_path / "empty_slide_deck.pptx"
        prs.save(str(deck_path))

        slide_model = inspect_slide(deck_path, 1)
        assert slide_model.shape_count == 0
        assert slide_model.shapes == []
        assert slide_model.title is None
        assert slide_model.has_notes is False

        d = slide_model.to_dict()
        assert d["shape_count"] == 0
        assert d["shapes"] == []

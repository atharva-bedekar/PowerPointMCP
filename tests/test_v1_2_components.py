"""Comprehensive tests for PowerPoint MCP v1.2: Semantic Components and Cross-Slide Consistency."""

import os
from pathlib import Path
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.models.component import ComponentModel, ComponentType
from powerpoint_mcp.pptx.components import detect_slide_components, inspect_components
from powerpoint_mcp.pptx.cross_slide import compare_cross_slides
from powerpoint_mcp.tools.editing import (
    ppt_create_stepper,
    ppt_create_structured_card_list,
    ppt_move_component,
    ppt_resize_component,
    ppt_sync_component,
    ppt_sync_layout,
    ppt_sync_slide_chrome,
    ppt_update_stepper,
)
from powerpoint_mcp.tools.inspection import (
    ppt_compare_slides,
    ppt_inspect_components,
    ppt_validate_slide,
)
from powerpoint_mcp.tools.rendering import ppt_render_slides
from powerpoint_mcp.tools.versioning import get_session_manager, open_presentation, save_session


@pytest.fixture(autouse=True)
def reset_session_state():
    """Ensure no global session leaks between tests."""
    sm = get_session_manager()
    sm.close_session()
    yield
    sm.close_session()


@pytest.fixture
def sample_flow_deck(tmp_path):
    """Create a 4-slide test presentation representing a flow (Analyze, Connect, Configure, Run)."""
    deck_path = tmp_path / "flow_deck.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    steps = ["ANALYZE", "CONNECT", "CONFIGURE", "RUN"]

    for i in range(4):
        slide = prs.slides.add_slide(blank_layout)

        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.6))
        title_box.name = f"Slide Title {i+1}"
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Phase {i+1}: {steps[i]}"
        p.runs[0].font.size = Pt(24)
        p.runs[0].font.bold = True

        # Stepper
        for j, step in enumerate(steps):
            bx = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(0.8 + j * 2.8),
                Inches(1.2),
                Inches(2.4),
                Inches(0.4),
            )
            bx.name = f"Step Node {j+1} - {step}"
            bx.text_frame.text = step

        # Content box
        content = slide.shapes.add_shape(
            1,
            Inches(0.8),
            Inches(2.0),
            Inches(11.7),
            Inches(4.5),
        )
        content.name = f"Content Card {i+1}"
        content.text_frame.text = f"Substantive Content for {steps[i]} with unique details."

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(5.0), Inches(0.4))
        footer.name = f"Footer {i+1}"
        footer.text_frame.text = "Confidential - Acme Corp Flow Deck"

    prs.save(str(deck_path))
    return str(deck_path)


def test_detect_components_header_footer_stepper(sample_flow_deck):
    """Test that component detection correctly identifies header, footer, stepper, and content area."""
    comps = detect_slide_components(sample_flow_deck, slide_number=1)
    comp_types = [c.type_str for c in comps]

    assert "header" in comp_types
    assert "footer" in comp_types
    assert "stepper" in comp_types
    assert "content_area" in comp_types or "card" in comp_types

    stepper_comp = next(c for c in comps if c.type_str == "stepper")
    assert len(stepper_comp.properties.get("steps", [])) == 4
    assert stepper_comp.properties["steps"] == ["ANALYZE", "CONNECT", "CONFIGURE", "RUN"]


def test_ppt_inspect_components_tool(sample_flow_deck):
    """Test ppt_inspect_components MCP tool endpoint."""
    res = ppt_inspect_components(slide_number=1, presentation_path=sample_flow_deck)
    assert res["success"] is True
    assert res["component_count"] >= 3
    comp_ids = [c["id"] for c in res["components"]]
    assert "header" in comp_ids
    assert "footer" in comp_ids
    assert "stepper" in comp_ids


def test_ppt_compare_slides_multi_target(sample_flow_deck):
    """Test cross-slide comparison of multiple targets against a reference slide."""
    res = ppt_compare_slides(
        reference_slide=1,
        target_slides=[2, 3, 4],
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True
    assert res["reference_slide"] == 1
    assert "summary" in res
    assert "REFERENCE: Slide 1" in res["summary"]
    assert "HEADER" in res["summary"]
    assert "STEPPER" in res["summary"]


def test_ppt_compare_slides_legacy_backward_compatibility(sample_flow_deck):
    """Test legacy 2-slide comparison using slide_a and slide_b."""
    res = ppt_compare_slides(
        slide_a=1,
        slide_b=2,
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True
    assert "overall_similarity_score" in res or "geometric_match_score" in res


def test_ppt_create_and_update_stepper_no_orphans(tmp_path):
    """Test stepper creation and active step updating, verifying complete cleanup of old shapes."""
    deck_path = tmp_path / "stepper_test.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.6))
    t.name = "Slide Title"
    t.text_frame.text = "Process Overview"
    c = slide.shapes.add_shape(1, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
    c.name = "Content Card"
    prs.save(str(deck_path))

    # Create clean stepper on slide 1
    res_create = ppt_create_stepper(
        slide_number=1,
        steps=["ANALYZE", "CONNECT", "CONFIGURE", "RUN"],
        active_step="ANALYZE",
        presentation_path=str(deck_path),
    )
    assert res_create["success"] is True

    prs1 = Presentation(str(deck_path))
    slide1_shapes_middle = [s.shape_id for s in prs1.slides[0].shapes]
    assert len(slide1_shapes_middle) == 2 + 7  # 2 base shapes + 4 steps + 3 arrows = 9 shapes

    # Update active step to CONNECT
    res_update = ppt_update_stepper(
        slide_number=1,
        active_step="CONNECT",
        presentation_path=str(deck_path),
    )
    assert res_update["success"] is True

    prs2 = Presentation(str(deck_path))
    slide1_shapes_after = list(prs2.slides[0].shapes)
    # Total shape count should remain exactly 9 (old stepper replaced cleanly, no orphaned shapes)
    assert len(slide1_shapes_after) == 9

    # Verify active step was updated to CONNECT
    stepper_names = [s.name for s in slide1_shapes_after]
    assert any("CONNECT" in name and "Active" in name for name in stepper_names)


def test_ppt_sync_slide_chrome(sample_flow_deck):
    """Test synchronizing shared slide chrome across slides 2, 3, 4 from reference slide 1."""
    res = ppt_sync_slide_chrome(
        reference_slide=1,
        target_slides=[2, 3, 4],
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True
    assert res["target_slides"] == [2, 3, 4]
    assert "header" in res["synchronized_components"]
    assert "footer" in res["synchronized_components"]


def test_ppt_sync_layout_preserves_content(sample_flow_deck):
    """Test synchronizing content layout while preserving target text."""
    res = ppt_sync_layout(
        reference_slide=1,
        target_slides=[2, 3],
        component="content_area",
        preserve_content=True,
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True

    # Verify target slide text was preserved
    prs = Presentation(sample_flow_deck)
    slide2_text = " ".join(s.text_frame.text for s in prs.slides[1].shapes if s.has_text_frame)
    assert "Substantive Content for CONNECT" in slide2_text


def test_ppt_create_structured_card_list(sample_flow_deck):
    """Test creating structured cards with row items and dividers."""
    res = ppt_create_structured_card_list(
        slide_number=3,
        container_bbox={"left": 2.0, "top": 2.0, "width": 8.0, "height": 3.5},
        items=[
            {"title": "Card Row 1", "description": "Details about component 1"},
            {"title": "Card Row 2", "description": "Details about component 2"},
            {"title": "Card Row 3", "description": "Details about component 3"},
        ],
        divider=True,
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True
    assert res["item_count"] == 3
    assert len(res["divider_shape_ids"]) == 2  # 2 dividers for 3 items


def test_ppt_move_and_resize_component(sample_flow_deck):
    """Test moving and resizing components atomically."""
    # Move stepper
    res_move = ppt_move_component(
        slide_number=1,
        component_id="stepper",
        dx=0.2,
        dy=0.1,
        presentation_path=sample_flow_deck,
    )
    assert res_move["success"] is True
    assert len(res_move["shape_ids"]) >= 2

    # Resize content area
    res_resize = ppt_resize_component(
        slide_number=1,
        component_id="content_area",
        scale_width=1.05,
        presentation_path=sample_flow_deck,
    )
    assert res_resize["success"] is True
    assert len(res_resize["shape_ids"]) >= 1


def test_ppt_render_slides_batch(sample_flow_deck):
    """Test batch slide rendering tool."""
    res = ppt_render_slides(
        slide_numbers=[1, 2, 3, 4],
        renderer="mock",
        presentation_path=sample_flow_deck,
    )
    assert res["success"] is True
    assert res["slide_count"] == 4
    assert len(res["rendered_slides"]) == 4
    for slide_info in res["rendered_slides"]:
        assert Path(slide_info["image_path"]).exists()


def test_session_lifecycle_with_components(sample_flow_deck):
    """Test complete isolated session lifecycle with component operations and persistence."""
    session = open_presentation(sample_flow_deck)
    assert session is not None

    # Update stepper in active session
    res_stepper = ppt_update_stepper(
        slide_number=2,
        active_step="CONNECT",
    )
    assert res_stepper["success"] is True
    assert res_stepper["target"] == "working"

    # Save session
    save_res = save_session()
    assert save_res["success"] is True

    # Verify change in saved original
    prs = Presentation(sample_flow_deck)
    comps = detect_slide_components(prs.slides[1], slide_number=2)
    stepper = next(c for c in comps if c.type_str == "stepper")
    assert stepper.properties.get("active_step") == "CONNECT"

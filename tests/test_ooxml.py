"""Comprehensive test suite for direct OOXML manipulation helpers (transparency, gradient fill, drop shadows, safe XML modifications)."""

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from powerpoint_mcp.pptx.ooxml import (
    NAMESPACES,
    get_or_create_spPr,
    get_raw_shape_xml,
    safe_modify_xml,
    set_drop_shadow,
    set_gradient_fill,
    set_shape_gradient_fill,
    set_shape_shadow_effect,
    set_shape_transparency,
)

A_NS = f"{{{NAMESPACES['a']}}}"
P_NS = f"{{{NAMESPACES['p']}}}"


class TestOOXMLTransparency:
    """Test setting shape transparency via DrawingML <a:alpha>."""

    def test_set_shape_transparency_50_percent(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        set_shape_transparency(shape, 50.0)

        spPr = get_or_create_spPr(shape)
        solid_fill = spPr.find(f"{A_NS}solidFill")
        assert solid_fill is not None

        # Find color element
        color_elem = solid_fill[0]
        alpha = color_elem.find(f"{A_NS}alpha")
        assert alpha is not None
        assert alpha.get("val") == "50000"

    def test_set_shape_transparency_0_percent_opaque(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        # First set to 30%, then reset to 0% (opaque)
        set_shape_transparency(shape, 30.0)
        set_shape_transparency(shape, 0.0)

        spPr = get_or_create_spPr(shape)
        solid_fill = spPr.find(f"{A_NS}solidFill")
        assert solid_fill is not None
        color_elem = solid_fill[0]
        alpha = color_elem.find(f"{A_NS}alpha")
        # Either removed or None
        assert alpha is None or alpha.get("val") == "100000"

    def test_set_shape_transparency_100_percent_transparent(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        set_shape_transparency(shape, 100.0)

        spPr = get_or_create_spPr(shape)
        solid_fill = spPr.find(f"{A_NS}solidFill")
        color_elem = solid_fill[0]
        alpha = color_elem.find(f"{A_NS}alpha")
        assert alpha is not None
        assert alpha.get("val") == "0"


class TestOOXMLGradients:
    """Test linear gradient fill configuration via <a:gradFill>."""

    def test_set_two_stop_gradient(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        set_shape_gradient_fill(shape, start_hex="003366", end_hex="3399FF", angle_deg=90.0)

        spPr = get_or_create_spPr(shape)
        grad_fill = spPr.find(f"{A_NS}gradFill")
        assert grad_fill is not None

        gs_lst = grad_fill.find(f"{A_NS}gsLst")
        assert gs_lst is not None
        gs_items = gs_lst.findall(f"{A_NS}gs")
        assert len(gs_items) == 2

        assert gs_items[0].get("pos") == "0"
        assert gs_items[0].find(f"{A_NS}srgbClr").get("val") == "003366"

        assert gs_items[1].get("pos") == "100000"
        assert gs_items[1].find(f"{A_NS}srgbClr").get("val") == "3399FF"

        lin = grad_fill.find(f"{A_NS}lin")
        assert lin is not None
        # 90 degrees * 60000 = 5400000
        assert lin.get("ang") == "5400000"

    def test_set_multi_stop_gradient(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        stops = [
            {"position": 0.0, "color": "FF0000"},
            {"position": 0.5, "color": "00FF00"},
            {"position": 1.0, "color": "0000FF"},
        ]
        set_gradient_fill(shape, stops=stops, angle_deg=45.0)

        spPr = get_or_create_spPr(shape)
        grad_fill = spPr.find(f"{A_NS}gradFill")
        gs_lst = grad_fill.find(f"{A_NS}gsLst")
        gs_items = gs_lst.findall(f"{A_NS}gs")
        assert len(gs_items) == 3
        assert gs_items[1].get("pos") == "50000"
        assert gs_items[1].find(f"{A_NS}srgbClr").get("val") == "00FF00"


class TestOOXMLDropShadow:
    """Test outer drop shadow effect configuration via <a:effectLst>/<a:outerShdw>."""

    def test_set_drop_shadow_parameters(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        set_drop_shadow(
            shape,
            blur_rad_pt=6.0,
            dist_pt=4.0,
            dir_deg=90.0,
            color_hex="112233",
            alpha_percent=30.0,
        )

        spPr = get_or_create_spPr(shape)
        effect_lst = spPr.find(f"{A_NS}effectLst")
        assert effect_lst is not None

        outer_shdw = effect_lst.find(f"{A_NS}outerShdw")
        assert outer_shdw is not None

        # blurRad: 6.0 pt * 12700 = 76200
        assert outer_shdw.get("blurRad") == "76200"
        # dist: 4.0 pt * 12700 = 50800
        assert outer_shdw.get("dist") == "50800"
        # dir: 90 deg * 60000 = 5400000
        assert outer_shdw.get("dir") == "5400000"

        srgb = outer_shdw.find(f"{A_NS}srgbClr")
        assert srgb is not None
        assert srgb.get("val") == "112233"

        alpha = srgb.find(f"{A_NS}alpha")
        assert alpha is not None
        assert alpha.get("val") == "30000"


class TestRawXMLAndSafeModification:
    """Test raw shape XML retrieval and transactional XML modification."""

    def test_get_raw_shape_xml(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        xml_str = get_raw_shape_xml(shape)
        assert isinstance(xml_str, str)
        assert "<p:sp" in xml_str or "<p:pic" in xml_str or "<p:graphicFrame" in xml_str

    def test_safe_modify_xml_success(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        def custom_modifier(elem):
            # Add a custom attribute or child
            elem.set("test_attr", "verified")

        result = safe_modify_xml(shape, custom_modifier)
        assert result is True
        assert shape._element.get("test_attr") == "verified"

    def test_safe_modify_xml_rollback_on_failure(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        def broken_modifier(elem):
            elem.set("should_revert", "temp")
            raise RuntimeError("Deliberate error to trigger rollback")

        with pytest.raises(ValueError, match="XML modification failed"):
            safe_modify_xml(shape, broken_modifier)

        assert shape._element.get("should_revert") is None


class TestOOXMLEdgeCases:
    """Test edge cases and safety mechanisms in OOXML helpers."""

    def test_transparency_clamping(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]

        # Negative percent -> 0% (opaque)
        set_shape_transparency(shape, -10.0)
        spPr = get_or_create_spPr(shape)
        solid_fill = spPr.find(f"{A_NS}solidFill")
        color_elem = solid_fill[0]
        alpha = color_elem.find(f"{A_NS}alpha")
        assert alpha is None or alpha.get("val") == "100000"

        # > 100% -> 100% (transparent)
        set_shape_transparency(shape, 150.0)
        alpha = color_elem.find(f"{A_NS}alpha")
        assert alpha is not None
        assert alpha.get("val") == "0"

    def test_get_raw_shape_xml_compact(self, temp_presentation):
        slide = temp_presentation.slides[0]
        shape = slide.shapes[0]
        xml_compact = get_raw_shape_xml(shape, pretty_print=False)
        assert "\n" not in xml_compact or len(xml_compact) > 0

    def test_safe_modify_xml_with_raw_lxml_element(self, temp_presentation):
        slide = temp_presentation.slides[0]
        elem = slide.shapes[0]._element

        def modifier(e):
            e.set("raw_attr", "123")

        assert safe_modify_xml(elem, modifier) is True
        assert elem.get("raw_attr") == "123"


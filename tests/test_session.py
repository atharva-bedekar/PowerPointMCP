"""Comprehensive test suite for session management, working-copy tracking, backups, and revert."""

from datetime import datetime
import json
import os
from pathlib import Path
import shutil
import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from powerpoint_mcp.pptx.editor import modify_shape, modify_text
from powerpoint_mcp.pptx.inspector import inspect_slide
from powerpoint_mcp.tools.versioning import (
    Session,
    SessionManager,
    create_backup,
    get_current_session,
    get_session,
    get_session_manager,
    open_presentation,
    revert_session,
    save_as,
    save_session,
)
from powerpoint_mcp.utils.paths import (
    cleanup_old_sessions,
    cleanup_session,
    ensure_session_dirs,
    generate_backup_filename,
    get_session_dir,
    get_session_metadata_path,
    get_session_working_path,
    get_workspace_dir,
)


class TestSessionLifecycle:
    """Test opening presentation sessions, workspace directory structures, and metadata persistence."""

    def test_open_presentation_initializes_session(self, temp_deck_path: Path, clean_ppt_env):
        """open_presentation provisions isolated directory, copies original and working copies, writes metadata."""
        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        assert isinstance(session, Session)
        assert session.session_id is not None
        assert session.slide_count == 3
        assert Path(session.working_path).exists()
        assert Path(session.source_path) == temp_deck_path.resolve()

        # Check directory hierarchy
        session_dir = get_session_dir(session.session_id)
        assert session_dir.exists()
        assert (session_dir / "working.pptx").exists()
        assert (session_dir / "original.pptx").exists()
        assert (session_dir / "metadata.json").exists()
        assert (session_dir / "backups").is_dir()
        assert (session_dir / "renders").is_dir()
        assert (session_dir / "diffs").is_dir()

        # Verify metadata.json content
        with open(session_dir / "metadata.json", "r", encoding="utf-8") as f:
            meta = json.load(f)
        assert meta["session_id"] == session.session_id
        assert meta["slide_count"] == 3
        assert meta["source_path"] == str(temp_deck_path.resolve())

    def test_non_destructive_working_copy_editing(self, temp_deck_path: Path, clean_ppt_env):
        """Edits applied to working.pptx do not alter original.pptx or the source presentation."""
        source_original_bytes = temp_deck_path.read_bytes()

        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        # Mutate the working copy
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="NEW MODIFIED TITLE TEST")
        working_prs.save(session.working_path)

        # Verify working.pptx has new title
        inspected_working = inspect_slide(session.working_path, 1)
        title_shape = inspected_working.shapes[0]
        assert "NEW MODIFIED TITLE TEST" in title_shape.text_frame.text

        # Verify original.pptx in session is untouched
        session_dir = get_session_dir(session.session_id)
        inspected_orig = inspect_slide(str(session_dir / "original.pptx"), 1)
        assert "Quarterly Performance Overview" in inspected_orig.shapes[0].text_frame.text

        # Verify source presentation on disk is untouched
        assert temp_deck_path.read_bytes() == source_original_bytes


class TestBackupsAndRevert:
    """Test timestamped backup creation and revert mechanisms."""

    def test_create_backup_timestamped_naming(self, temp_deck_path: Path, clean_ppt_env):
        """create_backup creates a snapshot with presentation.backup-YYYYMMDD-HHMMSS.pptx naming."""
        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        backup_path_str = manager.create_backup(session.session_id, operation="test_op", label="Test Label")
        backup_path = Path(backup_path_str)

        assert backup_path.exists()
        assert backup_path.name.startswith(f"{temp_deck_path.stem}.backup-")
        assert backup_path.name.endswith(".pptx")
        assert backup_path.parent == get_session_dir(session.session_id) / "backups"

        # Check session metadata updated
        session_meta = manager.get_session(session.session_id)
        assert len(session_meta.backups) == 1
        assert session_meta.backups[0]["backup_path"] == str(backup_path)
        assert session_meta.backups[0]["operation"] == "test_op"

    def test_revert_to_original(self, temp_deck_path: Path, clean_ppt_env):
        """revert_session('original') restores working.pptx to pristine original state."""
        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        # Edit working copy
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="MODIFIED TEXT FOR REVERT TEST")
        working_prs.save(session.working_path)

        # Confirm modified
        slide_mod = inspect_slide(session.working_path, 1)
        assert "MODIFIED TEXT FOR REVERT TEST" in slide_mod.shapes[0].text_frame.text

        # Revert to original
        res = manager.revert_session(session.session_id, "original")
        assert res["success"] is True
        assert res["reverted_to"] == "original"

        # Confirm restored
        slide_restored = inspect_slide(session.working_path, 1)
        assert "Quarterly Performance Overview" in slide_restored.shapes[0].text_frame.text

    def test_revert_to_specific_backup(self, temp_deck_path: Path, clean_ppt_env):
        """revert_session(backup_path) restores working.pptx to a specific intermediate snapshot."""
        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        # Step 1: Version 1
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="VERSION 1 STATE")
        working_prs.save(session.working_path)
        b1_path = manager.create_backup(session.session_id, label="v1")

        # Step 2: Version 2
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="VERSION 2 STATE")
        working_prs.save(session.working_path)
        b2_path = manager.create_backup(session.session_id, label="v2")

        # Confirm current is Version 2
        assert "VERSION 2 STATE" in inspect_slide(session.working_path, 1).shapes[0].text_frame.text

        # Revert to Version 1
        res = manager.revert_session(session.session_id, b1_path)
        assert res["success"] is True
        assert "VERSION 1 STATE" in inspect_slide(session.working_path, 1).shapes[0].text_frame.text


class TestSaveAndSaveAs:
    """Test committing changes to source or saving to new destinations."""

    def test_save_session_overwrites_source_with_presave_backup(self, temp_deck_path: Path, clean_ppt_env):
        """save_session creates a pre-save backup before committing working copy to source."""
        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        # Make an edit in working copy
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="COMMITTED TO SOURCE")
        working_prs.save(session.working_path)

        # Save session
        save_res = manager.save_session(session.session_id)
        assert save_res["success"] is True
        assert save_res["saved_path"] == str(temp_deck_path.resolve())
        assert save_res["backup_path"] is not None

        # Verify pre-save backup exists
        assert Path(save_res["backup_path"]).exists()

        # Verify source file now contains the changes
        source_inspected = inspect_slide(str(temp_deck_path), 1)
        assert "COMMITTED TO SOURCE" in source_inspected.shapes[0].text_frame.text

    def test_save_as_new_destination(self, temp_deck_path: Path, tmp_path: Path, clean_ppt_env):
        """save_as writes working copy to new destination, leaves source intact, and enforces overwrite."""
        source_bytes = temp_deck_path.read_bytes()
        new_output_path = tmp_path / "exported_presentation.pptx"

        manager = SessionManager()
        session = manager.open_presentation(temp_deck_path)

        # Make an edit
        working_prs = Presentation(session.working_path)
        modify_text(working_prs, 1, 2, text="SAVE AS COPY")
        working_prs.save(session.working_path)

        # Save as new file
        res = manager.save_as(session.session_id, new_output_path)
        assert res["success"] is True
        assert new_output_path.exists()

        # New file has edit
        assert "SAVE AS COPY" in inspect_slide(str(new_output_path), 1).shapes[0].text_frame.text

        # Source file is completely unmodified
        assert temp_deck_path.read_bytes() == source_bytes

        # Attempting save_as without overwrite raises FileExistsError
        with pytest.raises(FileExistsError):
            manager.save_as(session.session_id, new_output_path, overwrite=False)

        # Save as with overwrite=True succeeds
        res2 = manager.save_as(session.session_id, new_output_path, overwrite=True)
        assert res2["success"] is True
        assert res2["backup_path"] is not None


class TestSessionRetrievalAndHelpers:
    """Test global helpers, session retrieval, and path utilities."""

    def test_global_helper_functions(self, temp_deck_path: Path, clean_ppt_env):
        """open_presentation, get_session, get_current_session, save_session work via module functions."""
        sess = open_presentation(temp_deck_path)
        assert sess is not None

        curr = get_current_session()
        assert curr.session_id == sess.session_id

        by_id = get_session(sess.session_id)
        assert by_id.session_id == sess.session_id

    def test_ensure_session_dirs_and_cleanup(self, tmp_path: Path):
        """ensure_session_dirs creates all paths and cleanup_session safely deletes them."""
        test_sid = "test-session-12345"
        dirs = ensure_session_dirs(test_sid, base_dir=tmp_path)

        assert dirs["session_dir"].exists()
        assert dirs["backups_dir"].exists()
        assert dirs["renders_dir"].exists()
        assert dirs["diffs_dir"].exists()

        # Cleanup
        assert cleanup_session(test_sid, base_dir=tmp_path) is True
        assert not dirs["session_dir"].exists()

    def test_generate_backup_filename_format(self):
        """generate_backup_filename produces correct pattern."""
        dt = datetime(2026, 8, 21, 14, 30, 0)
        name = generate_backup_filename("sample.pptx", dt)
        assert name == "sample.backup-20260821-143000.pptx"

    def test_error_conditions(self, tmp_path: Path):
        """SessionManager properly raises descriptive exceptions for invalid inputs."""
        manager = SessionManager()

        # Non-existent presentation
        with pytest.raises(FileNotFoundError):
            manager.open_presentation(tmp_path / "non_existent.pptx")

        # Non-existent session revert
        with pytest.raises(ValueError):
            manager.revert_session("non-existent-session-id")
